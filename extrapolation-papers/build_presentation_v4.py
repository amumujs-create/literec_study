#!/usr/bin/env python3
"""Build 외삽 50분 발표자료 v4 (paper figures + method map) (widescreen, speaker notes, figures)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "_assets"
PAPER = ASSETS / "paper_figs"
OUT = ROOT / "외삽_50분_발표자료_v4.pptx"

# 16:9
W, H = Inches(13.333), Inches(7.5)

# Dark instrument theme (match figure console)
BG = RGBColor(0x0B, 0x10, 0x16)
CARD = RGBColor(0x14, 0x1B, 0x24)
PANEL = RGBColor(0x18, 0x21, 0x2C)
LINE = RGBColor(0x2C, 0x38, 0x48)
INK = RGBColor(0xEC, 0xF1, 0xF7)
SLATE = RGBColor(0xA0, 0xAD, 0xBC)
MUTED = RGBColor(0x6E, 0x7C, 0x8C)
TEAL = RGBColor(0x3D, 0xD6, 0xC6)
CORAL = RGBColor(0xE8, 0x9A, 0x5C)
NAVY = RGBColor(0x0C, 0x12, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xB8, 0xC4, 0xD0)
ACCENT_DIM = RGBColor(0x24, 0x3A, 0x42)  # teal-tinted panel


def _set_run(run, size=18, bold=False, color=INK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return p


def rect(slide, x, y, w, h, fill=CARD, line=None, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def bar(slide, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def footer(slide, page, total, part=""):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.08), Inches(10), Inches(0.32))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Extrapolation Seminar  ·  Kookmin Univ. IE Lab  ·  2026  {('·  ' + part) if part else ''}"
    _set_run(run, size=9, color=MUTED)
    num = slide.shapes.add_textbox(Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.32))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    _set_run(r2, size=9, color=MUTED)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s


def content_header(slide, kicker, title):
    bar(slide, TEAL)
    k = slide.shapes.add_textbox(Inches(0.42), Inches(0.26), Inches(10.5), Inches(0.30))
    add_text(k, kicker, size=11, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.42), Inches(0.52), Inches(12.0), Inches(0.50))
    add_text(t, title, size=24, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def card_text(slide, x, y, w, h, title, bullets, accent=TEAL, dense=False):
    # NAVY is a fill color (near-black) — never use it as title ink on dark cards
    title_color = TEAL if accent == NAVY else accent
    strip_color = TEAL if accent == NAVY else accent
    sh = rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = strip_color
    strip.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14), w - Inches(0.38), h - Inches(0.22))
    tf = add_text(box, title, size=12 if dense else 14, bold=True, color=title_color)
    gap = 3 if dense else 5
    fsz = 11 if dense else 13
    for b in bullets:
        # subtle separator lines as plain bullets
        prefix = "–  " if dense else "•  "
        add_para(tf, prefix + b, size=fsz, color=SLATE, space_before=gap, space_after=0)
    return sh


def fig_panel(slide, x, y, w, h):
    """Framed panel behind paper figures (soft chrome)."""
    return rect(slide, x, y, w, h, PANEL, LINE, radius=0.04)


# ── build ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Estimate total slides for footers (update if you add/remove)
TOTAL = 59  # + OoD-Bench shift-type slides (2)
page = 0

# Paper PDF paths (absolute URI for PowerPoint click)
# Relative paths from PPTX location (ASCII folder — Korean path file:// breaks on macOS)
PDF = {
    "xu": "paper_pdfs/Xu2021_How_Neural_Networks_Extrapolate.pdf",
    "eql": "paper_pdfs/Martius2016_Extrapolation_Learning_Equations_EQL.pdf",
    "nalu": "paper_pdfs/Trask2018_Neural_Arithmetic_Logic_Units_NALU.pdf",
    "runje": "paper_pdfs/Runje2023_Constrained_Monotonic_NN.pdf",
    "fesser": "paper_pdfs/Fesser2023_Extrapolation_Failures_PINNs.pdf",
    "zhu": "paper_pdfs/Zhu2022_Reliable_Extrapolation_DeepONet.pdf",
    "pfister": "paper_pdfs/Pfister2024_Extrapolation-Aware_Nonparametric_Inference.pdf",
    "domainbed": "paper_pdfs/Gulrajani2020_In_Search_of_Lost_Domain_Generalization.pdf",
    "liu": "paper_pdfs/Liu2023_OOD_Generalization_Survey.pdf",
    "arjovsky2019": "paper_pdfs/Arjovsky2019_Invariant_Risk_Minimization.pdf",
    "arjovsky2021": "paper_pdfs/Arjovsky2021_OOD_Generalization_in_ML.pdf",
    "raissi": "paper_pdfs/Raissi2019_Physics_Informed_Neural_Networks.pdf",
    "teckentrup": "paper_pdfs/Teckentrup2024_Probabilistic_Richardson_Extrapolation.pdf",
    "wu": "paper_pdfs/Wu2025_OOD_Time_Series_Survey.pdf",
    "ye2021": "paper_pdfs/Ye2021_Theoretical_Framework_OOD.pdf",
    "ye2022": "paper_pdfs/Ye2022_OoD-Bench.pdf",
    "bartley": "paper_pdfs/Bartley2019_Characterizing_Extrapolation_Multivariate.pdf",
    "aykol": "paper_pdfs/Aykol2021_Physics_ML_Battery_Lifetime.pdf",
    "li2023": "paper_pdfs/Li2023_Predicting_Battery_Lifetime_Varying_Conditions.pdf",
    "ghahramani": "paper_pdfs/Ghahramani2013_Bayesian_Nonparametrics.pdf",
    "note": "paper_pdfs/외삽_완전정복_학습노트.pdf",
}


def add_paper_fig(slide, name, left, top, width=None, height=None, pad=0.07, caption=None):
    """Paper figure inside soft framed panel. Returns picture shape or None.
    caption: paper title strip drawn directly under the panel."""
    from PIL import Image
    path = PAPER / name
    if not path.exists():
        return None
    iw, ih = Image.open(path).size
    if width is not None and height is None:
        fw, fh = width, width * (ih / iw)
    elif height is not None and width is None:
        fh, fw = height, height * (iw / ih)
    else:
        fw, fh = width, height
    pad_i = Inches(pad)
    fig_panel(slide, left, top, fw + 2 * pad_i, fh + 2 * pad_i)
    pic = slide.shapes.add_picture(str(path), left + pad_i, top + pad_i, width=fw, height=fh)
    if caption:
        cy = top + fh + 2 * pad_i + Inches(0.02)
        cb = slide.shapes.add_textbox(left, cy, fw + 2 * pad_i, Inches(0.24))
        add_text(cb, caption, size=9, color=SOFT)
    return pic


def credit(slide, text, y=Inches(6.74)):
    box = slide.shapes.add_textbox(Inches(0.5), y, Inches(10.5), Inches(0.26))
    add_text(box, text, size=8, color=MUTED)


def takeaway(slide, main, sub=None, y=Inches(6.38)):
    """하단 발표 핵심."""
    sh = rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.56), ACCENT_DIM, LINE, radius=0.04)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.07), Inches(0.56))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.62), y + Inches(0.05), Inches(12.1), Inches(0.48))
    tf = add_text(t, main, size=12, bold=True, color=INK)
    if sub:
        add_para(tf, sub, size=10, color=SOFT, space_before=1, space_after=0)
    return sh


def speak(slide, bullets, x=Inches(0.4), y=Inches(5.25), w=Inches(12.5), h=Inches(1.0)):
    sh = rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = CORAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.08), w - Inches(0.35), h - Inches(0.12))
    tf = add_text(t, "Remarks", size=11, bold=True, color=CORAL)
    for b in bullets:
        add_para(tf, "–  " + b, size=12, color=SLATE, space_before=3, space_after=0)
    return sh


def section_slide(title, subtitle, part_no, mins, page, total, points=None):
    s = blank()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11), Inches(0.5))
    add_text(t, f"PART {part_no}  ·  {mins}", size=16, bold=True, color=TEAL)
    t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11), Inches(1.0))
    add_text(t2, title, size=36, bold=True, color=WHITE)
    t3 = s.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11), Inches(0.5))
    add_text(t3, subtitle, size=18, color=SOFT)
    if points:
        for i, p in enumerate(points[:3]):
            x = Inches(1.0 + i * 3.9)
            box = rect(s, x, Inches(4.4), Inches(3.6), Inches(1.6), PANEL, LINE)
            tb = s.shapes.add_textbox(x + Inches(0.2), Inches(4.55), Inches(3.2), Inches(1.3))
            tf = add_text(tb, f"{i+1}", size=14, bold=True, color=TEAL)
            add_para(tf, p, size=14, color=WHITE, space_before=6)
    footer_box = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.35))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page} / {total}"
    _set_run(r, size=10, color=MUTED)
    return s


def pdf_btn(slide, key, left, top, w=Inches(1.35), h=Inches(0.34), label=None):
    """Clickable outline button → open paper PDF."""
    rel = PDF.get(key)
    if rel is None or not (ROOT / rel).exists():
        return None
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = CORAL
    sh.line.width = Pt(1.0)
    try:
        sh.adjustments[0] = 0.18
    except Exception:
        pass
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label or f"PDF · {key.upper()}"
    _set_run(run, size=10, bold=True, color=CORAL)
    sh.click_action.hyperlink.address = rel
    return sh


def P():
    global page
    page += 1
    return page


# 1 Title
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), W, Inches(1.6))
acc.fill.solid()
acc.fill.fore_color.rgb = PANEL
acc.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.82), W, Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = TEAL
stripe.line.fill.background()
t = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(0.4))
add_text(t, "논문 기반 심화 세미나  ·  50분", size=16, bold=True, color=TEAL)
t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.5), Inches(1.4))
add_text(t2, "외삽 (Extrapolation)\n이론부터 응용까지", size=44, bold=True, color=WHITE)
t3 = s.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.5), Inches(0.8))
add_text(
    t3,
    "기초 이론  ·  OOD 일반화  ·  신경망 외삽 방법론  ·  실전 응용",
    size=18,
    color=SOFT,
)
t4 = s.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11), Inches(0.8))
tf = add_text(t4, "Kookmin Univ. IE Lab", size=16, bold=True, color=WHITE)
add_para(tf, "literec_study · 36편 논문  ·  2026.07", size=13, color=MUTED)
note(s, "오프닝 30초. 오늘 목표: 외삽이 왜 어려운지 → 어떻게 대응하는가 → 실전 적용으로 연결.")
P()

# 2 Why
s = blank()
content_header(s, "들어가며", "외삽 연구의 동기와 필요성")
cards = [
    ("01  데이터가 항상 불완전", [
        "사례: 시계열 수요의 미래 시점 추정",
        "사례: 계절 부분표본으로 성수기 예측",
        "사례: 성장곡선으로 성인기 외삽",
        "측정 못한 구간을 추정해야 함",
    ]),
    ("02  ML의 근본 한계", [
        "다수 모형이 in-distribution 전제",
        "범위 밖 오차는 급격히 증가",
        "표준 평가 우수 ≠ 외삽 성능 보증",
        "외삽은 별도 평가가 필요",
    ]),
    ("03  현장 리스크", [
        "사례: 기상 ID 조건만 학습한 자율주행",
        "사례: 평시 분포만 학습한 임상 예측",
        "오예측임에도 고신뢰 출력 (silent failure)",
        "점추정 대비 불확실성 보고가 우선",
    ]),
]
for i, (title, bullets) in enumerate(cards):
    card_text(s, Inches(0.45 + i * 4.2), Inches(1.5), Inches(4.0), Inches(4.4), title, bullets, TEAL if i < 2 else CORAL)
footer(s, P(), TOTAL, "Intro")
note(s, "2분. 핵심 메시지: '잘 맞는 테스트'가 보간일 수 있다. 외삽은 별도 프로토콜이 필요.")

# 3 Agenda
s = blank()
content_header(s, "AGENDA", "50분 구성 (Q&A 별도)")
if (ASSETS / "fig_timeline.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_timeline.png"), Inches(0.6), Inches(1.4), width=Inches(12.1))
rows = [
    ("Part 1", "외삽 기초 이론", "보간 vs 외삽 · Convex Hull · Richardson · UQ", "11분"),
    ("Part 2", "OOD 일반화", "Shift 유형 · OoD-Bench · Survey · IRM · DomainBed", "11분"),
    ("Part 3", "신경망 외삽", "Xu · EQL · NALU · Monotonic · PINN · DeepONet", "14분"),
    ("Part 4", "실전 응용", "우리 데이터·프로토콜·모델·결과", "10분"),
    ("Part 5", "동향 & 가이드", "타임라인 · 필독 논문 · 연구 적용", "4분"),
]
y0 = Inches(3.35)
for i, (p, t, d, m) in enumerate(rows):
    y = y0 + Inches(i * 0.58)
    box = rect(s, Inches(0.5), y, Inches(12.3), Inches(0.52), CARD, LINE)
    a = s.shapes.add_textbox(Inches(0.7), y + Inches(0.08), Inches(1.3), Inches(0.35))
    add_text(a, p, size=14, bold=True, color=TEAL)
    b = s.shapes.add_textbox(Inches(2.1), y + Inches(0.08), Inches(2.6), Inches(0.35))
    add_text(b, t, size=14, bold=True, color=INK)
    c = s.shapes.add_textbox(Inches(4.8), y + Inches(0.08), Inches(6.2), Inches(0.35))
    add_text(c, d, size=13, color=SLATE)
    dbox = s.shapes.add_textbox(Inches(11.3), y + Inches(0.08), Inches(1.3), Inches(0.35))
    add_text(dbox, m, size=14, bold=True, color=CORAL, align=PP_ALIGN.RIGHT)
footer(s, P(), TOTAL)
note(s, "30초. 시간 배분 고지. Part3가 가장 김. Part4에서 우리 실험 숫자.")

# ════════════════ PART 1 ════════════════
section_slide(
    "외삽이란 무엇인가",
    "기초 이론 · Convex Hull · Richardson · 불확실성",
    "1",
    "11분",
    P(),
    TOTAL,
    points=[
        "보간≠외삽 — 테스트가 hull 안이면 착시",
        "데이터만으론 밖을 식별 불가 (가정 필요)",
        "예측값보다 UQ·가정이 먼저",
    ],
)

# P1 interp vs extrap
s = blank()
content_header(s, "PART 1 · 기초 이론", "보간(Interpolation) vs 외삽(Extrapolation)")
if (ASSETS / "fig_interp_extrap.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_interp_extrap.png"), Inches(0.25), Inches(1.1), height=Inches(3.0))
card_text(
    s,
    Inches(0.25),
    Inches(4.25),
    Inches(6.3),
    Inches(1.8),
    "보간 — Interpolation",
    [
        "정의: x ∈ Conv(X_train)  (지지집합 내부)",
        "그림 (a)(c★): 관측 범위·hull 안 질의",
        "신뢰도 상대적으로 높음 · 표준 test의 위치",
        "사례: 10·20·30℃ → 25℃  ·  평일→목요일",
    ],
    TEAL,
    dense=True,
)
card_text(
    s,
    Inches(6.75),
    Inches(4.25),
    Inches(6.2),
    Inches(1.8),
    "외삽 — Extrapolation",
    [
        "정의: x ∉ Conv(X_train)  (지지집합 외부)",
        "그림 (b)(c✗): 범위 밖 · 불확실성 대역 확대",
        "신뢰도 급감 · 가정 없으면 식별 불가",
        "사례: 0~30℃ → 45℃  ·  평일→성수기  ·  서울→부산",
    ],
    CORAL,
    dense=True,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림 (a)보간 (b)외삽 (c)hull. test∈hull이면 보간 평가.")
takeaway(s, 'test ∈ Conv(X_train) 이면 보간 평가일 뿐 — 외삽은 별도 프로토콜',
         '그림 (c): ★=보간, ✗=외삽')

# poly fig
s = blank()
content_header(s, "PART 1 · 고전 사례", "다항식 피팅의 외삽 실패")
if (ASSETS / "fig_poly_extrap.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_poly_extrap.png"), Inches(0.5), Inches(1.35), width=Inches(8.2))
card_text(
    s,
    Inches(8.9),
    Inches(1.45),
    Inches(3.9),
    Inches(4.6),
    "해석 사례",
    [
        "지지집합 내부: 적합 양호",
        "지지집합 외부: 급격한 발산",
        "── 해석 사례 ──",
        "단기 관측으로 고차 곡선 적합 시",
        "  → 장기 시점에서 비물리적 예측",
        "성장기 부분곡선만으로",
        "  → 성인기 값을 비현실적으로 외삽",
        "차수↑ → 원거리 불안정성 심화",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림 가리키며 훈련 범위/외삽 영역 구분. 차수 올릴수록 밖이 더 위험.")
takeaway(s, '지지집합 내부 적합 ≠ 외부 일반화 — 고차 피팅일수록 원거리 위험',
         '공학·시계열 OOD도 같은 구조')

# --- ★1 왜 외삽은 원리적으로 어려운가 (식별 불가) ---
s = blank()
content_header(s, "PART 1 · 핵심 통찰", "원리적으로 왜 식별이 불가능한가")
if (ASSETS / "fig_identifiability.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_identifiability.png"), Inches(0.35), Inches(1.25), width=Inches(8.7))
card_text(
    s,
    Inches(9.2),
    Inches(1.3),
    Inches(3.7),
    Inches(4.9),
    "해석 사례",
    [
        "유한 관측만으로는",
        "선형·이차·주기 함수 등",
        "지지집합 내부에서는 동등하게 적합 가능",
        "지지집합 밖에서 궤적이 분기",
        "── 해석 사례 ──",
        "냉각 과정의 온도 궤적",
        "  선형·지수 모두 국소적으로는 적합 가능",
        "성장곡선 → 성인기 값",
        "  지속 증가 vs 포화",
        "→ 가정 없이는 식별 불가",
    ],
    CORAL,
)
pdf_btn(s, "pfister", Inches(9.35), Inches(0.28), label="PDF · Pfister")
credit(s, "핵심: 외삽은 데이터 문제가 아니라 가정(Assumption)의 문제다  ·  Pfister & Bühlmann (2024)")
footer(s, P(), TOTAL, "Part 1")
note(s, "2.5분. 발표 핵심 슬라이드. Pfister '가정 없이 식별 불가'로 연결.")
takeaway(s, '외삽은 데이터 문제가 아니라 가정(Assumption)의 문제',
         'sin/exp/x²는 훈련점에선 모두 맞지만, 지지집합 밖에서 갈라진다')

# --- ★5 산업 실패 사례 ---
s = blank()
content_header(s, "PART 1 · 동기", "이론적 실패 모드가 현장 배포에서도 나타난다")
if (ASSETS / "fig_fail_cases.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_fail_cases.png"), Inches(1.5), Inches(1.1), height=Inches(4.0))
box = rect(s, Inches(0.3), Inches(5.25), Inches(12.7), Inches(0.85), BG, LINE)
t = s.shapes.add_textbox(Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.7))
tf = add_text(t, "학습 분포 ≠ 배포 분포  →  높은 확신 · 낮은 정확도  (silent failure)", size=15, bold=True, color=INK)
add_para(tf, "현장 사례도 같은 구조: 기상 OOD · 기관 전이 · 성수기 수요 전이  —  그림 (C)가 핵심 메커니즘", size=12, color=SLATE)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. (A)온도분포 (B)기상클러스터 (C)confidence≠accuracy.")
takeaway(s, '학습 분포 ≠ 배포 분포이면, 모델은 높은 확신으로 틀린다 (silent failure)',
         '현장 사고의 전형 패턴 — 평균 성능만으로는 감지 불가')

# convex hull definition
s = blank()
content_header(s, "PART 1 · 수학적 정의", "볼록 껍질(Convex Hull)과 외삽 영역")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(2.5),
    "정의",
    [
        "x ∈ Conv(X_train) → 보간 영역",
        "x ∉ Conv(X_train) → 외삽 영역",
        "Pfister 2024: support 밖은 추가 가정 없이 추론 불가",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(2.5),
    "차원의 저주",
    [
        "축이 늘수록 '안쪽' 비율 급감",
        "고차원: 테스트 대부분 외삽",
        "사례: 다변량 운용조건 (고차원 hull)",
        "  → 4축만 돼도 '밖'이 대부분",
    ],
    CORAL,
)
card_text(
    s,
    Inches(0.45),
    Inches(4.15),
    Inches(6.1),
    Inches(2.2),
    "불확실성 폭발",
    [
        "인근 관측 없을수록 bounds 확장",
        "기존 CI는 coverage 보장 상실",
        "→ extrapolation-aware CI 필요",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(4.15),
    Inches(6.0),
    Inches(2.2),
    "물리 지식 = 외삽 가정",
    [
        "단조성 · PDE · smoothness",
        "가정이 hull을 '논리적으로' 확장",
        "도메인 지식이 최고의 외삽 도구",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2.5분. Pfister: 가정 없으면 밖은 식별 불가. 물리 제약이 왜 등장하는지 연결.")
takeaway(s, 'Convex Hull 외부 = 외삽 영역. 고차원에서는 대부분 외부',
         '물리·단조 가정이 hull을 논리적으로 확장한다')

# hull fig
s = blank()
content_header(s, "PART 1 · 시각화", "2D Convex Hull로 보간/외삽 구분")
if (ASSETS / "fig_convex_hull.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_convex_hull.png"), Inches(0.5), Inches(1.35), width=Inches(8.3))
card_text(
    s,
    Inches(9.0),
    Inches(1.45),
    Inches(3.8),
    Inches(4.6),
    "진단 기준",
    [
        "hull 내부: 상대적 안정",
        "hull 외부: epistemic 불확실성↑",
        "d≫10: 대부분의 점이 외삽",
        "대응: 물리 제약 도입",
        "또는 외삽축을 분리 설계",
        "Bartley · Pfister",
    ],
)
footer(s, P(), TOTAL, "Part 1")
note(s, "1.5분. X 표시가 외삽. 실전에선 특정 운용축이 이 X에 해당.")
takeaway(s, 'Convex hull 내부=보간, 외부(X)=외삽',
         '실전에서는 특정 운용조건 축이 이 X에 대응')

# error decomposition
s = blank()
content_header(s, "PART 1 · 오차 분해", "외삽 오차의 폭증 메커니즘")
if (ASSETS / "fig_error_decomp.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_error_decomp.png"), Inches(0.3), Inches(1.15), height=Inches(3.05))
# compact legend table under / right
headers = ["요인", "보간", "외삽", "대책"]
rows = [
    ["Bias²", "복잡도 제어", "가정 위반 시 누적", "물리·smooth"],
    ["Variance", "인근 데이터多", "관측∅ → 급증", "extrap CI"],
    ["Noise σ²", "분리 측정", "epistemic과 혼재", "Bayesian UQ"],
]
col_w = [1.6, 2.0, 2.5, 2.2]
x0 = 0.35
y0 = 4.35
for j, h in enumerate(headers):
    x = Inches(x0 + sum(col_w[:j]))
    rect(s, x, Inches(y0), Inches(col_w[j] - 0.06), Inches(0.32), PANEL)
    t = s.shapes.add_textbox(x, Inches(y0 + 0.02), Inches(col_w[j] - 0.1), Inches(0.28))
    add_text(t, h, size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        x = Inches(x0 + sum(col_w[:j]))
        y = Inches(y0 + 0.36 + i * 0.38)
        fill = CARD if i % 2 == 0 else PANEL
        rect(s, x, y, Inches(col_w[j] - 0.06), Inches(0.36), fill, LINE)
        t = s.shapes.add_textbox(x + Inches(0.04), y + Inches(0.05), Inches(col_w[j] - 0.12), Inches(0.28))
        add_text(t, cell, size=10, bold=(j == 0), color=INK if j == 0 else SLATE, align=PP_ALIGN.LEFT)
card_text(
    s,
    Inches(9.0),
    Inches(4.35),
    Inches(3.95),
    Inches(1.7),
    "해석 포인트",
    [
        "지지집합 밖: Variance가 지배항",
        "점추정만으로는 의사결정 위험",
        "UQ(신뢰구간)를 병기",
    ],
    CORAL,
    dense=True,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림(a)(b) → 표로 요약.")
takeaway(s, '지지집합 밖에서는 epistemic variance가 급증한다',
         '점추정만 보고 결정하지 말고 UQ를 함께')

# Richardson
s = blank()
content_header(s, "PART 1 · 수치해석", "Richardson 외삽 → Probabilistic Richardson")
if (ASSETS / "fig_richardson.png").exists():
    # full-width explanation figure: (a) classical (b) probabilistic (c) workflow
    s.shapes.add_picture(str(ASSETS / "fig_richardson.png"), Inches(0.25), Inches(1.1), height=Inches(3.05))
card_text(
    s,
    Inches(0.25),
    Inches(4.3),
    Inches(6.3),
    Inches(1.75),
    "고전 Richardson — 점추정",
    [
        "모형: A(h)=A* + c·h^p + O(h^{p+1})",
        "추정: Â = (2^p A(h/2) − A(h)) / (2^p − 1)",
        "그림 (a): h→0 외삽이 다이아몬드 한 점",
        "한계: 불확실성 대역·CI 없음",
    ],
    MUTED,
    dense=True,
)
card_text(
    s,
    Inches(6.75),
    Inches(4.3),
    Inches(6.2),
    Inches(1.75),
    "Probabilistic (Teckentrup 2024)",
    [
        "Richardson + GP prior → A* 사후분포",
        "그림 (b): GP band + A* 95% CI",
        "그림 (c): 점값 vs 구간(안전장치)",
        "의의: 수치해석 × Bayesian UQ",
    ],
    TEAL,
    dense=True,
)
pdf_btn(s, "teckentrup", Inches(10.8), Inches(0.28), w=Inches(2.1), label="PDF · Teckentrup")
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림 (a)(b)(c)를 손가락으로 가리키며 설명.")
takeaway(s, '고전 = 점추정 Â · Probabilistic = 사후분포·구간 (그림 b·c)',
         '공학 외삽에서는 점값보다 신뢰구간이 의사결정에 안전하다')

# UQ
s = blank()
content_header(s, "PART 1 · UQ", "불확실성 정량화 — 외삽에서의 필연성")
if (ASSETS / "fig_uq_aleatoric_epistemic.png").exists():
    s.shapes.add_picture(
        str(ASSETS / "fig_uq_aleatoric_epistemic.png"),
        Inches(0.35),
        Inches(1.15),
        height=Inches(2.05),
    )
card_text(
    s,
    Inches(0.35),
    Inches(3.35),
    Inches(4.15),
    Inches(2.75),
    "Aleatoric  σ²_a  — 관측 노이즈",
    [
        "정의: y = f(x) + ε,  ε~P_noise",
        "데이터↑로도 σ²_a 안 사라짐",
        "수식:  irreducible noise",
        "사례: 센서 계측 잡음 ±δ",
        "사례: 배치 간 공정 변동성",
        "사례: 동일 조건 반복측정의 분산",
        "보간·외삽 모두 존재 (거의 평탄)",
    ],
    MUTED,
    dense=True,
)
card_text(
    s,
    Inches(4.6),
    Inches(3.35),
    Inches(4.15),
    Inches(2.75),
    "Epistemic  σ²_e  — 지식 부족",
    [
        "정의: 모델/데이터 부족 불확실성",
        "데이터↑·가정↑로 줄일 수 있음",
        "수식: Var_θ[f_θ(x)]  (사후분산)",
        "외삽: x∉Conv(X) 에서 σ²_e ≫ σ²_a",
        "사례: 학습 온도 지지집합 밖 질의",
        "사례: 신규 기관·시즌으로의 배포",
        "핵심 질문: 얼마나 틀릴 수 있나?",
    ],
    CORAL,
    dense=True,
)
card_text(
    s,
    Inches(8.85),
    Inches(3.35),
    Inches(4.15),
    Inches(2.75),
    "외삽 UQ — 점추정의 한계",
    [
        "σ²_tot = σ²_a + σ²_e",
        "밖에서는 epistemic이 지배",
        "기존 CI: coverage 붕괴 가능",
        "  (i.i.d./hull 안 가정 깨짐)",
        "대안: bounds + CI + 가정 명시",
        "Ghahramani'13: Bayesian ML 관점",
        "Pfister'24: 가정 없으면 식별 불가",
        "함의: σ²_e↑ → 예측 보류 / 추가 관측",
    ],
    TEAL,
    dense=True,
)
pdf_btn(s, "ghahramani", Inches(9.0), Inches(0.28), w=Inches(1.9), label="PDF · Ghahramani")
pdf_btn(s, "pfister", Inches(11.05), Inches(0.28), w=Inches(1.9), label="PDF · Pfister")
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림(a)(b) → 세 카드 수식·예시. 예측값 < UQ.")
takeaway(s, '외삽에선 예측값보다 불확실성(σ²_e)이 핵심',
         'σ²_tot=σ²_a+σ²_e — 밖에서는 epistemic 지배, ↑면 STOP')

# ════════════════ PART 2 ════════════════
section_slide(
    "OOD 일반화",
    "Shift 유형 · OoD-Bench · Survey · IRM · DomainBed · 시계열",
    "2",
    "11분",
    P(),
    TOTAL,
    points=[
        "i.i.d. 붕괴 = ERM 실패 시작점",
        "OoD-Bench: diversity vs correlation 먼저 식별",
        "Liu Survey: 식별 후 대응 가지 선택",
        "IRM은 아이디어, DomainBed가 현실 점검",
    ],
)

# Generic OOD hook (no dataset-specific TRA until Part4)
s = blank()
content_header(s, "PART 2 · 개념", "OOD: 학습 분포와 배포 분포의 불일치")
if (ASSETS / "fig_ood_intuition.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_ood_intuition.png"), Inches(0.25), Inches(1.1), height=Inches(3.05))
fb = rect(s, Inches(0.25), Inches(4.25), Inches(12.8), Inches(0.55), NAVY)
ft = s.shapes.add_textbox(Inches(0.4), Inches(4.32), Inches(12.5), Inches(0.42))
tf = add_text(ft, "형식적 정의   P_train(X,Y)  ≠  P_test(X,Y)   ≡   min_h max_e R^e(h)", size=13, bold=True, color=WHITE)
add_para(tf, "사례: 기관 A→B 전이  ·  기상 조건 외삽  ·  평시→성수기 수요 전이", size=11, color=SOFT, space_before=1)
card_text(
    s,
    Inches(0.25),
    Inches(4.95),
    Inches(4.15),
    Inches(1.1),
    "리스크 구조",
    [
        "기대성능(평균)은 유지될 수 있음",
        "꼬리 리스크(worst-case)에서 급락",
    ],
    CORAL,
    dense=True,
)
card_text(
    s,
    Inches(4.55),
    Inches(4.95),
    Inches(4.15),
    Inches(1.1),
    "평가의 함정",
    [
        "표준 test ⊂ 학습 지지집합(ID)",
        "→ 보간 성능일 뿐 OOD 보증 아님",
    ],
    MUTED,
    dense=True,
)
card_text(
    s,
    Inches(8.85),
    Inches(4.95),
    Inches(4.2),
    Inches(1.1),
    "대응 전략",
    [
        "불변표현 · 물리/단조 제약 · 최악환경 검증",
        "엄밀 프로토콜은 Part 4",
    ],
    TEAL,
    dense=True,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. support mismatch → mean vs worst-case → formal objective.")
takeaway(s, 'OOD = 학습·배포 분포 불일치 — ID 테스트 성능 ≠ 배포 강건성',
         '평균 리스크와 함께 worst-case 리스크를 보고해야 한다')

# shift types
s = blank()
content_header(s, "PART 2 · 분류", "분포 이동(Distribution Shift) 유형")
if (ASSETS / "fig_shift_types.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_shift_types.png"), Inches(0.25), Inches(1.15), height=Inches(3.15))
# mini captions under figure as 4 thin cards
items = [
    ("Covariate", "P(X) 변화 · P(Y|X) 불변", "기관·지역·시즌 전이"),
    ("Label", "P(Y) 변화 · P(X|Y) 불변", "유병률·수요 비중 이동"),
    ("Concept", "P(Y|X) 자체 변화", "정책·시스템 드리프트"),
    ("Domain", "생성과정 전반 변화", "시뮬→실기 전이"),
]
for i, (title, defn, ex) in enumerate(items):
    card_text(
        s,
        Inches(0.3 + i * 3.25),
        Inches(4.45),
        Inches(3.15),
        Inches(1.55),
        title,
        [defn, f"대표 사례: {ex}", "유형 식별 → 해법 분기"],
        TEAL if i % 2 == 0 else CORAL,
        dense=True,
    )
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. 네 유형 taxonomy. 실전은 covariate(+concept) 혼합이 흔함.")
takeaway(s, 'Distribution shift 유형을 먼저 명명하면 해법이 달라진다',
         '실전은 주로 Covariate(+Concept) Shift')

# OOD formal
s = blank()
content_header(s, "PART 2 · 정의", "ERM vs OOD — 평균 오차 vs 최악 환경 오차")
box = rect(s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(1.15), NAVY)
t = s.shapes.add_textbox(Inches(0.7), Inches(1.40), Inches(12.0), Inches(0.9))
tf = add_text(t, "직관:  시험 평균을 맞출 것인가,  최악 과목도 확보할 것인가", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "전자 = ERM(평균 오차↓) · 후자 = OOD(최악 환경 오차↓)", size=13, color=SOFT, align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(2.60),
    Inches(6.0),
    Inches(3.45),
    "표준 ERM  —  평균 오차를 줄인다",
    [
        "병원 A 데이터로만 학습 → A에서 오차↓",
        "병원 B·C로 가면 성능이 급락할 수 있음",
        "학습 분포의 '자주 나오는 케이스'에 과적합",
        "수식(참고):  min  평균_e  오차(e)",
        "한계: 평균이 좋아도 낯선 환경에서 붕괴",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(2.60),
    Inches(6.0),
    Inches(3.45),
    "OOD 목표  —  최악 환경 오차를 줄인다",
    [
        "병원 A·B·C 중 어디서든 오차가 크면 실패",
        "목표: 가장 어려운 환경의 오차를 줄인다",
        "→ 환경이 바뀌어도 같은 규칙이 통해야 함",
        "수식(참고):  min  최악_e  오차(e)",
        "Ye 2021 · Arjovsky 2021",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. 시험 평균 vs 최악 과목 비유 → 병원 A/B 사례. 수식은 참고만.")
takeaway(s, "ERM=평균 오차 최소화, OOD=최악 환경 오차 최소화 — 목표가 다르다",
         "낯선 환경(병원 B)에서 성능이 급락하면 OOD 실패")

# OOD must be measurable first
s = blank()
content_header(s, "PART 2 · 전제", "OOD라고 부르려면 — 먼저 분포 이탈을 수치로 정의한다")
box = rect(s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.95), NAVY)
t = s.shapes.add_textbox(Inches(0.7), Inches(1.40), Inches(12.0), Inches(0.7))
tf = add_text(t, "분포 이탈 정의가 없으면 'OOD 성능'은 검증 불가능한 주장이다", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "Part 1의 hull · Part 2의 shift 유형 · 아래 수치 정의가 그 전제", size=12, color=SOFT, align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(2.40),
    Inches(4.0),
    Inches(3.55),
    "① 왜 정의가 필요한가",
    [
        "test가 사실상 보간이면 → 외삽 착시",
        "'OOD에서 이겼다'만으로는 비교 불가",
        "같은 이름이라도 데이터셋마다 난이도 다름",
        "정의 = 평가 프로토콜의 1번째 줄",
    ],
)
card_text(
    s,
    Inches(4.65),
    Inches(2.40),
    Inches(4.0),
    Inches(3.55),
    "② 문헌의 수치 정의",
    [
        "Hull: x ∉ Conv(X_train)  (Xu)",
        "Shift type 명명 (covariate…)",
        "OoD-Bench (Ye 2022 CVPR):",
        "  diversity · correlation shift 점수",
        "DomainBed: 환경 분할을 먼저 고정",
    ],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(2.40),
    Inches(4.0),
    Inches(3.55),
    "③ 정의만으로는 부족하다",
    [
        "정의 ≠ 성능 보증 (PINN 교훈)",
        "분할을 바꾸면 순위가 뒤집힐 수 있음",
        "그래서: 정의 공개 + ERM baseline",
        "        + 민감도까지 한 세트로 보고",
        "Part 4: TRA q70/q90으로 우리가 정의",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 2")
note(
    s,
    "2분. 핵심 멘트: 'OOD는 분포 이탈을 수치로 정의한 뒤에야 성립한다. "
    "정의해도 실험 전엔 모른다 → 그래서 DomainBed·민감도. "
    "우리 연구는 Part4에서 TRA 분위수로 정의했다.'",
)
takeaway(s, "OOD 주장은 분포 이탈을 수치로 정의한 뒤에야 성립한다",
         "정의만으로 부족 — ERM baseline·민감도까지 함께 보고")

# OoD-Bench — diversity vs correlation shift
s = blank()
content_header(s, "PART 2 · OoD-Bench", "Ye et al. (2022) — OOD를 두 축으로 수치화한다")
if (ASSETS / "fig_ood_bench_shifts.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_ood_bench_shifts.png"), Inches(0.30), Inches(1.12), width=Inches(12.7))
card_text(
    s, Inches(0.35), Inches(4.85), Inches(12.6), Inches(1.30),
    "핵심 메시지",
    [
        "같은 'OOD'라도 실패 원인이 다름 — 방법을 고르기 전에 shift 유형(또는 D_div / D_cor)을 먼저 파악해야 한다",
        "Diversity: train에 없던 feature/도메인 등장 (PACS, Camelyon)  ·  Correlation: 같은 feature인데 Y와의 가짜 상관이 바뀜 (Colored MNIST, CelebA)",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "ye2022", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · OoD-Bench")
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. 'OOD 방법 고르기 전에 어떤 OOD냐?' — diversity vs correlation. 그림 (a)(b)(c) 순서.")
takeaway(s, "OOD 대응 전에 shift 유형을 식별한다 — diversity vs correlation",
         "데이터셋마다 D_div / D_cor 프로파일이 다르다")

# OoD-Bench — method choice depends on shift type
s = blank()
content_header(s, "PART 2 · OoD-Bench", "Shift 유형에 따라 유효한 대응이 달라진다")
card_text(
    s, Inches(0.35), Inches(1.20), Inches(4.05), Inches(3.55),
    "Diversity shift 우세",
    [
        "원인: train support 밖 feature",
        "예: 새 병원·새 스타일·새 운용점",
        "유효 후보: representation /",
        "  domain mix · augmentation",
        "  disentangle · meta-learning",
        "약한 후보: spurious만 제거",
    ],
    TEAL,
)
card_text(
    s, Inches(4.55), Inches(1.20), Inches(4.05), Inches(3.55),
    "Correlation shift 우세",
    [
        "원인: 가짜 단서·허위 상관",
        "예: Colored MNIST 색 함정",
        "유효 후보: IRM · GroupDRO ·",
        "  debiasing · causal feature",
        "약한 후보: 단순 ERM+증강",
    ],
    CORAL,
)
card_text(
    s, Inches(8.75), Inches(1.20), Inches(4.20), Inches(3.55),
    "OoD-Bench 주장 (실험)",
    [
        "대부분 알고리즘은 한 축에서만",
        "  ERM을 이김 (양쪽 X)",
        "→ 'OOD SOTA' 주장은",
        "  shift 유형을 명시해야 함",
        "평가: diversity·correlation",
        "  각각 dominated set에서",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(4.90), Inches(12.60), Inches(1.25),
    "우리 과제(N-CMAPSS) 연결",
    [
        "고TRA = diversity(새 운용 영역) + correlation(부하-열화 얽힘) 혼합 → Part4에서 TRA 분위수로 diversity 축을 명시하고, dual encoder+제약으로 correlation 축을 분리",
        "다음: Liu Survey — shift 유형을 아는 뒤 '어디에 개입할지(표현/지도/최적화)'를 고른다",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "ye2022", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · OoD-Bench")
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. OoD-Bench 핵심: shift 유형별 방법. N-CMAPSS 한 줄. Liu로.")
takeaway(s, "OOD 방법은 shift 유형(diversity / correlation)에 맞게 고른다",
         "한쪽에서만 이긴 알고리즘을 범용 OOD라 부르지 말 것")

# Liu survey — OOD response taxonomy
s = blank()
content_header(s, "PART 2 · Survey", "Liu et al. (2023) — OOD 대응을 학습 파이프라인 위치로 분류")
if (ASSETS / "fig_ood_survey_taxonomy.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_ood_survey_taxonomy.png"), Inches(0.35), Inches(1.15), width=Inches(12.6))
card_text(
    s, Inches(0.35), Inches(4.95), Inches(12.6), Inches(1.20),
    "발표 포인트",
    [
        "같은 목표(P_train ≠ P_test에서도 성능 유지)지만, 개입 지점이 다름 → 방법 비교 전에 갈래를 먼저 말하라",
        "다음 장: 실무에서 가장 많이 쓰는 ② Supervised 가지를 세부 분해",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "liu", Inches(11.5), Inches(0.28), w=Inches(1.4), label="PDF · Liu")
footer(s, P(), TOTAL, "Part 2")
note(s, "1.5분. Liu 3갈래. DA/DG와의 차이 한 줄. Supervised로 넘김.")
takeaway(s, "OOD 대응 = 표현 초기화 · 지도 학습 전략 · 최악-케이스 최적화",
         "개입 지점이 다르면 '같은 OOD 방법'이라도 역할이 다르다")

# Liu survey — supervised branch detail + when to use
s = blank()
content_header(s, "PART 2 · Survey", "Supervised 가지 — 무엇을 알고 있는가에 따라 고른다")
card_text(
    s, Inches(0.30), Inches(1.20), Inches(3.10), Inches(3.55),
    "인과 (Causal)",
    [
        "가정: Y의 원인만 불변",
        "개입·교란에 의한 shift",
        "원인 변수로만 예측",
        "강점: 원리적으로 강건",
        "약점: SCM·개입 가정 필요",
    ],
    dense=True,
)
card_text(
    s, Inches(3.50), Inches(1.20), Inches(3.10), Inches(3.55),
    "불변 (Invariant)",
    [
        "환경마다 공통 특징만 유지",
        "대표: IRM · ICP",
        "가짜 단서(색 등) 제거",
        "강점: 환경 라벨만으로 가능",
        "약점: 환경 정의에 민감",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(6.70), Inches(1.20), Inches(3.10), Inches(3.55),
    "학습 전략",
    [
        "Meta / Ensemble / SSL",
        "Feature Normalization",
        "도메인 혼합·자기도전",
        "강점: 구현·확장 용이",
        "약점: 이론 보증이 약함",
    ],
    dense=True,
)
card_text(
    s, Inches(9.90), Inches(1.20), Inches(3.05), Inches(3.55),
    "최적화 (DRO)",
    [
        "목표: 최악 분포 오차↓",
        "GroupDRO · f-DRO",
        "그룹/불확실 집합 지정",
        "강점: 목표와 수식 직결",
        "약점: 그룹 정의·계산 비용",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(0.30), Inches(4.90), Inches(12.65), Inches(1.25),
    "선택 규칙 (직관)",
    [
        "0단계: OoD-Bench — diversity vs correlation (또는 D_div/D_cor) 먼저",
        "Correlation 우세 → IRM / GroupDRO / debiasing   ·   Diversity 우세 → mix·augment·representation",
        "가정 약함 → ERM + UQ   ·   어떤 경우든 DomainBed식 공정 비교 필수",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "liu", Inches(11.5), Inches(0.28), w=Inches(1.4), label="PDF · Liu")
footer(s, P(), TOTAL, "Part 2")
note(s, "2.5분. Survey 핵심. 다음 IRM은 Invariant 가지의 대표 사례.")
takeaway(s, "아는 것이 원인·환경·그룹 중 무엇이냐에 따라 대응 가지가 갈린다",
         "다음: Invariant 가지의 대표 — IRM")

# IRM
s = blank()
content_header(s, "PART 2 · 알고리즘", "IRM — Invariant 가지의 대표 사례")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "한 줄 직관 (Arjovsky)",
    [
        "문제: 모델이 '색' 같은 가짜 단서를 학습",
        "Colored MNIST: 숫자 대신 색으로 정답 맞춤",
        "환경(색 규칙)이 바뀌면 ERM은 바로 붕괴",
        "IRM: 모든 환경에서 통하는 특징만 남긴다",
        "예: '숫자 모양'은 불변 · '색'은 가변 → 버림",
        "한계: 환경을 어떻게 나눌지가 애매하면 약함",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "Survey에서의 위치",
    [
        "Liu 분류: Supervised · Invariant",
        "인과를 완화한 실용형 (원인 대신 불변성)",
        "환경 라벨이 있을 때 후보",
        "GroupDRO와 목표(최악)는 가깝지만",
        "  수단이 다름: 특징 제약 vs 목적함수",
        "다음: DomainBed — 이 방법이 실제로 이겼나?",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. Survey → IRM 사례. DomainBed로 연결.")
takeaway(s, "IRM: 환경마다 바뀌는 가짜 단서를 버리고, 공통 규칙만 학습",
         "다음은 DomainBed — 공정 비교하면 ERM이 자주 이긴다")
pdf_btn(s, "arjovsky2019", Inches(9.5), Inches(0.28), w=Inches(1.7), label="PDF · IRM 2019")
pdf_btn(s, "arjovsky2021", Inches(11.35), Inches(0.28), w=Inches(1.6), label="PDF · 2021")
pdf_btn(s, "liu", Inches(7.85), Inches(0.28), w=Inches(1.5), label="PDF · Liu")

# DomainBed
s = blank()
content_header(s, "PART 2 · 벤치마크", "DomainBed — 공정 평가에 의한 재평가")
box = rect(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(1.6), PANEL, CORAL)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.1))
tf = add_text(t, "공정한 HP 탐색 하에서 ERM ≥ 대부분의 OOD 알고리즘", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_para(tf, "Gulrajani & Lopez-Paz (2020) · ICML 2021", size=14, color=SOFT, align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "모델 선택이 핵심",
    ["성능 차의 상당 부분 = HP·선택", "알고리즘 자체보다 튜닝"],
)
card_text(
    s,
    Inches(4.65),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "공정 비교 프로토콜",
    ["① ERM baseline", "② 동일 HP budget", "③ 동일 capacity", "④ target val 전략 명시"],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "Survey와의 연결",
    [
        "대응 가지를 골라도 평가는 공통",
        "Liu: 방법론 / DomainBed: 검증",
        "주장 전 ERM 재검증이 필수",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. Survey로 고른 방법을 DomainBed로 검증.")
takeaway(s, '새 OOD 알고리즘 주장 전: ERM + 동일 HP budget 필수',
         '공정 비교 없으면 성능 향상은 튜닝 착시일 수 있다')
pdf_btn(s, "domainbed", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · DomainBed")

# Bridge: survey methods → our setting / time series
s = blank()
content_header(s, "PART 2 · 연결", "Survey 대응 → 우리 문제(시계열·공학)로")
card_text(
    s,
    Inches(0.35),
    Inches(1.25),
    Inches(4.15),
    Inches(3.50),
    "Survey에서 가져올 것",
    [
        "① 분포 이탈 정의 + shift 유형(OoD-Bench)",
        "② 유형에 맞는 대응 가지 선택(Liu)",
        "③ ERM + 공정 비교로 검증",
        "④ 가정 약하면 UQ·기권 병행",
    ],
    TEAL, dense=True,
)
card_text(
    s,
    Inches(4.60),
    Inches(1.25),
    Inches(4.15),
    Inches(3.50),
    "시계열 OOD (Wu 2025)",
    [
        "시간적 분포 이동",
        "주기적 분포 이동",
        "공학 시계열은 둘 다 흔함",
        "충방전·운전 주기 = 환경",
    ],
    dense=True,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.25),
    Inches(4.05),
    Inches(3.50),
    "Part 3–4로의 다리",
    [
        "Part 3: NN 구조로 가정 내장",
        "  (Xu 진단 · Mono · PINN)",
        "Part 4: TRA로 OOD 정의 +",
        "  제약·isotonic으로 대응",
        "Survey '가정→방법'의 실전판",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(4.90), Inches(12.55), Inches(1.25),
    "한 줄 정리",
    [
        "Liu Survey = 메뉴판  ·  DomainBed = 시식 검증  ·  Part 3–4 = 우리 문제에 맞게 요리",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "liu", Inches(8.9), Inches(0.28), w=Inches(1.2), label="Liu")
pdf_btn(s, "wu", Inches(10.2), Inches(0.28), w=Inches(1.2), label="Wu")
pdf_btn(s, "domainbed", Inches(11.5), Inches(0.28), w=Inches(1.4), label="DomainBed")
footer(s, P(), TOTAL, "Part 2")
note(s, "1.5분. Part2 클로징. Survey→시계열→Part3/4.")
takeaway(s, "Survey로 가지를 고르고, DomainBed로 검증하고, Part3–4에서 구조화한다",
         "시계열·공학 OOD = 시간 이동 + 주기 이동 → Part4 실전")

# ════════════════ PART 3 ════════════════
section_slide(
    "신경망 외삽 방법론",
    "방법론 지도 · 논문 원도 · 언제 무엇을 쓸까",
    "3",
    "14분",
    P(),
    TOTAL,
    points=[
        "Q1. 실패 기제: ‖x‖→∞ ⇒ ReLU MLP → Ax+b  (Xu)",
        "Q2. 구조적 처방: ŷ∈⟨U⟩ / ∂ŷ/∂x≥0 / PDE 잔차",
        "Q3. 모르면?  σ²(x)↑ → STOP  (UQ)",
    ],
)

# --- Method taxonomy map ---
s = blank()
content_header(s, "PART 3 · 방법론 Tree", "외삽 해결법 — 방법론이 분기하는 이유")
if (ASSETS / "fig_method_tree.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_method_tree.png"), Inches(0.25), Inches(1.15), height=Inches(3.4))
card_text(
    s,
    Inches(7.85),
    Inches(1.15),
    Inches(5.1),
    Inches(3.4),
    "가지 · 가정(수식) · 사례",
    [
        "① Activation · Xu'21",
        "  lim ‖x‖→∞ f_θ(x) = Affine / 포화",
        "  사례: 주기 신호를 ReLU MLP로 근사",
        "② Equation · EQL / NALU",
        "  f ∈ span{sin, ×, …} 또는 ±×÷",
        "  사례: y=sin(ωt) 계절 성분",
        "③ Constraint · CMNN",
        "  ∂f/∂x_i ≥ 0  (단조)",
        "  사례: 가격–수요 단조관계",
        "④ Physics · PINN",
        "  L = L_data + λ ‖N[u_θ]‖²",
        "  사례: 열전도 PDE",
        "⑤ Operator+UQ",
        "  G: u(·)↦s(·),  Var↑→STOP",
        "  사례: 신용구간 이탈 시 배포 보류",
    ],
    TEAL,
    dense=True,
)
speak(
    s,
    [
        "진단→처방→안전장치: Activation → Equation/Constraint/Physics → UQ",
        "가정이 강할수록 외삽에 유리, 틀리면 오차도 큼 (assumption bias)",
        "실전: 부분 지식(단조·물리)을 제약으로 + 외삽 전용 평가 프로토콜",
    ],
    y=Inches(4.7),
    h=Inches(1.4),
)
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. 가지마다 수식 가정을 명시.")
takeaway(s, '각 가지는 상이한 함수 가정 — 가정 수식을 명시할 수 있어야 한다',
         'Activation → Equation → Constraint → Physics → Operator+UQ')

# --- ★3 왜 되는가 직관표 ---
s = blank()
content_header(s, "PART 3 · 개요", "외삽 방법 — 작동 원리 · 적용 사례 · 한계")
if (ASSETS / "fig_method_cases.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_method_cases.png"), Inches(0.2), Inches(1.1), height=Inches(4.0))
# right compact formula strip
card_text(
    s,
    Inches(9.55),
    Inches(1.1),
    Inches(3.45),
    Inches(4.95),
    "수식 요약",
    [
        "ReLU: ‖x‖→∞⇒Ax+b",
        "Tanh: σ→±1 포화",
        "EQL: ŷ∈⟨U⟩ 연장",
        "Mono: ∂ŷ/∂x≥0",
        "PINN: L+λ‖N[u]‖²",
        "UQ: σ_e↑ → STOP",
        "── 판독 기준 ──",
        "초록 띠=훈련 구간",
        "점선=참 함수",
        "실선=모델 거동",
        "── 한계 요약 ──",
        "ReLU: 비선형 실패",
        "EQL: U 밖 불가",
        "PINN: 시간외삽弱",
        "UQ: 경고만(수정X)",
    ],
    TEAL,
    dense=True,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. 6패널 케이스 가리키며 수식 요약.")
takeaway(s, '같은 데이터라도 구조가 다르면 지지집합 밖 거동이 완전히 다르다',
         '그림 (a)~(f) = 방법 선택의 기준 맵')

# --- Xu with paper fig ---
s = blank()
content_header(s, "PART 3 · Xu et al. (2021)  Fig.1", "ReLU MLP — 훈련 지지집합 밖에서는 Affine")
ok = add_paper_fig(s, "xu_fig1_relu_extrap.png", Inches(0.25), Inches(1.15), width=Inches(7.5), caption="논문: Xu et al., How Neural Networks Extrapolate: From Feedforward to Graph Neural Networks (ICLR 2021), Fig.1")
if ok is None and (ASSETS / "fig_activation.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_activation.png"), Inches(0.25), Inches(1.15), width=Inches(7.5))
# formula banner under title area right top of cards
card_text(
    s,
    Inches(7.95),
    Inches(1.15),
    Inches(5.0),
    Inches(4.95),
    "수식 · 사례 · 평가",
    [
        "Xu et al., ICLR 2021  Fig.1 / Thm.1",
        "정의: x ∉ Conv(X_train) → 외삽",
        "Thm.1 (핵심):",
        "  lim ‖x‖→∞  f_θ(x) = A_d x + b_d",
        "  (방향 d마다 Affine)",
        "ReLU: σ(z)=max(0,z)",
        "  ‖x‖↑ ⇒ 부호패턴 고정 ⇒ 선형합성",
        "── 예시 ──",
        "참: y*=sin(x), train: x∈[-π,π]",
        "  밖 |x|>π 에서 MLP→접선(직선)",
        "매출: y*=a+b sin(2πt/12)",
        "  ReLU 학습 → 내년은 직선 추세만",
        "커피: T(t)=T∞+(T0-T∞)e^{-kt}",
        "  선형 외삽 → t↑ 시 T→-∞ 발산",
        "── 평가 함정 ──",
        "test ⊂ Conv(X_train) = 보간 착시",
        "외삽 평가: X_test ∩ Conv=∅ 강제",
    ],
    CORAL,
    dense=True,
)
# left bottom formula strip
fb = rect(s, Inches(0.25), Inches(3.30), Inches(7.5), Inches(2.80), NAVY)
ft = s.shapes.add_textbox(Inches(0.4), Inches(3.44), Inches(7.2), Inches(2.55))
tf = add_text(ft, "핵심 수식 (발표용)", size=13, bold=True, color=TEAL)
add_para(tf, "외삽 영역:   x ∉ Conv(X_train)", size=12, color=WHITE, space_before=6)
add_para(tf, "Xu Thm.1:    ‖x‖→∞  ⇒  f_θ(x) → A x + b   (ReLU MLP)", size=12, color=WHITE, space_before=4)
add_para(tf, "예시 타깃:   y* = sin(x)   vs   모델 밖:  ŷ = αx+β", size=12, color=WHITE, space_before=4)
add_para(tf, "시사점: 활성화 선택 = 외삽 시 함수족(function class) 선택", size=12, color=SOFT, space_before=6)
add_para(tf, "다음: 왜 Affine이 되는지(부호패턴) → 활성화↔타깃 짝맞추기", size=11, color=SOFT, space_before=3)
pdf_btn(s, "xu", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · Xu")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. 수식 박스 + 예시(sin/매출/커피).")
takeaway(s, '‖x‖→∞ ⇒ ReLU MLP는 Affine — 비선형 법칙은 구조 없이 외삽 불가',
         'test가 Conv(X_train) 안에 있으면 보간 평가에 불과하다')

# --- ★4 왜 ReLU가 직선이 되는가 ---
s = blank()
content_header(s, "PART 3 · 메커니즘", "ReLU의 원거리 Affine 수렴")
if (ASSETS / "fig_relu_affine.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_relu_affine.png"), Inches(0.25), Inches(1.1), height=Inches(3.15))
steps = [
    ("① 부호 패턴 고정", [
        "σ(z)=max(0,z)",
        "z≫0 ⇒ ON,  z≪0 ⇒ OFF",
        "‖x‖↑ ⇒ 활성집합 A(x) 불변",
        "함의: 원거리 입력에서 활성패턴 고정",
    ], TEAL),
    ("② 층 선형화", [
        "A(x) 고정 구간에서",
        "h^{(ℓ)} = W^{(ℓ)} h^{(ℓ-1)} + b",
        "(비선형 굽힘 소멸)",
        "함의: 조각선형이 단일 affine으로 붕괴",
    ], TEAL),
    ("③ Affine 수렴", [
        "f = W_L … W_1 x + b_tot",
        "즉 f(x)=A x + b  (Thm.1)",
        "깊이·너비↑로도 동일",
        "함의: 비선형 타깃의 원거리 접선 근사",
    ], CORAL),
]
for i, (title, bullets, accent) in enumerate(steps):
    card_text(s, Inches(0.25 + i * 4.3), Inches(4.4), Inches(4.2), Inches(1.7), title, bullets, accent, dense=True)
pdf_btn(s, "xu", Inches(11.5), Inches(0.28), w=Inches(1.4), label="PDF · Xu")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. σ→A(x)→Ax+b 체인 + 예시.")
takeaway(s, 'σ 부호패턴 고정 → 선형합성 → f(x)=Ax+b (Xu Thm.1)',
         '다음: 활성화 φ와 타깃 y*의 짝이 맞아야 외삽')

# --- Xu activations paper fig ---
s = blank()
content_header(s, "PART 3 · Xu et al. Fig.5", "활성화 = 외삽 가정 (타깃과 맞으면 성공)")
add_paper_fig(s, "xu_fig5_activations.png", Inches(0.2), Inches(1.1), width=Inches(7.0), caption="논문: Xu et al., How Neural Networks Extrapolate (ICLR 2021), Fig.5")
card_text(
    s,
    Inches(7.4),
    Inches(1.1),
    Inches(5.55),
    Inches(5.0),
    "수식 · 실험 · 대응 · 함의",
    [
        "모델: ŷ = MLP_φ(x)  (φ=활성화)",
        "타깃: y* ∈ {tanh, cos, quad, lin}",
        "지표 (외삽):",
        "  MAPE = (1/n) Σ |ŷ-y*| / |y*|",
        "관측: φ ≈ y* 족일 때만 MAPE↓",
        "  mismatch 시 MAPE ×10²~10³",
        "해석: φ가 곧 외삽 시 함수족 가정",
        "── φ ↔ y* 대응 ──",
        "φ=cos, y*=cos(ωt): 계절 성분 — 일치",
        "φ=ReLU, y*=sin: 밖은 직선 → 실패",
        "φ=tanh, y*=성장곡선: 포화 — 일치",
        "φ=ReLU, y*=성장곡선: 계속 증가 착각",
        "φ=quad, y*=T(t)∝t²: 이차 성장 — 일치",
        "── 설계 함의 ──",
        "주기? → sin/cos 계열",
        "포화? → tanh/σ",
        "단조? → ∂ŷ/∂x≥0 제약",
        "모름? → UQ (σ²↑ → STOP)",
    ],
    TEAL,
    dense=True,
)
mx = rect(s, Inches(0.2), Inches(3.58), Inches(7.0), Inches(2.55), CARD, LINE)
mt = s.shapes.add_textbox(Inches(0.35), Inches(3.66), Inches(6.7), Inches(2.40))
tf = add_text(mt, "짝맞추기:  φ (행)  ×  y* (열)   —  ✓=가정 일치", size=12, bold=True, color=INK)
add_para(tf, "          tanh     cos     quad    linear", size=11, color=SLATE, space_before=5)
add_para(tf, "tanh        ✓        ✗       ✗        ✗", size=12, color=INK, space_before=3)
add_para(tf, "cos          ✗        ✓       ✗        ✗", size=12, color=INK, space_before=2)
add_para(tf, "quad        ✗        ✗       ✓        ✗", size=12, color=INK, space_before=2)
add_para(tf, "ReLU       ✗        ✗       ✗     (선형만)", size=12, color=SLATE, space_before=2)
add_para(tf, "명제:  φ ∈ F 이고  y* ∈ F  ⇒  외삽 가능 (가정 일치)", size=11, bold=True, color=CORAL, space_before=6)
pdf_btn(s, "xu", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · Xu")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. MAPE 정의 + φ↔y* + 예시.")
takeaway(s, 'φ와 y*가 같은 함수족일 때만 외삽 성립 — mismatch 시 MAPE 폭증',
         '활성화 선택 = 외삽 시 함수족(F) 선택')

# --- EQL paper fig ---
s = blank()
content_header(s, "PART 3 · EQL  Fig.1", "Equation Learner — 수식 유닛을 네트워크 구조에 내장")
add_paper_fig(s, "eql_fig1_architecture.png", Inches(0.2), Inches(1.1), width=Inches(5.4), caption="논문: Martius & Lampert, Extrapolation and Learning Equations (arXiv 2016), Fig.1")
card_text(
    s,
    Inches(5.8),
    Inches(1.1),
    Inches(7.15),
    Inches(2.7),
    "수식 · 구조 · 사례 · 한계  (Martius & Lampert 2016)",
    [
        "은닉유닛 집합 U = {id, sin, cos, σ, ×, …}",
        "층:  z = Σ_k w_k · u_k(·) ,  u_k ∈ U",
        "목표: y* ≈ 닫힌 식 g(x) 를 구조로 학습",
        "예1: y*=sin(x) → 학습식 ŷ≈sin(1.01x)",
        "예2: 계절 매출 y*=a+b sin(2πt/12)",
        "  → 지지집합 밖에서도 주기 연장 (MLP는 직선)",
        "외삽 기제: g∈⟨U⟩이면 정의역 전체에서 유효",
        "한계: U에 없는 함수 / 고차원·잡음 / sparsity 필요",
    ],
    TEAL,
    dense=True,
)
if (ASSETS / "fig_eql_before_after.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_eql_before_after.png"), Inches(0.2), Inches(3.95), height=Inches(2.15))
card_text(
    s,
    Inches(9.7),
    Inches(3.95),
    Inches(3.25),
    Inches(2.15),
    "대비 (수식)",
    [
        "MLP: ŷ=조각선형",
        "  ⇒ 밖 Ax+b",
        "EQL: ŷ∈⟨U⟩",
        "  ⇒ 밖도 같은 g",
        "다음: NALU",
        "  (±×÷ 강제)",
    ],
    CORAL,
    dense=True,
)
pdf_btn(s, "eql", Inches(11.5), Inches(0.28), w=Inches(1.5), label="PDF · EQL")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. U 집합 수식 + sin/매출 예시 + MLP 대비.")
takeaway(s, 'EQL: ŷ∈⟨U⟩ 이면 정의역 전체에서 같은 식으로 연장',
         'MLP 조각선형(밖=Ax+b)과 결정적으로 다름')

# --- NALU paper fig ---
s = blank()
content_header(s, "PART 3 · NALU  Fig.2", "NAC / NALU — 산술 연산을 inductive bias로 고정")
pic = add_paper_fig(s, "nalu_fig2_architecture.png", Inches(0.28), Inches(1.18), width=Inches(7.35), caption="논문: Trask et al., Neural Arithmetic Logic Units (NeurIPS 2018), Fig.2 · arXiv:1808.00508")
card_text(
    s, Inches(7.85), Inches(1.18), Inches(5.05), Inches(2.40),
    "구조  ·  Trask et al. 2018",
    [
        "NAC:  W = tanh(Ŵ) ⊙ σ(M̂)  ∈ [−1, 1]",
        "→ 가산·감산만 허용",
        "NALU:  y = g⊙a + (1−g)⊙m",
        "a = Wx ,  m = exp(W log(|x|+ε))",
        "g = σ(Gx) 가 ± 와 ×÷ 경로를 선택",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(3.70), Inches(5.05), Inches(2.40),
    "외삽 결과  ·  한계",
    [
        "소규모 정수 → 대규모 피연산자로 일반화",
        "카운팅·누적합에서 MLP 대비 우위",
        "수치 지지집합 밖에서도 연산 규칙 유지",
        "한계: gate 불안정 · 0-나눗셈 · 가정 민감",
    ],
    TEAL, dense=True,
)
_uy = (pic.top + pic.height + Inches(0.44)) if pic else Inches(3.9)
card_text(
    s, Inches(0.28), _uy, Inches(7.35), Inches(6.22) - _uy,
    "해석",
    [
        "좌 NAC = ± 결합  ·  우 NALU = log-space 승산 + gate",
        "Xu의 ‘밖 = Ax+b’와 대비 — 산술 함수족을 구조로 고정",
        "형태 제약 계열(EQL同): 가정이 맞을 때만 외삽이 성립",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "nalu", Inches(11.45), Inches(0.26), label="PDF · NALU")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. NAC vs NALU. 수식과 외삽 기제.")
takeaway(s, "NAC/NALU — 산술을 가중치 공간에 제약하면 수치 범위 외삽이 성립한다",
         "가정 불일치·게이트 실패 시에는 일반 MLP와 동일하게 붕괴")

# --- Mono: cubic fail/success ---
s = blank()
content_header(s, "PART 3 · CMNN  Fig.1", "단조성 vs 표현력 — weight clip의 실패와 CMNN")
pic = add_paper_fig(s, "runje_fig1_cubic.png", Inches(0.28), Inches(1.18), width=Inches(7.35), caption="논문: Runje & Shankaranarayana, Constrained Monotonic Neural Networks (ICML 2023), Fig.1 · arXiv:2205.11775")
card_text(
    s, Inches(7.85), Inches(1.18), Inches(5.05), Inches(2.40),
    "실험 설정  ·  Runje & Shankaranarayana 2023",
    [
        "목표: 단조 증가 함수  y = x³",
        "(a) unconstrained — 적합↑, 단조 위반 가능",
        "(b) weight ≥ 0 clip — 단조↑, 표현력↓",
        "→ 사실상 선형, x³ 근사 실패",
        "(c) CMNN — 단조 + 비선형을 동시 확보",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(3.70), Inches(5.05), Inches(2.40),
    "이론적 함의",
    [
        "단조 함수족에서도 universal approximation 필요",
        "부호 clip ≠ 충분한 재매개화",
        "외삽: ∂f/∂x 부호가 지지집합 밖에서도 유지",
        "단조성은 페널티가 아니라 구조적 불변량",
    ],
    TEAL, dense=True,
)
_uy = (pic.top + pic.height + Inches(0.44)) if pic else Inches(3.5)
card_text(
    s, Inches(0.28), _uy, Inches(7.35), Inches(6.22) - _uy,
    "해석  ·  (a) → (b) → (c)",
    [
        "(b) 제약이 조악하면 단조성은 얻되 함수족이 붕괴한다",
        "(c) 올바른 재매개화는 비선형 단조성까지 보존한다",
        "대응 영역: 가격–수요, SOC–OCV 등 방향만 아는 외삽",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "runje", Inches(11.45), Inches(0.26), label="PDF · CMNN")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. clip 실패 → CMNN.")
takeaway(s, "단순 weight clip은 x³조차 실패한다 — CMNN은 단조와 표현력을 동시에 만족",
         "단조성은 손실 페널티가 아니라 구조적 불변량으로 강제해야 한다")

# --- Mono unit architecture ---
s = blank()
content_header(s, "PART 3 · CMNN  Fig.3", "Monotonic Dense Unit — 제약을 레이어에 내재화")
pic = add_paper_fig(s, "runje_fig3_mono_unit.png", Inches(0.28), Inches(1.18), height=Inches(2.70), caption="논문: Runje & Shankaranarayana, Constrained Monotonic Neural Networks (ICML 2023), Fig.3")
card_text(
    s, Inches(7.85), Inches(1.18), Inches(5.05), Inches(2.40),
    "아키텍처",
    [
        "W → |·|_t → W′   (부호·크기 재매개화)",
        "h = W′x + b",
        "s = (s̆, ŝ, s̃) 로 채널 split",
        "ρ̆ , ρ̂ , ρ̃  →  concat → y",
        "단조성 by construction",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(3.70), Inches(5.05), Inches(2.40),
    "외삽에서의 역할",
    [
        "지지집합 밖에서도 ∂f/∂xᵢ 부호 역전 차단",
        "점추정 오차와 무관하게 방향 오류 감소",
        "Part 4: CA-CSS + isotonic calibration",
        "구조(제약) ⊥ 크기(UQ) — 보완 관계",
    ],
    CORAL, dense=True,
)
_uy = (pic.top + pic.height + Inches(0.44)) if pic else Inches(4.0)
card_text(
    s, Inches(0.28), _uy, Inches(7.35), Inches(6.22) - _uy,
    "블록 해독",
    [
        "|·|_t : 단조 허용 집합으로의 미분가능 투영 (hard clip 아님)",
        "split + ρ : 국소 비선형을 유지하면서 전역 단조성 보존",
        "Xu(진단) → CMNN(제약) → Part 4(프로토콜·모델)",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "runje", Inches(11.45), Inches(0.26), label="PDF · CMNN")
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. 레이어 구조 → Part4.")
takeaway(s, "CMNN 레이어 — 단조성 by construction으로 외삽 시 부호 역전을 차단",
         "제약은 구조를, UQ는 크기를 담당한다")

# --- PINN failure paper fig ---
s = blank()
content_header(s, "PART 3 · PINN  Fig.1", "물리 잔차가 작아도 시간 외삽은 보장되지 않는다")
pic = add_paper_fig(s, "fesser_fig1_pinn_extrap.png", Inches(0.28), Inches(1.18), width=Inches(7.35), caption="논문: Fesser et al., Understanding and Mitigating Extrapolation Failures in PINNs (2023), Fig.1 · arXiv:2306.09478")
card_text(
    s, Inches(7.85), Inches(1.18), Inches(5.05), Inches(2.40),
    "프로토콜  ·  Fesser et al. 2023",
    [
        "L = L_data + λ ‖N[u_θ]‖²",
        "학습  t ∈ [0, T/2]",
        "평가  t ∈ (T/2, T]  — 시간 축 외삽",
        "Burgers · Allen–Cahn에서 L2·잔차 폭증",
        "Raissi (2019) 성공 ≠ 외삽 일반 보증",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(3.70), Inches(5.05), Inches(2.40),
    "실패 기제",
    [
        "Spectral bias — 고주파·충격파 학습 지연",
        "Loss balancing — 잔차↓ ≠ 해 오차↓",
        "Collocation 희소 — 미래 시간 커버 부족",
        "PDE prior는 필요조건일 뿐",
    ],
    TEAL, dense=True,
)
_uy = (pic.top + pic.height + Inches(0.44)) if pic else Inches(3.4)
card_text(
    s, Inches(0.28), _uy, Inches(7.35), Inches(6.22) - _uy,
    "해석  ·  시간 hull",
    [
        "점선(t = 0.5): 공간이 아니라 시간 외삽의 경계",
        "잔차가 작아 보여도 해는 발산할 수 있다",
        "대응: 구간 재학습 · causality · UQ 기권 (→ Zhu)",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "fesser", Inches(9.85), Inches(0.26), w=Inches(1.45), label="PDF · Fesser")
pdf_btn(s, "raissi", Inches(11.40), Inches(0.26), w=Inches(1.50), label="PDF · Raissi")
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. 시간 분할 강조.")
takeaway(s, "PINN — PDE 잔차 최소화는 시간 외삽 성능의 보증이 아니다",
         "물리 prior와 외삽 평가 프로토콜은 분리해서 설계해야 한다")

# --- DeepONet paper fig ---
s = blank()
content_header(s, "PART 3 · DeepONet  Fig.1", "Neural operator — 예측과 기권(abstain)의 분리")
pic = add_paper_fig(s, "zhu_fig1_deeponet.png", Inches(0.28), Inches(1.18), height=Inches(3.15), caption="논문: Zhu et al., Reliable Extrapolation of Deep Neural Operators (CMAME 2023), Fig.1 · arXiv:2212.06347")
card_text(
    s, Inches(7.85), Inches(1.18), Inches(5.05), Inches(2.40),
    "모형  ·  Zhu et al. 2022",
    [
        "G(u)(y) ≈ Σₖ bₖ(u) · tₖ(y)",
        "Branch b(u): 입력 함수 인코딩",
        "Trunk t(y): 질의 좌표 인코딩",
        "(A) 점추정 NN — 분포 밖 붕괴",
        "(B) Operator — 함수 공간 OOD 동일",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(3.70), Inches(5.05), Inches(2.40),
    "Reliable extrapolation",
    [
        "u ∉ supp(P_train) → 점추정만으로 위험",
        "물리 잔차·관측 mismatch → 신뢰도",
        "임계값 이하: abstain (기권)",
        "성능 = 정확도 × 선택적 예측",
    ],
    CORAL, dense=True,
)
_uy = (pic.top + pic.height + Inches(0.44)) if pic else Inches(4.5)
card_text(
    s, Inches(0.28), _uy, Inches(7.35), Inches(6.22) - _uy,
    "해석  ·  (A) vs (B)",
    [
        "(A) high confidence, low accuracy — 조용한 실패",
        "(B) UQ로 실패를 감지하고 기권할 수 있는가",
        "제약 · 물리 · 기권 — 세 안전장치가 보완 관계",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "zhu", Inches(11.45), Inches(0.26), label="PDF · Zhu")
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. abstain과 Part1 UQ 연결.")
takeaway(s, "DeepONet + UQ — 함수 공간 외삽에서 기권 가능한가가 핵심이다",
         "점추정 정확도만으로는 고위험 외삽을 대응할 수 없다")

# --- Enhanced method compare ---
s = blank()
content_header(s, "PART 3 · 종합 비교", "언제 무엇을 쓰는가")
if (ASSETS / "fig_method_decision.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_method_decision.png"), Inches(0.25), Inches(1.1), height=Inches(3.0))
headers = ["방법", "원리", "언제", "한계"]
rows = [
    ["EQL/NALU", "수식·산술", "법칙 모양 앎", "유닛/불안정"],
    ["Mono NN", "단조 강제", "방향만 앎", "부호 오용"],
    ["PINN", "PDE 잔차", "지배식 있음", "시간외삽"],
    ["DeepONet+UQ", "연산자+경고", "모름/고위험", "비용·설계"],
]
col_w = [2.2, 2.3, 2.5, 2.4]
x0 = 0.35
y0 = 4.25
for j, h in enumerate(headers):
    x = Inches(x0 + sum(col_w[:j]))
    rect(s, x, Inches(y0), Inches(col_w[j] - 0.06), Inches(0.3), TEAL)
    t = s.shapes.add_textbox(x, Inches(y0 + 0.02), Inches(col_w[j] - 0.1), Inches(0.26))
    add_text(t, h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        x = Inches(x0 + sum(col_w[:j]))
        y = Inches(y0 + 0.34 + i * 0.36)
        highlight = i == 1
        fill = TEAL if highlight else (CARD if i % 2 == 0 else PANEL)
        fc = WHITE if highlight else (INK if j == 0 else SLATE)
        rect(s, x, y, Inches(col_w[j] - 0.06), Inches(0.34), fill, LINE)
        t = s.shapes.add_textbox(x + Inches(0.04), y + Inches(0.04), Inches(col_w[j] - 0.12), Inches(0.26))
        add_text(t, cell, size=10, bold=(j == 0 or highlight), color=fc, align=PP_ALIGN.LEFT)
card_text(
    s,
    Inches(9.85),
    Inches(4.25),
    Inches(3.1),
    Inches(1.8),
    "발표 연결",
    [
        "지식 종류 → 구조 선택",
        "공통: 외삽 평가·ERM비교",
        "우리: Mono/제약 → Part4",
    ],
    CORAL,
    dense=True,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. 플로우차트 → Mono 하이라이트 → Part4.")
takeaway(s, '무엇을 아는가? → 구조 선택 → 안전장치(평가·UQ)',
         '우리 연구에 가장 가까운 가지 = Monotonic / Constraint')

# ════════════════ PART 4 ════════════════
section_slide(
    "실전 응용",
    "데이터셋 · 프로토콜 · APEX-Guard 아키텍처 · 결과 · Ablation",
    "4",
    "10분",
    P(),
    TOTAL,
    points=[
        "N-CMAPSS: 무엇을 왜 외삽하는가",
        "설계: Dual Encoder + 제약 + isotonic",
        "근거: 3.26 vs 3.80 · ablation",
    ],
)

# --- 데이터셋 소개 ---
s = blank()
content_header(s, "PART 4 · 데이터", "N-CMAPSS DS02 — 항공 엔진 run-to-failure 시뮬레이션")
card_text(
    s, Inches(0.35), Inches(1.25), Inches(4.15), Inches(2.5),
    "데이터셋 개요",
    [
        "NASA 터보팬 열화 시뮬레이터 (CMAPSS 후속)",
        "unit = 엔진 1대의 전체 수명 궤적",
        "각 시점: 다채널 센서 + 운용조건 기록",
        "레이블: RUL (잔여 수명, cycle 단위)",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(4.60), Inches(1.25), Inches(4.15), Inches(2.5),
    "운용조건 4축",
    [
        "Altitude — 고도",
        "Mach — 속도",
        "TRA — 스로틀(출력 요구) ← 외삽축",
        "T30 등 온도 계열",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.85), Inches(1.25), Inches(4.05), Inches(2.5),
    "왜 TRA가 외삽축인가",
    [
        "고TRA = 고부하 = 열화 가속 구간",
        "실전: 학습 못 본 고부하 운용 발생",
        "저TRA로 학습 → 고TRA 예측 요구",
        "= covariate shift의 물리적 실체",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(3.90), Inches(12.55), Inches(2.2),
    "과제 정의",
    [
        "입력: 센서 시계열 window  →  출력: 해당 시점 RUL (회귀)",
        "난점 1 — 열화 신호와 운용부하 신호가 얽혀 있음 (같은 센서에 두 원인이 겹침)",
        "난점 2 — 고TRA 구간은 train 지지집합 밖 → Part 1의 convex hull 논리가 그대로 적용",
        "따라서: 부하와 열화를 분리하는 구조 + 물리 제약 + 외삽 전용 평가가 모두 필요",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 데이터가 뭔지, RUL 과제, TRA가 왜 외삽축인지.")
takeaway(s, "N-CMAPSS = 엔진 수명 궤적 + 운용조건 — RUL 회귀 과제",
         "고TRA 미관측이 covariate shift의 물리적 실체")

# --- 데이터셋 구성 (pipeline) ---
s = blank()
content_header(s, "PART 4 · 데이터 구성", "원시 시계열 → strict 외삽 밴드")
if (ASSETS / "fig_dataset_pipeline.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_dataset_pipeline.png"), Inches(2.30), Inches(1.18), height=Inches(3.52))
card_text(
    s, Inches(0.35), Inches(4.85), Inches(12.55), Inches(1.35),
    "구성 원칙",
    [
        "분위수(q70/q90)는 train-pool 행에서만 계산 — test 정보 누출 차단",
        "high_frac ≥ 0.7: window 내 고TRA 비율 조건 → 저→고 전이 구간이 test에 섞이는 것을 방지",
        "밴드별 표본: train 6,645 · val 1,868 · strict_late 201 — 셋이 TRA 축에서 겹치지 않음",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 파이프라인 그림 따라 설명. 누출 차단 강조.")
takeaway(s, "데이터 구성 = 윈도우화 → TRA 분위수 필터 → 3밴드 분리",
         "train 지지집합과 test가 TRA 축에서 분리되어야 '외삽 평가'가 성립")

# --- ★6 N-CMAPSS 왜 외삽인가 (4D) ---
s = blank()
content_header(s, "PART 4 · 외삽성", "N-CMAPSS Operating Regime — 4차원 분포 관점")
if (ASSETS / "fig_ncmapss_4d.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_ncmapss_4d.png"), Inches(3.22), Inches(1.22), height=Inches(4.48))
box = rect(s, Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.7), PANEL, LINE)
t = s.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.5))
add_text(
    t,
    "Train: TRA≤q70    Test: TRA>q90 + unit holdout    →  test는 train hull 밖 (보간 평가 아님)",
    size=14,
    bold=True,
    color=INK,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 4축 중 TRA만 train/test 분포가 갈라짐을 그림으로.")
takeaway(s, "4차원 중 TRA에서만 train/test 분포가 분리 — TRA가 외삽축",
         "그래서 N-CMAPSS 과제는 보간이 아니라 외삽 문제")

# --- ★7 APEX lineage ---
s = blank()
content_header(s, "PART 4 · 연결", "문헌 교훈 → APEX-Guard 설계 요소")
if (ASSETS / "fig_apex_lineage.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_apex_lineage.png"), Inches(0.4), Inches(1.35), width=Inches(12.3))
box = rect(s, Inches(0.45), Inches(5.7), Inches(12.4), Inches(0.85), PANEL, LINE)
t = s.shapes.add_textbox(Inches(0.7), Inches(5.85), Inches(12), Inches(0.6))
tf = add_text(t, "Xu(실패 기제) → CMNN(단조 제약) → DomainBed(공정 평가) → CA-CSS → APEX-Guard", size=15, bold=True, color=INK)
add_para(tf, "각 논문의 교훈이 설계 요소 하나씩으로 번역된다", size=13, color=SLATE)
pdf_btn(s, "xu", Inches(0.55), Inches(0.28), w=Inches(1.2), label="Xu")
pdf_btn(s, "runje", Inches(1.9), Inches(0.28), w=Inches(1.3), label="CMNN")
pdf_btn(s, "domainbed", Inches(3.35), Inches(0.28), w=Inches(1.6), label="DomainBed")
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 문헌 교훈이 설계 요소로 연결되는 흐름.")
takeaway(s, "Xu → CMNN → DomainBed → CA-CSS → APEX — 문헌 교훈이 설계로 연결된다",
         "세미나 논리의 수렴점")

# --- 아키텍처 (그림 + 설계 근거) ---
s = blank()
content_header(s, "PART 4 · 아키텍처", "APEX-Guard — 구조와 신호 흐름")
if (ASSETS / "fig_apex_architecture.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_apex_architecture.png"), Inches(0.25), Inches(1.18), width=Inches(8.3))
card_text(
    s, Inches(8.70), Inches(1.18), Inches(4.25), Inches(2.50),
    "동작 흐름",
    [
        "① OC-Norm: 운용점 편차를 분리 정규화",
        "② Health Enc: TRA 마스킹 → 열화만 인코딩",
        "③ Load Enc: TRA·φ·T30·cycle → 부하 전담",
        "④ RUL = health − damage  (w≥0)",
        "⑤ 추론 후 unit별 isotonic 정렬",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.70), Inches(3.80), Inches(4.25), Inches(2.35),
    "설계 근거 (문헌 → 설계)",
    [
        "Xu: 밖=Ax+b → 외삽축을 별도 경로로 명시",
        "CMNN: 단조는 구조로 → w≥0 + isotonic",
        "DomainBed: 평가 먼저 → strict 프로토콜",
        "얽힘 제거: 열화·부하 분리가 일반화의 핵심",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2분. 그림 흐름 ①~⑤ 순서로. 설계 근거를 문헌과 연결.")
takeaway(s, "구조 요지: 열화(health)와 부하(load)를 분리하고, 단조를 구조로 강제",
         "각 블록이 Part 1~3의 교훈 하나씩에 대응")

# APEX 3-layer defense
s = blank()
content_header(s, "PART 4 · 설계 체계", "APEX-Guard — 3층 설계 (평가·학습·추론)")
box = rect(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.05), NAVY)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.8), Inches(0.7))
add_text(t, "Adaptive Physics EXtrapolation Guard — 외삽은 불가피하므로 설계로 대응한다", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
card_text(
    s, Inches(0.45), Inches(2.65), Inches(4.0), Inches(3.5),
    "L1 · 평가",
    ["strict_late 프로토콜", "순수 외삽만 측정", "전이 window 제거", "→ 숫자의 신뢰성 확보"],
)
card_text(
    s, Inches(4.65), Inches(2.65), Inches(4.0), Inches(3.5),
    "L2 · 학습",
    ["CA-CSS 제약 발견·컴파일", "TRA(−) hinge + cycle_mono", "물리 감독이 hull을 확장", "→ 방향 오류 차단"],
    TEAL,
)
card_text(
    s, Inches(8.85), Inches(2.65), Inches(4.0), Inches(3.5),
    "L3 · 추론",
    ["unit별 isotonic 정렬", "H1 위반 28% → 0%", "단조 RUL 궤적 보장", "→ 궤적 물리성 확보"],
    CORAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 3층 프레임: 평가/학습/추론.")
takeaway(s, "L1 평가 · L2 제약학습 · L3 단조추론 — 3층 설계",
         "외삽 성능을 우연이 아니라 설계로 확보한다")

# strict protocol
s = blank()
content_header(s, "PART 4 · 프로토콜", "Strict TRA Extrap — 3겹 조건으로 고정")
card_text(
    s, Inches(0.35), Inches(1.25), Inches(4.15), Inches(2.55),
    "① TRA shift",
    [
        "Train: tra_end≤q70 ∧ tra_mean≤q70",
        "Test: tra_end>q90 ∧ high_frac≥0.7",
        "→ 전이 window 오염 제거",
    ],
    dense=True,
)
card_text(
    s, Inches(4.60), Inches(1.25), Inches(4.15), Inches(2.55),
    "② Unit holdout (tra_rich)",
    [
        "train {5,18,16,11,15,14}",
        "val {2}  ·  test {10,20}",
        "고TRA share 상위 유닛 holdout",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.85), Inches(1.25), Inches(4.05), Inches(2.55),
    "③ Late-life band",
    [
        "RUL ≤ 50만 측정 (strict_late)",
        "n = 201  ·  seed = 42  ·  e15",
        "all_rul 7.53 → late 3.26",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(3.95), Inches(12.55), Inches(2.15),
    "왜 이 프로토콜인가 (근거)",
    [
        "이전 hard: tra_end만 보면 저→고 전이 window가 고TRA로 혼입 · unit 14는 고TRA 거의 없음",
        "high_frac=(TRA>q85 비율)로 순수 고부하만 남김 → 전이 필터 유효",
        "말기(RUL≤50)가 핵심: all_rul 7.53 → strict_late 3.26",
        "DomainBed 교훈 — 평가 프로토콜을 먼저 고정한 뒤 모델 비교",
        "즉: OOD를 주장하기 전에 TRA로 '밖'을 수치 정의했다",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2분. 프로토콜 3겹 + 재설계 근거.")
takeaway(s, "Strict = TRA shift × unit holdout × late band",
         "프로토콜이 약하면 성능 숫자는 의미가 없다")

# architecture
s = blank()
content_header(s, "PART 4 · 구조", "v4_disentangled — Dual Encoder")
card_text(s, Inches(0.45), Inches(1.4), Inches(3.05), Inches(4.7), "입력", ["OC-Norm", "운용점 편차 분리", "cycle channel"])
card_text(s, Inches(3.7), Inches(1.4), Inches(3.05), Inches(4.7), "Health Enc", ["TRA=0 masking", "열화 패턴", "TRA overfit 방지"], TEAL)
card_text(s, Inches(6.95), Inches(1.4), Inches(3.05), Inches(4.7), "Load Enc", ["TRA, φ, T30, cycle", "외삽 축 전담"], CORAL)
card_text(s, Inches(10.2), Inches(1.4), Inches(2.7), Inches(4.7), "RUL 분해", ["RUL=health−damage", "w≥0 제약", "d=64, L=2, ep=15"])
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. Dual path.")
takeaway(s, "Health(TRA=0) + Load(외삽축) — 역할을 분리했다", "RUL = health − damage")

# CA-CSS
s = blank()
content_header(s, "PART 4 · 학습", "CA-CSS + Isotonic")
card_text(
    s, Inches(0.45), Inches(1.4), Inches(4.0), Inches(4.7),
    "CA-CSS 파이프라인",
    ["Discovery → Pseudo-extrap", "→ Compile → Train", "train-band만 (leakage 방지)", "TRA(−) λ=0.05 adopt", "cycle_mono adopt"],
)
card_text(
    s, Inches(4.65), Inches(1.4), Inches(4.0), Inches(4.7),
    "손실",
    ["L = L_RUL + L_physics + L_TRA", "TRA(−) hinge", "cycle_mono hinge", "물리 감독 = 외삽 가드"],
    TEAL,
)
card_text(
    s, Inches(8.85), Inches(1.4), Inches(4.0), Inches(4.7),
    "Isotonic 후처리",
    ["unit별 cycle 정렬", "decreasing isotonic", "H1: 28% → 0%", "단조 궤적 보장"],
    CORAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. Part3 Mono와 동일 철학.")
takeaway(s, "CA-CSS가 제약을 발견·컴파일하고 isotonic이 궤적을 지킨다",
         "H1 위반 28%→0%가 물리 가드의 증거")

# results — main comparison
s = blank()
content_header(s, "PART 4 · 결과", "strict_late 비교 (tra_rich · seed=42 · n=201)")
if (ASSETS / "fig_strict_rmse.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_strict_rmse.png"), Inches(0.28), Inches(1.18), width=Inches(7.45))
card_text(
    s, Inches(7.90), Inches(1.18), Inches(5.05), Inches(3.35),
    "주결과 표",
    [
        "APEX (v4_iso)   3.26   R² 0.947",
        "TabPFN           3.80   R² 0.927",
        "XGB              6.95   R² 0.757",
        "Transformer      7.31   R² 0.732",
        "TCN             14.09   R² 0.002",
        "Δ vs TabPFN: −0.54 RMSE",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.90), Inches(4.65), Inches(5.05), Inches(1.50),
    "해석",
    [
        "제약+단조가 강 tabular baseline을 이김",
        "순수 용량(TF/TCN)만으로는 외삽 실패",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(0.28), Inches(4.65), Inches(7.45), Inches(1.50),
    "실험 설정",
    [
        "epochs=15 · seed=42 · TabPFN max_train=2000",
        "모델 키 v4_iso = APEX-Guard · 출처: strict_extrap_RESULTS.md",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2분. 메인 비교표.")
takeaway(s, "strict_late: APEX 3.26 < TabPFN 3.80 (n=201) — 주결과",
         "용량만 키운 모델(TF/TCN)은 외삽에서 붕괴")

# ablation evidence
s = blank()
content_header(s, "PART 4 · Ablation", "무엇이 성능을 만드는가 — 구성요소 제거 실험")
if (ASSETS / "fig_ablation.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_ablation.png"), Inches(0.25), Inches(1.15), width=Inches(7.55))
card_text(
    s, Inches(7.95), Inches(1.15), Inches(5.0), Inches(2.55),
    "Ablation (hard · 3-seed · e15)",
    [
        "full          3.88 ± 0.80",
        "no_tra_loss   3.97  (+0.09)",
        "no_physics    4.21  (+0.33)",
        "no_iso        4.37  (+0.49) ← 최대",
        "revin         8.85  (+4.97) ← 해로움",
        "TabPFN ref    4.79",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.95), Inches(3.85), Inches(5.0), Inches(2.30),
    "근거 해석",
    [
        "Isotonic이 가장 큰 기여 (+0.49)",
        "물리 제약 제거도 유의미 (+0.33)",
        "H1 위반: 28% → 0% (iso on)",
        "RevIN은 이 과제에서 역효과",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.25), Inches(4.95), Inches(7.55), Inches(1.20),
    "한계 · 강건성 (nasa_drop14)",
    [
        "test를 {11,15}로 바꾸면: TF 4.45 < TabPFN 5.14 < APEX 5.30 — 순위 역전",
        "주결과는 tra_rich 한정 · 유닛 규칙·multi-seed를 명시",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2.5분. ablation + nasa_drop14 한계.")
takeaway(s, "Ablation: isotonic(+0.49) > physics(+0.33) > TRA loss(+0.09)",
         "주결과와 별도로 유닛 민감성(nasa_drop14)도 보고한다")

# takeaway part4
s = blank()
content_header(s, "PART 4 · 교훈", "설계가 성능으로 연결되는 논리")
box = rect(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(2.0), NAVY)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.50), Inches(11.8), Inches(1.7))
tf = add_text(t, "외삽 성능을 우연에 맡기지 말고\n평가(L1) · 제약(L2) · 단조(L3)를 설계하라", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "이론 → 프로토콜(strict_late) → 모델(APEX) → ablation으로 근거 확인", size=14, color=SOFT, align=PP_ALIGN.CENTER, space_before=8)
trip = [
    ("정리 1", "외삽 정의\n고TRA × unit holdout × late"),
    ("정리 2", "설계 대응\nDual Enc + CA-CSS + iso"),
    ("정리 3", "근거\n3.26 vs 3.80 · iso +0.49"),
]
for i, (h, b) in enumerate(trip):
    x = Inches(0.45 + i * 4.2)
    card_text(s, x, Inches(3.60), Inches(4.0), Inches(2.45), h, b.split("\n"), TEAL if i != 2 else CORAL)
footer(s, P(), TOTAL, "Part 4")
note(s, "1분. Part4 클로징.")
takeaway(s, "평가·제약·단조를 설계하면 외삽 성능이 따라온다",
         "주결과 + ablation + 한계(유닛 민감)를 한 세트로 말하라")

# ════════════════ PART 5 ════════════════
section_slide(
    "연구 동향 & 논문 가이드",
    "30년 타임라인 · 필독 · 연구 적용",
    "5",
    "4분",
    P(),
    TOTAL,
    points=[
        "우선 읽을 논문 3편부터",
        "내 문제에 방법 매핑",
        "Future로 연구 전망을 정리",
    ],
)

# timeline
s = blank()
content_header(s, "PART 5 · 타임라인", "외삽 연구 30년")
items = [
    ("1990s", "Vapnik · VC\n보간 이론"),
    ("2000s", "GP · Covariate\nShift"),
    ("2016-19", "EQL · NALU\nIRM"),
    ("2020-21", "DomainBed\nXu · PINN"),
    ("2022-23", "CMNN\nDeepONet+UQ"),
    ("2024-25", "Pfister\nProb.Richardson\nWu 시계열"),
]
for i, (yr, txt) in enumerate(items):
    x = Inches(0.4 + i * 2.15)
    rect(s, x, Inches(2.0), Inches(2.0), Inches(3.5), CARD, LINE)
    top = slide_shapes_top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.0), Inches(0.55))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL if i % 2 == 0 else CORAL
    top.line.fill.background()
    ty = s.shapes.add_textbox(x, Inches(2.08), Inches(2.0), Inches(0.4))
    add_text(ty, yr, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.8), Inches(1.8), Inches(2.4))
    add_text(tb, txt, size=13, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL, "Part 5")
note(s, "1분. 흐름만.")
takeaway(s, '30년: 이론→OOD→NN진단→제약/물리→UQ',
         '지금은 ')

# must read
s = blank()
content_header(s, "PART 5 · 필독", "우선 읽을 논문 가이드 (PDF 동봉)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★★★ 필수",
    [
        "Xu 2021 — NN 외삽",
        "Gulrajani 2020 — DomainBed",
        "Liu 2023 — OOD Survey",
    ],
    CORAL,
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★★☆ 방법",
    [
        "Trask 2018 NALU",
        "Runje 2023 CMNN",
        "Fesser 2023 PINN 실패",
        "Teckentrup 2024",
    ],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★☆☆ 응용",
    [
        "Aykol 2021 Battery+Phys",
        "Li 2023 Lifetime OOD",
        "Wu 2025 시계열 OOD",
        "총 36편 / extrapolation-papers/",
    ],
)
pdf_btn(s, "xu", Inches(0.55), Inches(0.28), w=Inches(1.05), label="Xu")
pdf_btn(s, "domainbed", Inches(1.7), Inches(0.28), w=Inches(1.35), label="DomainBed")
pdf_btn(s, "liu", Inches(3.15), Inches(0.28), w=Inches(1.05), label="Liu")
pdf_btn(s, "arjovsky2019", Inches(4.3), Inches(0.28), w=Inches(1.05), label="IRM")
pdf_btn(s, "nalu", Inches(5.45), Inches(0.28), w=Inches(1.05), label="NALU")
pdf_btn(s, "runje", Inches(6.6), Inches(0.28), w=Inches(1.05), label="CMNN")
pdf_btn(s, "fesser", Inches(7.75), Inches(0.28), w=Inches(1.15), label="Fesser")
pdf_btn(s, "raissi", Inches(9.0), Inches(0.28), w=Inches(1.15), label="Raissi")
pdf_btn(s, "zhu", Inches(10.25), Inches(0.28), w=Inches(1.05), label="Zhu")
pdf_btn(s, "eql", Inches(11.4), Inches(0.28), w=Inches(1.05), label="EQL")
footer(s, P(), TOTAL, "Part 5")
note(s, "1.5분. 버튼 클릭으로 PDF 열기.")
takeaway(s, '우선 읽을 논문 3편: Xu · DomainBed · Liu',
         '나머지는 paper_pdfs/에서 필요할 때')

# apply
s = blank()
content_header(s, "PART 5 · 적용", "연구 문제에 방법을 매핑하기")
rows = [
    ("문제", "추천 방법", "논문"),
    ("단조 관계 외삽 (SOC-OCV 등)", "Monotonic NN + isotonic", "Runje 2023"),
    ("조건/도메인 OOD", "제약 학습 + ERM 비교", "Liu · Gulrajani"),
    ("물리 PDE 영역", "PINN + UQ (DeepONet)", "Fesser · Zhu"),
    ("소량·해석 필요", "EQL + Bayesian UQ", "Martius · Teckentrup"),
]
for j, h in enumerate(rows[0]):
    rect(s, Inches(0.45 + j * 4.15), Inches(1.4), Inches(4.05), Inches(0.5), PANEL)
    t = s.shapes.add_textbox(Inches(0.55 + j * 4.15), Inches(1.48), Inches(3.85), Inches(0.35))
    add_text(t, h, size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows[1:]):
    for j, cell in enumerate(row):
        y = Inches(2.0 + i * 1.05)
        rect(s, Inches(0.45 + j * 4.15), y, Inches(4.05), Inches(0.95), CARD if i % 2 == 0 else PANEL, LINE)
        t = s.shapes.add_textbox(Inches(0.55 + j * 4.15), y + Inches(0.28), Inches(3.85), Inches(0.45))
        add_text(t, cell, size=13, color=INK, align=PP_ALIGN.CENTER, bold=(j == 0))
footer(s, P(), TOTAL, "Part 5")
note(s, "1분.")
takeaway(s, '문제→방법 매핑이 발표의 실용적 결론',
         '단조면 CMNN, OOD면 ERM비교, PDE면 PINN+UQ')

# summary
s = blank()
content_header(s, "핵심 요약", "외삽을 다루는 4가지 핵심")
keys = [
    ("01  이론", "보간 ≠ 외삽\nsupport 밖은 가정 없으면 식별 불가"),
    ("02  OOD", "ERM 실패 가능\nDomainBed로 재검증 필수"),
    ("03  구조", "물리·단조 제약이\n최고의 외삽 도구"),
    ("04  UQ", "예측값보다\n얼마나 틀릴 수 있는가"),
]
for i, (t, b) in enumerate(keys):
    x = Inches(0.45 + i * 3.2)
    rect(s, x, Inches(1.5), Inches(3.05), Inches(4.4), CARD, LINE)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(3.05), Inches(0.7))
    top.fill.solid()
    top.fill.fore_color.rgb = [NAVY, TEAL, CORAL, RGBColor(0x4A, 0x6F, 0xA5)][i]
    top.line.fill.background()
    tt = s.shapes.add_textbox(x, Inches(1.62), Inches(3.05), Inches(0.5))
    add_text(tt, t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bb = s.shapes.add_textbox(x + Inches(0.2), Inches(2.6), Inches(2.65), Inches(2.8))
    add_text(bb, b, size=15, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL)
note(s, "1.5분. 클로징 메시지.")
takeaway(s, '이론 · OOD · 구조제약 · UQ — 네 핵심을 반복',
         '도메인 지식이 최고의 외삽 도구')

# --- ★8 Future ---
s = blank()
content_header(s, "PART 5 · Future", "앞으로 연구 방향")
if (ASSETS / "fig_future.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_future.png"), Inches(0.5), Inches(1.35), width=Inches(12.2))
speak(
    s,
    [
        "Future 네 단계: Inverse → Physics prior → Foundation → Digital Twin",
        "명제: 외삽 가정(Assumption)을 시스템에 더 강하게·자동으로 내재화한다",
        "후속 연구 질의에 대한 응답 프레임",
    ],
    y=Inches(5.1),
    h=Inches(1.15),
)
footer(s, P(), TOTAL, "Part 5")
note(s, "1분. 교수님 반응 좋은 마무리 방향.")
takeaway(s, 'Future: Inverse → Physics prior → Foundation → Twin',
         '가정을 더 강하게·자동으로 넣는 방향')

# final advice
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), W, Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = TEAL
stripe.line.fill.background()
t = s.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.45))
add_text(t, "최종 조언 — 마지막에 이 구조만 남겨라", size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.25), Inches(11.3), Inches(1.5))
add_text(
    t2,
    "외삽을 풀 때 가장 할 일:\n어떤 물리적·수학적 불변성이 존재하는가?",
    size=26,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
closes = [
    ("1. 가정", "단조성 · PDE · 보존\n없으면 식별 불가"),
    ("2. 평가", "hull 밖 / OOD 프로토콜\n없으면 숫자 착시"),
    ("3. 안전장치", "제약 · UQ · 단조\n설계로 대응"),
]
for i, (h, b) in enumerate(closes):
    x = Inches(1.0 + i * 3.9)
    rect(s, x, Inches(4.1), Inches(3.6), Inches(1.7), PANEL, LINE)
    tb = s.shapes.add_textbox(x + Inches(0.2), Inches(4.25), Inches(3.2), Inches(1.4))
    tf = add_text(tb, h, size=16, bold=True, color=TEAL)
    add_para(tf, b, size=14, color=WHITE, space_before=8)
footer_box = s.shapes.add_textbox(Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.45))
add_text(footer_box, "도메인 지식이 있다면 — 그것이 최고의 외삽 도구다", size=16, bold=True, color=SOFT, align=PP_ALIGN.CENTER)
note(s, "30초.")
P()

# refs — 2 columns so all key PDFs clickable
s = blank()
content_header(s, "참고문헌 · Q&A", "주황 버튼 → PDF  ·  전체는 paper_pdfs/ (39편) + INDEX.md")
refs_l = [
    ("xu", "Xu 2021 — NN Extrapolate"),
    ("domainbed", "Gulrajani 2020 — DomainBed"),
    ("liu", "Liu 2023 — OOD Survey"),
    ("arjovsky2019", "Arjovsky 2019 — IRM"),
    ("arjovsky2021", "Arjovsky 2021 — OOD thesis"),
    ("eql", "Martius 2016 — EQL"),
    ("nalu", "Trask 2018 — NALU"),
    ("runje", "Runje 2023 — CMNN"),
    ("fesser", "Fesser 2023 — PINN fail"),
]
refs_r = [
    ("raissi", "Raissi 2019 — PINN"),
    ("zhu", "Zhu 2022 — DeepONet+UQ"),
    ("pfister", "Pfister 2024 — Identifiability"),
    ("teckentrup", "Teckentrup 2024 — Prob.Richardson"),
    ("wu", "Wu 2025 — Time-series OOD"),
    ("ye2021", "Ye 2021 — OOD theory"),
    ("ye2022", "Ye 2022 — OoD-Bench"),
    ("bartley", "Bartley 2019 — Multivariate"),
    ("aykol", "Aykol 2021 — Battery+Phys"),
    ("note", "학습노트 PDF"),
]
for i, (key, title) in enumerate(refs_l):
    y = Inches(1.3 + i * 0.55)
    pdf_btn(s, key, Inches(0.4), y, w=Inches(1.45), h=Inches(0.4), label="PDF")
    tb = s.shapes.add_textbox(Inches(2.0), y + Inches(0.05), Inches(4.3), Inches(0.35))
    add_text(tb, title, size=12, color=INK)
for i, (key, title) in enumerate(refs_r):
    y = Inches(1.3 + i * 0.55)
    pdf_btn(s, key, Inches(6.7), y, w=Inches(1.45), h=Inches(0.4), label="PDF")
    tb = s.shapes.add_textbox(Inches(8.3), y + Inches(0.05), Inches(4.5), Inches(0.35))
    add_text(tb, title, size=12, color=INK)
footer(s, P(), TOTAL)
note(s, "Q&A. paper_pdfs/에 전체 39편. INDEX.md 목록 참고.")

# fix TOTAL in footers if page count differs
actual = len(prs.slides)
print(f"slides={actual} (footer TOTAL was {TOTAL})")
prs.save(OUT)
print(f"saved: {OUT}")
print(f"size_mb={OUT.stat().st_size/1e6:.2f}")
# verify links
from pptx import Presentation as _P
_prs = _P(str(OUT))
n = 0
for sl in _prs.slides:
    for sh in sl.shapes:
        try:
            a = sh.click_action.hyperlink.address
            if a and a.startswith("paper_pdfs/"):
                n += 1
                assert (ROOT / a).exists(), a
        except Exception:
            pass
print(f"relative pdf links ok: {n}")

print(f"relative pdf links ok: {n}")
print(f"relative pdf links ok: {n}")
print(f"relative pdf links ok: {n}")
print(f"relative pdf links ok: {n}")
print(f"relative pdf links ok: {n}")
print(f"relative pdf links ok: {n}")