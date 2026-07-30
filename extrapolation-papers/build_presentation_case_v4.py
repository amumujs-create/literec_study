#!/usr/bin/env python3
"""Build CA-CSS v4 case appendix deck — matches v5 dark theme."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
CASE_ASSETS = ROOT / "_assets" / "case_v4"
OUT = ROOT / "외삽_50분_사례_CA-CSS_v4_부록.pptx"

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(0x0B, 0x10, 0x16)
CARD = RGBColor(0x14, 0x1B, 0x24)
PANEL = RGBColor(0x18, 0x21, 0x2C)
LINE = RGBColor(0x2C, 0x38, 0x48)
INK = RGBColor(0xEC, 0xF1, 0xF7)
SLATE = RGBColor(0xA0, 0xAD, 0xBC)
MUTED = RGBColor(0x6E, 0x7C, 0x8C)
TEAL = RGBColor(0x3D, 0xD6, 0xC6)
CORAL = RGBColor(0xE8, 0x9A, 0x5C)
NAVY = RGBColor(0x0C, 0x12, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xB8, 0xC4, 0xD0)
ACCENT_DIM = RGBColor(0x24, 0x3A, 0x42)

TOTAL = 5
page = 0


def fig_panel(slide, x, y, w, h):
    return rect(slide, x, y, w, h, PANEL, LINE, radius=0.04)


def add_fig(slide, name, left, top, width=None, height=None, pad=0.06, caption=None):
    from PIL import Image
    path = CASE_ASSETS / name
    if not path.exists():
        return None
    iw, ih = Image.open(path).size
    if width is not None and height is None:
        fw, fh = width, width * (ih / iw)
    elif height is not None and width is None:
        fh, fw = height, height * (iw / ih)
    else:
        fw, fh = width or Inches(5), height or Inches(3)
    pad_i = Inches(pad)
    fig_panel(slide, left, top, fw + 2 * pad_i, fh + 2 * pad_i)
    pic = slide.shapes.add_picture(str(path), left + pad_i, top + pad_i, width=fw, height=fh)
    if caption:
        cy = top + fh + 2 * pad_i + Inches(0.02)
        cb = slide.shapes.add_textbox(left, cy, fw + 2 * pad_i, Inches(0.24))
        add_text(cb, caption, size=9, color=SOFT)
    return pic


def _set_run(run, size=18, bold=False, color=INK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return p


def rect(slide, x, y, w, h, fill=CARD, line=None, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def bar(slide, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def footer(slide, part=""):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.08), Inches(10), Inches(0.32))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"CA-CSS v4 Case Appendix  ·  Kookmin Univ. IE Lab  ·  2026  {('·  ' + part) if part else ''}"
    _set_run(run, size=9, color=MUTED)
    num = slide.shapes.add_textbox(Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.32))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {TOTAL}"
    _set_run(r2, size=9, color=MUTED)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s


def content_header(slide, kicker, title):
    bar(slide, TEAL)
    k = slide.shapes.add_textbox(Inches(0.42), Inches(0.26), Inches(10.5), Inches(0.30))
    add_text(k, kicker, size=11, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.42), Inches(0.52), Inches(12.0), Inches(0.50))
    add_text(t, title, size=24, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def card_text(slide, x, y, w, h, title, bullets, accent=TEAL, dense=False):
    sh = rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14), w - Inches(0.38), h - Inches(0.22))
    tf = add_text(box, title, size=12 if dense else 14, bold=True, color=accent)
    fsz = 11 if dense else 13
    for b in bullets:
        add_para(tf, "•  " + b, size=fsz, color=SLATE, space_before=3 if dense else 5)
    return sh


def takeaway(slide, main, sub=None):
    y = Inches(6.38)
    rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.56), ACCENT_DIM, LINE, radius=0.04)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.07), Inches(0.56))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.62), y + Inches(0.05), Inches(12.1), Inches(0.48))
    tf = add_text(t, main, size=12, bold=True, color=INK)
    if sub:
        add_para(tf, sub, size=10, color=SOFT, space_before=1)


def section_slide(prs, part_no, title, subtitle, points=None):
    global page
    page += 1
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11), Inches(0.5))
    add_text(t, f"APPENDIX · PART {part_no}", size=16, bold=True, color=TEAL)
    t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11), Inches(1.0))
    add_text(t2, title, size=36, bold=True, color=WHITE)
    t3 = s.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11), Inches(0.5))
    add_text(t3, subtitle, size=18, color=SOFT)
    if points:
        for i, p in enumerate(points[:3]):
            x = Inches(1.0 + i * 3.9)
            rect(s, x, Inches(4.4), Inches(3.6), Inches(1.6), PANEL, LINE)
            tb = s.shapes.add_textbox(x + Inches(0.2), Inches(4.55), Inches(3.2), Inches(1.3))
            tf = add_text(tb, f"{i+1}", size=14, bold=True, color=TEAL)
            add_para(tf, p, size=14, color=WHITE, space_before=6)
    footer(s)
    return s


def slide(prs, kicker, title, part="", note_text=""):
    global page
    page += 1
    s = blank(prs)
    content_header(s, kicker, title)
    footer(s, part)
    if note_text:
        note(s, note_text)
    return s


def add_table(slide, x, y, w, headers, rows, col_widths=None, row_h=0.38, body_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    rh = Inches(row_h)
    th = rh * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, th)
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    hdr_size = body_size if body_size <= 10 else body_size - 1
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _set_run(r, size=hdr_size, bold=True, color=TEAL)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            bold = j == 0
            color = INK if bold else SLATE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_run(r, size=body_size, bold=bold, color=color)
    return shape


def load_json_metrics(name: str) -> list[dict]:
    import json as _json
    p = CASE_ASSETS / name
    if not p.exists():
        return []
    return _json.loads(p.read_text())


COMP_ROWS = [
    ["Inputs X_seq", "엔진·운용→센서", "window 시계열 입력"],
    ["OC norm · cycle", "스케일 차", "정규화·학습 안정"],
    ["TRA Mask", "H ⊥ 부하", "health 경로 TRA→0"],
    ["Transformer×2", "열화 시간 누적", "MHA 시계열 의존"],
    ["Health Head → H", "H > 0", "softplus 잠재 건강"],
    ["MonotoneLoadHead", "TRA↑→D↑", "∂RUL/∂TRA<0 **구조 항등식**"],
    ["RUL = H − D", "건강−손상", "분해 가정 (disentangled)"],
    ["Isotonic", "cycle↑→RUL↓", "엔진별 시간 단조 후처리"],
    ["L_RUL", "—", "MSE · 라벨 fit (주 loss)"],
    ["L_physics", "TRA Jacobian soft", "λ·L_phys — **보조** (구조와 중복)"],
]


def fmt_pm(mean: float, std: float, digits: int = 2) -> str:
    if std and std > 0:
        return f"{mean:.{digits}f}±{std:.{digits}f}"
    return f"{mean:.{digits}f}"


def build():
    """Case deck — 5 slides: 문제 · 구조 · physics/구조 · 결과 · 검증."""
    global page
    page = 0
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 · 문제 — 시험(외삽) 구간만
    s = slide(prs, "사례 1/5", "시험(외삽) 구간에서 터보팬 수명을 맞추는 문제",
              "N-CMAPSS · 처음 보는 엔진 + 고부하(TRA>q90)만 평가.")
    card_text(s, Inches(0.42), Inches(1.12), Inches(4.5), Inches(2.3), "다루는 문제",
              [
                  "항공 터보팬 남은 수명(RUL) 예측",
                  "N-CMAPSS — 엔진 9대, 센서 시계열",
                  "부하·스로틀(TRA)이 바뀌면 분포도 같이 바뀜",
                  "라벨은 ‘몇 cycle 남았나’ — 순간 TRA와 무관 (∂RUL/∂TRA=0)",
              ], TEAL, dense=True)
    card_text(s, Inches(0.42), Inches(3.55), Inches(4.5), Inches(2.5), "평가 구간 (시험·외삽)",
              [
                  "엔진 11·14·15 — 훈련에 없던 3대",
                  "TRA > q90 — 본편 hard와 같은 고부하 밴드",
                  "새 엔진 + 새 부하 동시 (n≈159 windows)",
                  "fair epoch: v4 e15 · LSTM/Trans e120 → TabPFN 4.8대",
              ], CORAL, dense=True)
    add_fig(s, "fig_tra_hard_split.png", Inches(5.05), Inches(1.1), width=Inches(7.85),
            caption="시험(외삽)만 — 엔진 11·14·15 · TRA > q90")
    takeaway(s, "hard = unit×regime 결합 외삽 — 여기서만 headline 주장")

    # 2 · 구조 + 설계 의도
    s = slide(prs, "사례 2/5", "CA-CSS v4 구조 — 방향은 아키텍처에 내장",
              "본편 S20 CMNN echo: loss가 아니라 구조가 prior.")
    add_fig(s, "fig_v4_architecture.png", Inches(0.42), Inches(1.0), height=Inches(2.45),
            caption="RUL=H−D · health TRA-blind · MonotoneLoadHead")
    half = len(COMP_ROWS) // 2 + 1
    add_table(s, Inches(0.42), Inches(3.62), Inches(6.15),
              ["구성요소", "가정", "설계 의도"],
              COMP_ROWS[:half],
              col_widths=[Inches(1.45), Inches(1.75), Inches(2.95)],
              row_h=0.22, body_size=7.5)
    add_table(s, Inches(6.75), Inches(3.62), Inches(6.15),
              ["구성요소", "가정", "설계 의도"],
              COMP_ROWS[half:],
              col_widths=[Inches(1.45), Inches(1.75), Inches(2.95)],
              row_h=0.22, body_size=7.5)
    takeaway(s, "부하→수명 방향은 MonotoneLoadHead 항등식 — TRA mask + RUL=H−D",
             sub="∂RUL/∂TRA ≤ −softplus(w) < 0 (구조상 위반 불가)")

    # 3 · physics loss vs 구조 prior (핵심 수정)
    s = slide(prs, "사례 3/5", "Physics loss vs 구조 prior — 무엇이 방향을 만드는가",
              "PINN式 soft loss ≠ 여기서의 방향 메커니즘.")
    add_fig(s, "fig_physics_vs_structure.png", Inches(0.42), Inches(1.05), width=Inches(12.5),
            caption="λ_tra 스윕: RMSE·adherence 모두 무반응 · no_physics +0.33")
    card_text(s, Inches(0.42), Inches(4.05), Inches(6.0), Inches(2.15), "L = L_RUL + λ·L_physics",
              [
                  "L_physics: TRA Jacobian·단조 hinge (soft constraint)",
                  "no_physics: RMSE +0.33 — fit 보조 역할",
                  "λ_tra 0 / 0.05 / 0.5: RMSE ±0.09 · adherence 100%",
              ], TEAL, dense=True)
    card_text(s, Inches(6.55), Inches(4.05), Inches(6.35), Inches(2.15), "구조 ablation → 방향 인과",
              [
                  "full: RMSE 3.88 · ΔRUL<0 100%",
                  "free_load: RMSE 4.24 · adherence 0%",
                  "mono head 제거 = worst seed 6.1",
                  "레버 = MonotoneLoadHead, not λ_tra",
              ], CORAL, dense=True)
    takeaway(s, "physics loss는 fit 보조 — counterfactual 방향은 구조 prior",
             sub="‘물리 법칙 준수’ ✗ · ‘도메인 prior + OOD RMSE’ ✓")

    # 4 · 결과
    s = slide(prs, "사례 4/5", "hard 구간에서 다른 모델과 비교하면",
              "모델마다 학습 epoch 맞춘 공정 비교.")
    add_fig(s, "fig_baseline_fair_hard.png", Inches(0.42), Inches(1.0), width=Inches(8.3),
            caption="각 모델에 맞게 epoch 조정 후 hard 구간 RMSE·R²")
    fair_panels = load_json_metrics("fair_baseline_metrics.json")
    tbl_fair = []
    for panel in fair_panels:
        if panel.get("dataset") != "ncmapss_hard":
            continue
        for m in panel.get("models", []):
            ep = m.get("epoch", "")
            ep_s = f"e{ep}" if isinstance(ep, int) else str(ep)
            tbl_fair.append([
                m["label"],
                fmt_pm(m["rmse"], m.get("rmse_std", 0)),
                fmt_pm(m["r2"], m.get("r2_std", 0), 3),
            ])
    if not tbl_fair:
        tbl_fair = [
            ["v4+iso", "3.43±0.73", "0.966"],
            ["TabPFN", "4.79", "0.937"],
            ["Trans", "4.99±0.46", "0.931"],
            ["LSTM", "4.79±0.55", "0.936"],
        ]
    add_table(s, Inches(8.85), Inches(1.08), Inches(4.05),
              ["모델", "오차(RMSE)", "R²"], tbl_fair,
              col_widths=[Inches(1.15), Inches(1.05), Inches(0.85)])
    card_text(s, Inches(8.85), Inches(3.55), Inches(4.05), Inches(2.5), "읽는 포인트",
              [
                  "LSTM/Trans도 fair epoch면 4.8~5.0",
                  "v4+iso 3.43±0.73 · R² 0.97 (10-seed)",
                  "p=0.00017 vs TabPFN — hard만 차별화",
                  "C-MAPSS e120 ≈ Transformer (headline ✗)",
              ], CORAL, dense=True)
    takeaway(s, "결합 hard만 v4 유일 우위 — C-MAPSS 단일축은 수렴하면 비슷")

    # 5 · 검증 · 요약
    s = slide(prs, "사례 5/5", "무엇까지 말할 수 있고, 한 줄로 정리하면",
              "주장 범위를 좁히고, 본편 메시지로 회수.")
    card_text(s, Inches(0.42), Inches(1.1), Inches(5.9), Inches(3.2), "말해도 되는 것",
              [
                  "hard extrap RMSE ↓ (10-seed, p=0.00017)",
                  "구조 prior → model-input ΔRUL<0 100%",
                  "MonotoneLoadHead ablation으로 인과 확인",
                  "본편 S30 체크 4항 충족",
              ], TEAL, dense=True)
    card_text(s, Inches(6.5), Inches(1.1), Inches(6.4), Inches(3.2), "말하지 않는 것",
              [
                  "‘데이터가 TRA↑⇒RUL↓ 증명’ (r≈0, ∂RUL/∂TRA=0)",
                  "‘λ_tra·physics loss가 방향 생성’ (λ=0도 100%)",
                  "‘PINN/PDE 물리 법칙 준수’",
                  "‘C-MAPSS 전 split SOTA’",
              ], CORAL, dense=True)
    rows = [
        ("문제", "N-CMAPSS hard — unit×TRA 밖 (결합 외삽)"),
        ("가정", "counterfactual 부하 prior → MonotoneLoadHead 구조"),
        ("loss", "L_RUL fit + L_physics 보조 — 방향은 구조"),
        ("결과", "RMSE 3.43 vs TabPFN 4.79 · worst-seed 방어"),
    ]
    y0 = Inches(4.45)
    for i, (k, v) in enumerate(rows):
        rect(s, Inches(0.42), y0 + Inches(i * 0.48), Inches(12.5), Inches(0.42), PANEL, LINE, radius=0.03)
        add_text(s.shapes.add_textbox(Inches(0.58), y0 + Inches(i * 0.48) + Inches(0.06), Inches(1.6), Inches(0.32)),
                 k, size=11, bold=True, color=TEAL)
        add_text(s.shapes.add_textbox(Inches(2.1), y0 + Inches(i * 0.48) + Inches(0.06), Inches(10.5), Inches(0.32)),
                 v, size=11, color=SOFT)
    takeaway(s, "밖을 버티는 건 데이터가 아니라 구조 prior — physics loss는 fit 보조",
             sub="본편 echo: 가정을 구조에 넣고, hard에서 검증")

    prs.save(str(OUT))
    print(f"Saved {OUT} ({page} slides)")


if __name__ == "__main__":
    build()
