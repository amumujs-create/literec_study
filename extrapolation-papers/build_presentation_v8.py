#!/usr/bin/env python3
"""Build v8 from v7: replace assumption / identifiability figures (train-fit fix)."""
from __future__ import annotations

import hashlib
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "_assets"
SRC = ROOT / "외삽_50분_발표자료_v7 [Repaired].pptx"
DST = ROOT / "외삽_50분_발표자료_v8.pptx"

# slide index (0-based) -> asset filename
REPLACEMENTS = {
    2: "fig_assumption_defined.png",  # S03 가정
    9: "fig_identifiability.png",  # S10 식별불가
}


def _md5(data: bytes) -> str:
    return hashlib.md5(data).hexdigest()


def _largest_picture(slide):
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        raise RuntimeError("no picture on slide")
    return max(pics, key=lambda sh: int(sh.width) * int(sh.height))


def _replace_picture_blob(picture, new_blob: bytes) -> None:
    blip = picture._element.find(".//" + qn("a:blip"))
    if blip is None:
        raise RuntimeError("picture has no a:blip")
    r_id = blip.get(qn("r:embed"))
    if not r_id:
        raise RuntimeError("picture blip missing r:embed")
    image_part = picture.part.related_part(r_id)
    image_part._blob = new_blob  # noqa: SLF001


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"missing source: {SRC}")
    for name in REPLACEMENTS.values():
        if not (ASSETS / name).exists():
            raise SystemExit(f"missing asset: {ASSETS / name}")

    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))

    for slide_idx, asset_name in REPLACEMENTS.items():
        new_blob = (ASSETS / asset_name).read_bytes()
        slide = prs.slides[slide_idx]
        target = _largest_picture(slide)
        old = _md5(target.image.blob)
        _replace_picture_blob(target, new_blob)
        print(f"S{slide_idx + 1:02d} {asset_name}: {old[:8]} → {_md5(new_blob)[:8]}")

    prs.save(str(DST))
    print(f"wrote {DST} ({DST.stat().st_size} bytes)")

    prs2 = Presentation(str(DST))
    for slide_idx, asset_name in REPLACEMENTS.items():
        emb = _md5(_largest_picture(prs2.slides[slide_idx]).image.blob)
        asset = _md5((ASSETS / asset_name).read_bytes())
        if emb != asset:
            raise SystemExit(f"verify failed S{slide_idx + 1:02d}")
        print(f"verify S{slide_idx + 1:02d} OK")


if __name__ == "__main__":
    main()
