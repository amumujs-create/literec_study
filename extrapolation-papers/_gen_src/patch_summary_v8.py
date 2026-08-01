#!/usr/bin/env python3
"""Patch S25 (마무리) in v8[Repaired].pptx with cleaned 2×2 layout."""
from __future__ import annotations

import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
V8 = ROOT / "외삽_50분_발표자료_v8[Repaired].pptx"
SLIDE_IDX = 24  # 0-based → slide 25
TOTAL = 31

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(0x0B, 0x10, 0x16)
CARD = RGBColor(0x14, 0x1B, 0x24)
LINE = RGBColor(0x2C, 0x38, 0x48)
INK = RGBColor(0xEC, 0xF1, 0xF7)
SLATE = RGBColor(0xA0, 0xAD, 0xBC)
MUTED = RGBColor(0x6E, 0x7C, 0x8C)
TEAL = RGBColor(0x3D, 0xD6, 0xC6)
CORAL = RGBColor(0xE8, 0x9A, 0x5C)
NAVY = RGBColor(0x0C, 0x12, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xB8, 0xC4, 0xD0)
PANEL = RGBColor(0x18, 0x21, 0x2C)
ACCENT_DIM = RGBColor(0x24, 0x3A, 0x42)

PDF = {
    "xu": "paper_pdfs/Xu2021_How_Neural_Networks_Extrapolate.pdf",
    "runje": "paper_pdfs/Runje2023_Constrained_Monotonic_NN.pdf",
    "bartley": "paper_pdfs/Bartley2019_Characterizing_Extrapolation_Multivariate.pdf",
}

_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")


def _set_run(run, size=18, bold=False, color=INK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def _add_rich_runs(p, text, size, bold, color, font="Apple SD Gothic Neo"):
    pos = 0
    found = False
    for m in _BOLD_RE.finditer(text):
        found = True
        if m.start() > pos:
            run = p.add_run()
            run.text = text[pos:m.start()]
            _set_run(run, size=size, bold=bold, color=color, font=font)
        run = p.add_run()
        run.text = m.group(1)
        _set_run(run, size=size, bold=True, color=color, font=font)
        pos = m.end()
    if not found:
        run = p.add_run()
        run.text = text
        _set_run(run, size=size, bold=bold, color=color, font=font)
    elif pos < len(text):
        run = p.add_run()
        run.text = text[pos:]
        _set_run(run, size=size, bold=bold, color=color, font=font)


def add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _add_rich_runs(p, text, size, bold, color, font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    _add_rich_runs(p, text, size, bold, color)
    return p


def rect(slide, x, y, w, h, fill=CARD, line=None, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def bar(slide, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def content_header(slide, kicker, title):
    bar(slide, TEAL)
    k = slide.shapes.add_textbox(Inches(0.42), Inches(0.26), Inches(10.5), Inches(0.30))
    add_text(k, kicker, size=11, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.42), Inches(0.52), Inches(12.0), Inches(0.50))
    add_text(t, title, size=24, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def card_text(slide, x, y, w, h, title, bullets, accent=TEAL, dense=False):
    title_color = TEAL if accent == NAVY else accent
    strip_color = TEAL if accent == NAVY else accent
    rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = strip_color
    strip.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14), w - Inches(0.38), h - Inches(0.22))
    tf = add_text(box, title, size=12 if dense else 14, bold=True, color=title_color)
    gap = 3 if dense else 5
    fsz = 11 if dense else 13
    for b in bullets:
        add_para(tf, "–  " + b, size=fsz, color=SLATE, space_before=gap, space_after=0)


def pdf_btn(slide, key, left, top, w=Inches(1.35), h=Inches(0.34), label=None):
    rel = PDF.get(key)
    if rel is None or not (ROOT / rel).exists():
        return None
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = CORAL
    sh.line.width = Pt(1.0)
    try:
        sh.adjustments[0] = 0.18
    except Exception:
        pass
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label or f"PDF · {key.upper()}"
    _set_run(run, size=10, bold=True, color=CORAL)
    sh.click_action.hyperlink.address = rel
    return sh


def footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.08), Inches(10), Inches(0.32))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Extrapolation Seminar  ·  Kookmin Univ. IE Lab  ·  2026"
    _set_run(run, size=9, color=MUTED)
    num = slide.shapes.add_textbox(Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.32))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    _set_run(r2, size=9, color=MUTED)


def takeaway(slide, main, sub=None, y=Inches(6.28)):
    rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.52), ACCENT_DIM, LINE, radius=0.04)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.07), Inches(0.52))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.62), y + Inches(0.05), Inches(12.1), Inches(0.44))
    tf = add_text(t, main, size=12, bold=True, color=INK)
    if sub:
        add_para(tf, sub, size=10, color=SOFT, space_before=1, space_after=0)


def clear_slide(slide):
    bg_xml = None
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE.RECTANGLE and sh.width == W and sh.height == H:
            bg_xml = deepcopy(sh.element)
        sh.element.getparent().remove(sh.element)
    if bg_xml is not None:
        slide.shapes._spTree.insert_element_before(bg_xml, "p:extLst")
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG
        bg.line.fill.background()


def build_summary(slide):
    content_header(slide, "마무리", "**회귀 extrap** — 오늘 정리")

    rect(slide, Inches(0.45), Inches(1.15), Inches(12.43), Inches(0.82), NAVY, TEAL, radius=0.05)
    hero = slide.shapes.add_textbox(Inches(0.65), Inches(1.22), Inches(12.0), Inches(0.72))
    hero_tf = add_text(
        hero,
        "train **밖**은 **데이터**가 아니라 **가정**이 지탱한다",
        size=21,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_para(
        hero_tf,
        "“외삽 되나?” → **가정**이 맞고 · **밖에서 시험**했는가?",
        size=13,
        color=SOFT,
        align=PP_ALIGN.CENTER,
        space_before=5,
    )

    story = [
        ("① 1부 · 정의", ["외삽 = hull **밖** 예측", "연속 회귀 · 분류 DG와 다름"], TEAL),
        ("② 2부 · 실패", ["데이터만으론 밖 **모양** 못 정함", "NN도 밖에선 직선 · 불확실성↑"], CORAL),
        ("③ 3부 · 대응", ["아는 만큼 **가정**을 넣는다", "식 · 방향 · 물리 · 모르면 UQ"], TEAL),
        ("④ 실무 · 검증", ["**밖 holdout**으로 확인 후 배포", "못 믿으면 **예측 안 함**(기권)"], CORAL),
    ]
    cw = Inches(6.08)
    ch = Inches(1.18)
    cgx = Inches(0.14)
    cgy = Inches(0.12)
    sy = Inches(2.12)
    for i, (title, bullets, accent) in enumerate(story):
        col, row = i % 2, i // 2
        x = Inches(0.45 + col * (cw + cgx))
        y = sy + row * (ch + cgy)
        card_text(slide, x, y, cw, ch, title, bullets, accent, dense=True)

    lbl = slide.shapes.add_textbox(Inches(0.45), Inches(4.72), Inches(2.5), Inches(0.26))
    add_text(lbl, "필독 3편", size=11, bold=True, color=TEAL)

    trio = [
        ("진단", "Xu 2021", "ReLU·직선화 — **왜** 밖에서 깨지나", "xu", CORAL),
        ("처방", "Runje 2023", "CMNN — **방향** 가정을 구조에", "runje", TEAL),
        ("검증", "Bartley 2019", "hull·고차원 — **밖**인지 판정", "bartley", TEAL),
    ]
    tw = Inches(3.98)
    tg = Inches(0.14)
    ty = Inches(5.02)
    th = Inches(1.08)
    for i, (role, cite, desc, key, accent) in enumerate(trio):
        x = Inches(0.45 + i * (tw + tg))
        card_text(slide, x, ty, tw, th, f"{role} · {cite}", [desc], accent, dense=True)
        pdf_btn(slide, key, x + tw - Inches(1.08), ty + Inches(0.72), w=Inches(0.92), h=Inches(0.30), label="PDF")

    footer(slide, SLIDE_IDX + 1, TOTAL)
    takeaway(slide, "가정 맞추기 · **밖에서 시험** · UQ/기권", "Xu · Runje · Bartley", y=Inches(6.28))

    note = (
        "[1.5분]\n"
        "관통 문장(hero) → 2×2 스토리(①~④) → 필독 3편.\n"
        "더 읽기: Pfister·Fesser·Zhu·Li. Q&A로."
    )
    slide.notes_slide.notes_text_frame.text = note


def main():
    path = Path(sys.argv[1]) if len(sys.argv) > 1 else V8
    if not path.exists():
        raise SystemExit(f"missing: {path}")
    prs = Presentation(str(path))
    if len(prs.slides) <= SLIDE_IDX:
        raise SystemExit(f"expected slide {SLIDE_IDX + 1}, got {len(prs.slides)}")
    slide = prs.slides[SLIDE_IDX]
    clear_slide(slide)
    build_summary(slide)
    prs.save(str(path))
    print(f"patched slide {SLIDE_IDX + 1} in {path}")


if __name__ == "__main__":
    main()
