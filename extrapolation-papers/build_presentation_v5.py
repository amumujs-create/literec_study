#!/usr/bin/env python3
"""Build 외삽 50분 발표자료 v5 — 단일 축(가정) 심화판. 발표자 노트에 대본 전문 포함."""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "_assets"
PAPER = ASSETS / "paper_figs"
OUT = ROOT / "외삽_50분_발표자료_v5_심화.pptx"

# 16:9
W, H = Inches(13.333), Inches(7.5)

# Dark instrument theme (match figure console)
BG = RGBColor(0x0B, 0x10, 0x16)
CARD = RGBColor(0x14, 0x1B, 0x24)
PANEL = RGBColor(0x18, 0x21, 0x2C)
LINE = RGBColor(0x2C, 0x38, 0x48)
INK = RGBColor(0xEC, 0xF1, 0xF7)
SLATE = RGBColor(0xA0, 0xAD, 0xBC)
MUTED = RGBColor(0x6E, 0x7C, 0x8C)
TEAL = RGBColor(0x3D, 0xD6, 0xC6)
CORAL = RGBColor(0xE8, 0x9A, 0x5C)
NAVY = RGBColor(0x0C, 0x12, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xB8, 0xC4, 0xD0)
ACCENT_DIM = RGBColor(0x24, 0x3A, 0x42)  # teal-tinted panel


def _set_run(run, size=18, bold=False, color=INK, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")


def _add_rich_runs(p, text, size, bold, color, font="Apple SD Gothic Neo"):
    """Split text on **bold** markers into runs — no literal asterisks on slide."""
    pos = 0
    found = False
    for m in _BOLD_RE.finditer(text):
        found = True
        if m.start() > pos:
            run = p.add_run()
            run.text = text[pos:m.start()]
            _set_run(run, size=size, bold=bold, color=color, font=font)
        run = p.add_run()
        run.text = m.group(1)
        _set_run(run, size=size, bold=True, color=color, font=font)
        pos = m.end()
    if not found:
        run = p.add_run()
        run.text = text
        _set_run(run, size=size, bold=bold, color=color, font=font)
    elif pos < len(text):
        run = p.add_run()
        run.text = text[pos:]
        _set_run(run, size=size, bold=bold, color=color, font=font)


def add_text(shape, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _add_rich_runs(p, text, size, bold, color, font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    _add_rich_runs(p, text, size, bold, color)
    return p


def rect(slide, x, y, w, h, fill=CARD, line=None, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def bar(slide, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def footer(slide, page, total, part=""):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.08), Inches(10), Inches(0.32))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Extrapolation Seminar  ·  Kookmin Univ. IE Lab  ·  2026  {('·  ' + part) if part else ''}"
    _set_run(run, size=9, color=MUTED)
    num = slide.shapes.add_textbox(Inches(11.6), Inches(7.08), Inches(1.4), Inches(0.32))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    _set_run(r2, size=9, color=MUTED)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s


def content_header(slide, kicker, title):
    bar(slide, TEAL)
    k = slide.shapes.add_textbox(Inches(0.42), Inches(0.26), Inches(10.5), Inches(0.30))
    add_text(k, kicker, size=11, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.42), Inches(0.52), Inches(12.0), Inches(0.50))
    add_text(t, title, size=24, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def card_text(slide, x, y, w, h, title, bullets, accent=TEAL, dense=False):
    # NAVY is a fill color (near-black) — never use it as title ink on dark cards
    title_color = TEAL if accent == NAVY else accent
    strip_color = TEAL if accent == NAVY else accent
    sh = rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = strip_color
    strip.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14), w - Inches(0.38), h - Inches(0.22))
    tf = add_text(box, title, size=12 if dense else 14, bold=True, color=title_color)
    gap = 3 if dense else 5
    fsz = 11 if dense else 13
    for b in bullets:
        # subtle separator lines as plain bullets
        prefix = "–  " if dense else "•  "
        add_para(tf, prefix + b, size=fsz, color=SLATE, space_before=gap, space_after=0)
    return sh


def fig_panel(slide, x, y, w, h):
    """Framed panel behind paper figures (soft chrome)."""
    return rect(slide, x, y, w, h, PANEL, LINE, radius=0.04)


# ── build ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Estimate total slides for footers (update if you add/remove)
TOTAL = 30  # v5: UQ+기권 합침, PINN 1장, OoD-Bench·DomainBed 제거 (−2)
page = 0

# Paper PDF paths (absolute URI for PowerPoint click)
# Relative paths from PPTX location (ASCII folder — Korean path file:// breaks on macOS)
PDF = {
    "xu": "paper_pdfs/Xu2021_How_Neural_Networks_Extrapolate.pdf",
    "eql": "paper_pdfs/Martius2016_Extrapolation_Learning_Equations_EQL.pdf",
    "nalu": "paper_pdfs/Trask2018_Neural_Arithmetic_Logic_Units_NALU.pdf",
    "runje": "paper_pdfs/Runje2023_Constrained_Monotonic_NN.pdf",
    "fesser": "paper_pdfs/Fesser2023_Extrapolation_Failures_PINNs.pdf",
    "zhu": "paper_pdfs/Zhu2022_Reliable_Extrapolation_DeepONet.pdf",
    "pfister": "paper_pdfs/Pfister2024_Extrapolation-Aware_Nonparametric_Inference.pdf",
    "domainbed": "paper_pdfs/Gulrajani2020_In_Search_of_Lost_Domain_Generalization.pdf",
    "liu": "paper_pdfs/Liu2023_OOD_Generalization_Survey.pdf",
    "arjovsky2019": "paper_pdfs/Arjovsky2019_Invariant_Risk_Minimization.pdf",
    "arjovsky2021": "paper_pdfs/Arjovsky2021_OOD_Generalization_in_ML.pdf",
    "raissi": "paper_pdfs/Raissi2019_Physics_Informed_Neural_Networks.pdf",
    "teckentrup": "paper_pdfs/Teckentrup2024_Probabilistic_Richardson_Extrapolation.pdf",
    "wu": "paper_pdfs/Wu2025_OOD_Time_Series_Survey.pdf",
    "ye2021": "paper_pdfs/Ye2021_Theoretical_Framework_OOD.pdf",
    "ye2022": "paper_pdfs/Ye2022_OoD-Bench.pdf",
    "bartley": "paper_pdfs/Bartley2019_Characterizing_Extrapolation_Multivariate.pdf",
    "aykol": "paper_pdfs/Aykol2021_Physics_ML_Battery_Lifetime.pdf",
    "li2023": "paper_pdfs/Li2023_Predicting_Battery_Lifetime_Varying_Conditions.pdf",
    "muckley": "paper_pdfs/Muckley2023_Interpretable_Models_Extrapolation_SciML.pdf",
    "decugis": "paper_pdfs/Decugis2024_Extrapolation_Power_Implicit_Models.pdf",
    "ghahramani": "paper_pdfs/Ghahramani2013_Bayesian_Nonparametrics.pdf",
    "note": "paper_pdfs/외삽_완전정복_학습노트.pdf",
}


def add_paper_fig(slide, name, left, top, width=None, height=None, pad=0.07, caption=None,
                   caption_w=None, caption_h=Inches(0.24), align=PP_ALIGN.LEFT):
    """Paper figure inside soft framed panel. Returns picture shape or None.
    caption: paper title strip drawn directly under the panel.
    caption_w: override caption box width (e.g. a shared column width) — panel stays image-sized."""
    from PIL import Image
    path = PAPER / name
    if not path.exists():
        return None
    iw, ih = Image.open(path).size
    if width is not None and height is None:
        fw, fh = width, width * (ih / iw)
    elif height is not None and width is None:
        fh, fw = height, height * (iw / ih)
    else:
        fw, fh = width, height
    pad_i = Inches(pad)
    fig_panel(slide, left, top, fw + 2 * pad_i, fh + 2 * pad_i)
    pic = slide.shapes.add_picture(str(path), left + pad_i, top + pad_i, width=fw, height=fh)
    if caption:
        cy = top + fh + 2 * pad_i + Inches(0.02)
        cw = caption_w if caption_w is not None else (fw + 2 * pad_i)
        cx = left - (cw - (fw + 2 * pad_i)) / 2 if caption_w is not None else left
        cb = slide.shapes.add_textbox(cx, cy, cw, caption_h)
        cb.text_frame.word_wrap = True
        add_text(cb, caption, size=9, color=SOFT, align=align)
    return pic


def credit(slide, text, y=Inches(6.74)):
    box = slide.shapes.add_textbox(Inches(0.5), y, Inches(10.5), Inches(0.26))
    add_text(box, text, size=8, color=MUTED)


def takeaway(slide, main, sub=None, y=Inches(6.38)):
    """하단 발표 핵심."""
    sh = rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.56), ACCENT_DIM, LINE, radius=0.04)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y, Inches(0.07), Inches(0.56))
    strip.fill.solid()
    strip.fill.fore_color.rgb = TEAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.62), y + Inches(0.05), Inches(12.1), Inches(0.48))
    tf = add_text(t, main, size=12, bold=True, color=INK)
    if sub:
        add_para(tf, sub, size=10, color=SOFT, space_before=1, space_after=0)
    return sh


def speak(slide, bullets, x=Inches(0.4), y=Inches(5.25), w=Inches(12.5), h=Inches(1.0)):
    sh = rect(slide, x, y, w, h, CARD, LINE, radius=0.05)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = CORAL
    strip.line.fill.background()
    t = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.08), w - Inches(0.35), h - Inches(0.12))
    tf = add_text(t, "Remarks", size=11, bold=True, color=CORAL)
    for b in bullets:
        add_para(tf, "–  " + b, size=12, color=SLATE, space_before=3, space_after=0)
    return sh


def section_slide(title, subtitle, part_no, mins, page, total, points=None):
    s = blank()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11), Inches(0.5))
    add_text(t, f"PART {part_no}  ·  {mins}", size=16, bold=True, color=TEAL)
    t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11), Inches(1.0))
    add_text(t2, title, size=36, bold=True, color=WHITE)
    t3 = s.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11), Inches(0.5))
    add_text(t3, subtitle, size=18, color=SOFT)
    if points:
        for i, p in enumerate(points[:3]):
            x = Inches(1.0 + i * 3.9)
            box = rect(s, x, Inches(4.4), Inches(3.6), Inches(1.6), PANEL, LINE)
            tb = s.shapes.add_textbox(x + Inches(0.2), Inches(4.55), Inches(3.2), Inches(1.3))
            tf = add_text(tb, f"{i+1}", size=14, bold=True, color=TEAL)
            add_para(tf, p, size=14, color=WHITE, space_before=6)
    footer_box = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.35))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page} / {total}"
    _set_run(r, size=10, color=MUTED)
    return s


def pdf_btn(slide, key, left, top, w=Inches(1.35), h=Inches(0.34), label=None):
    """Clickable outline button → open paper PDF."""
    rel = PDF.get(key)
    if rel is None or not (ROOT / rel).exists():
        return None
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = CORAL
    sh.line.width = Pt(1.0)
    try:
        sh.adjustments[0] = 0.18
    except Exception:
        pass
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label or f"PDF · {key.upper()}"
    _set_run(run, size=10, bold=True, color=CORAL)
    sh.click_action.hyperlink.address = rel
    return sh


def P():
    global page
    page += 1
    return page


# ═══════════════════════════════════════════════════════
# v5 — 세 가지 질문으로 진행
#   Q1. 외삽이란 무엇인가 (5')
#   Q1. 외삽 정의 — 보간/외삽 · convex hull (5')
#   Q2. 실패 원인 (4) — 밖을 데이터만으로 못 정하는 이유 (15')
#   Q3. 대응법 + 성능 검증 — 가정 넣기 + 외삽 주장 검증 (27') + 엔딩 (3')
# 모든 내용 슬라이드에 그림 1개 이상. note = 발표 대본 전문.
# ═══════════════════════════════════════════════════════

# ── S1 타이틀 ──
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.82), W, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = TEAL; stripe.line.fill.background()
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), W, Inches(1.6))
acc.fill.solid(); acc.fill.fore_color.rgb = PANEL; acc.line.fill.background()
t = s.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11), Inches(0.4))
add_text(t, "논문 기반 세미나  ·  50분", size=16, bold=True, color=TEAL)
t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.7), Inches(1.6))
add_text(t2, "회귀 외삽 (Regression Extrapolation)\n훈련 범위 밖 **연속값** 예측 — 정의 · 실패 원인 · 대응법 · 성능 검증", size=32, bold=True, color=WHITE)
t3 = s.shapes.add_textbox(Inches(1.0), Inches(4.35), Inches(11.5), Inches(0.6))
add_text(t3, "범위: tabular·시계열 **회귀**  ·  분류·DG(OOD)는 본 세미나 **범위 밖**", size=16, color=SOFT)
t3b = s.shapes.add_textbox(Inches(1.0), Inches(4.95), Inches(11.5), Inches(0.45))
add_text(t3b, "1부 외삽 정의  ·  2부 실패 원인  ·  3부 대응법 + 성능 검증", size=15, color=MUTED)
t4 = s.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11), Inches(0.8))
tf = add_text(t4, "Kookmin Univ. IE Lab", size=16, bold=True, color=WHITE)
add_para(tf, "extrapolation-papers · 2026.07", size=13, color=MUTED)
note(s, "[30초]\n"
     "대본: 오늘은 **회귀 외삽** — 연속값 y, 센서→품질 같은 tabular·시계열. "
     "분류·domain generalization(OOD)은 범위 밖. "
     "네 파트: 정의, 실패 원인, 대응법, 성능 검증. "
     "결론 미리 — 밖을 지탱하는 건 데이터가 아니라 **가정**.")
P()

# ── S2 로드맵 ──
s = blank()
content_header(s, "로드맵", "발표 구조 — **회귀** 외삽 · 정의 · 실패 · 대응 · 검증")
if (ASSETS / "fig_three_questions.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_three_questions.png"), Inches(0.55), Inches(1.18), width=Inches(12.2))
card_text(
    s, Inches(0.40), Inches(4.28), Inches(4.05), Inches(2.0),
    "1부  외삽 정의 (5분)",
    [
        "OOD vs 외삽 — 오늘은 **외삽**에 초점",
        "보간 vs 외삽 — 새 입력이 범위 안/밖인지",
        "훈련 범위 = convex hull",
        "내부 적합도 ≠ 외삽 성능",
        "논문: Bartley 2019",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(4.60), Inches(4.28), Inches(4.25), Inches(2.0),
    "2부  실패 원인 (15분)",
    [
        "① 데이터만으론 못 고름 — Pfister 2024",
        "② 훈련 범위 밖이 기본 — Bartley 2019",
        "③ 밖에선 불확실성↑ — Ghahramani",
        "④ ReLU→직선 — Xu 2021",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.85), Inches(4.28), Inches(4.05), Inches(2.0),
    "3부  대응법·성능 검증 (23분)",
    [
        "**회귀** 가정 넣기: EQL / CMNN / Physics-ML",
        "모를 때 — UQ · 기권 (연속 y)",
        "검증 — hull 밖 holdout · baseline · UQ",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL)
note(s, "[1.5분]\n"
     "대본: 위 세 상자가 목차입니다. 1부 — **외삽 정의**(보간 vs 외삽, 훈련 범위). "
     "2부 — **외삽 실패 원인** 네 가지(데이터·차원·불확실성·ReLU). "
     "3부 — **외삽 대응법**(가정 넣기)과 **외삽 성능 검증**(주장이 맞는지 확인). "
     "다음 장에서 핵심 용어 '가정'을 정의합니다.")
takeaway(s, "외삽 정의 → 실패 원인 4 → 대응법 + 성능 검증",
         "핵심: 밖 예측을 정당화하는 건 데이터가 아니라 가정")

# ── S3 NEW: 가정이란 무엇인가 (용어 고정) ──
s = blank()
content_header(s, "용어 고정", "가정(Assumption) — 오늘 발표의 중심 개념")
if (ASSETS / "fig_assumption_defined.png").exists():
    s.shapes.add_picture(
        str(ASSETS / "fig_assumption_defined.png"),
        Inches(0.35), Inches(1.12), height=Inches(3.15),
    )
card_text(
    s, Inches(8.05), Inches(1.12), Inches(4.95), Inches(1.55),
    "① 가정이란?",
    [
        "훈련 범위 밖에서 함수가",
        "  어떻게 거동할지에 대한",
        "  사전 제약·구조적 지식",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.05), Inches(2.82), Inches(4.95), Inches(1.55),
    "② 왜 필요한가",
    [
        "데이터만으로는 밖을 유일하게",
        "  정할 수 없음 (Pfister 2024)",
        "→ 가정 없이는 외삽 불가",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.05), Inches(4.52), Inches(4.95), Inches(1.55),
    "③ 어디에 넣나",
    [
        "활성화 φ → 함수족 가정",
        "단조·PDE·UQ/기권 → 방향·물리·유보",
        "3부 방법과 1:1 대응",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(4.55), Inches(7.45), Inches(1.55),
    "④ 그림 해석",
    [
        "훈련점(●)은 세 곡선 공통 · x>2 밖: A선형·B지수·C비선형 → 예측 전부 다름",
        "데이터는 같고 가정만 다름 = 오늘 발표 출발점",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL)
note(s, "[2분] ★ 용어 장 — 반드시 천천히.\n"
     "대본: '가정'이 오늘 키워드입니다. 왼쪽 그림 — 훈련점(●)은 전부 같습니다. "
     "가정 A(선형)·B(지수)·C(비선형)만 바꾸면 x>2 밖에서 곡선이 완전히 달라집니다. "
     "오른쪽 ① 정의, ② 왜 필요(Pfister), ③ 주입 위치. 아래 ④ 그림 해석. "
     "이후 '가정을 넣는다'는 이 뜻으로만 쓰겠습니다.")
takeaway(s, "가정 = 밖에서의 함수 거동을 미리 제약하는 사전 지식",
         "같은 훈련 데이터, 다른 가정 → 다른 외삽")

# ════════ 1부 · 외삽이란 무엇인가 ════════
section_slide(
    "1부 · 외삽이란 무엇인가",
    "훈련 범위의 안과 밖 — 오늘 내내 쓸 좌표계",
    "1", "5분", P(), TOTAL,
    points=[
        "OOD vs 외삽 — 용어 구분, 오늘은 **외삽**",
        "보간 vs 외삽: 입력이 훈련 범위 안/밖인지",
        "훈련 범위의 정확한 정의 = convex hull",
        "대표 사례: 내부 적합과 외부 발산",
    ],
)
note(prs.slides[-1], "[20초]\n대본: 1부입니다. 먼저 OOD와 외삽 — 헷갈리기 쉬운 두 단어. "
     "이 발표는 **외삽** 좌표계에 집중합니다. 그다음 보간/외삽, hull, 사례. 5분입니다.")

# ── S5 OOD vs 외삽 ──
s = blank()
content_header(s, "1부 · 시작", "OOD vs 외삽 — 오늘은 **외삽**을 다룹니다")
fig_path = ASSETS / "fig_ood_vs_extrap.png"
if fig_path.exists():
    fig_panel(s, Inches(0.30), Inches(1.15), Inches(8.05), Inches(4.85))
    s.shapes.add_picture(str(fig_path), Inches(0.38), Inches(1.22), width=Inches(7.90))
card_text(
    s, Inches(8.55), Inches(1.15), Inches(4.45), Inches(1.55),
    "그림 읽기",
    [
        "(a) OOD — **분포**가 train≠test",
        "(b) 외삽 — **입력**이 hull 밖",
        "★ extrap · ■ concept shift (OOD만)",
        "extrap ⊂ OOD **종종**, 같지 않음",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.55), Inches(2.82), Inches(4.45), Inches(1.55),
    "질문이 다름",
    [
        "OOD → ‘분포가 다른가?’",
        "외삽 → ‘**입력**이 hull 밖인가?’",
        "concept shift → hull로 **판별 불가**",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.55), Inches(4.48), Inches(4.45), Inches(1.52),
    "오늘 발표 범위",
    [
        "**회귀 외삽** — hull 밖 + **연속 y**",
        "분류·DG · IRM/DomainBed → **다른 문제**",
        "다음 — 보간/외삽 **구체화**",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "liu", Inches(9.75), Inches(0.28), w=Inches(1.55), label="PDF · OOD survey")
pdf_btn(s, "bartley", Inches(11.40), Inches(0.28), w=Inches(1.55), label="PDF · Bartley")
footer(s, P(), TOTAL, "1부")
note(s, "[1.5분]\n"
     "대본: 1부 시작 — OOD와 외삽부터. **왼쪽 그림 (a)** — OOD. train과 test **분포**가 다릅니다. "
     "오른쪽 꼬리 query는 train이 본 적 없는 영역 — OOD. "
     "**오른쪽 그림 (b)** — 외삽. 훈련 점들의 convex hull 안/밖으로 판별. "
     "★ 표시는 hull **밖** extrap — 오늘의 좌표. "
     "보라 ■는 hull **안**인데 concept shift — OOD지만 extrap 판별은 안 됩니다. "
     "extrap은 종종 OOD 안쪽이지만, OOD 전체가 extrap은 아닙니다. "
     "오늘 50분은 **회귀 외삽** — hull·연속 y·가정. "
     "IRM·DomainBed는 **분류 DG** 쪽 — 오늘 범위 밖. 다음 장 보간/외삽.")
takeaway(s, "OOD=분포 밖(넓음) · 회귀 외삽=hull 밖 입력(오늘)",
         "분류·DG(OOD) ≠ 오늘 — 연속값 extrap")

# ── S6 보간 vs 외삽 ──
s = blank()
content_header(s, "1부 · 외삽 정의", "외삽 = 훈련 범위(hull) 밖에서의 예측")
if (ASSETS / "fig_interp_extrap.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_interp_extrap.png"), Inches(0.30), Inches(1.20), width=Inches(8.1))
card_text(
    s, Inches(0.35), Inches(4.65), Inches(8.05), Inches(1.55),
    "적용 예",
    [
        "키 160~180cm로 학습 → 175cm 예측 = 보간 (안전)",
        "같은 모델에 195cm 예측 = 외삽 (본 적 없음)",
        "여름 데이터로 학습한 전력수요 모델 → 겨울 예측 = 외삽 (시간축도 마찬가지)",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.60), Inches(1.20), Inches(4.35), Inches(2.6),
    "판별 기준",
    [
        "입력 x가 훈련 범위 안 → 보간",
        "입력 x가 훈련 범위 밖 → 외삽",
        "같은 모델도 입력마다",
        "  보간·외삽 동시 가능",
        "먼저 확인: '이 입력, 범위 밖?'",
        "피할 말: '모델이 외삽?' (입력 기준)",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.60), Inches(3.95), Inches(4.35), Inches(2.2),
    "그림·기호 약속",
    [
        "음영 = 훈련 범위 (hull)",
        "음영 밖 = 외삽 구간",
        "이후 모든 장표에서 동일 좌표계",
        "논문: Bartley 2019 (훈련 범위 정의)",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "1부")
note(s, "[1.5분]\n"
     "대본: 정의부터. 외삽은 '훈련 범위 밖에서의 예측'입니다. "
     "왼쪽 그림 — 같은 모델인데 **예측할 입력**이 훈련 범위 안이면 보간, 밖이면 외삽입니다. "
     "그래서 '모델이 외삽되나?'가 아니라 '**지금 넣는 입력이 범위 밖인가?**'가 먼저입니다. "
     "적용 예 — 키 160~180으로 학습한 몸무게 모델에 175를 물으면 보간, 195를 물으면 외삽입니다. "
     "여름 데이터로 겨울 수요를 예측하는 것도 외삽입니다 — 시간축도 똑같습니다. "
     "약속 하나: 오늘 모든 그림에서 음영이 훈련 범위, 그 밖이 외삽 영역입니다.")
takeaway(s, "입력이 훈련 범위 안 → 보간, 밖 → 외삽",
         "같은 모델도 **입력이 범위 안/밖**에 따라 달라짐")

# ── S5 훈련 범위의 정확한 정의: convex hull ──
s = blank()
content_header(s, "1부 · 외삽 정의", "'훈련 범위'를 수학적으로 — convex hull")
if (ASSETS / "fig_convex_hull.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_convex_hull.png"), Inches(0.30), Inches(1.20), width=Inches(6.8))
card_text(
    s, Inches(0.35), Inches(5.20), Inches(7.35), Inches(1.00),
    "예시 — 조합의 함정",
    [
        "20~30°C·습도 40~60%로 학습 → (35°C, 50%) 예측 = 밖",
        "함정: (22도, 59%)처럼 '변수별로는 범위 안'이어도 조합이 훈련 범위 밖일 수 있다",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.95), Inches(1.20), Inches(5.0), Inches(2.80),
    "convex hull — ① 개념 · ② 수식",
    [
        "① 개념 — 훈련 row(●)들을 고무줄로 감싼 영역",
        "   (외곽 끝 + 그 안 전체, 피처 하나씩 아님)",
        "② 수식 — Conv(X_train) = { x = Σ α_i x_i |",
        "   x_i = row (모든 컬럼),  α_i ≥ 0,  Σ α_i = 1 }",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.95), Inches(4.10), Inches(5.0), Inches(2.22),
    "③ 의미 · ④ 왜 정의?",
    [
        "hull 안 (x ∈) → 보간 — 훈련점 섞기로 만들 수 있음",
        "hull 밖 (x ∉) → 외삽 — 본 적 없는 입력·조합",
        "Σ=1 → 섞는 비율 100%, 고무줄 안 · Σ≠1 → 밖",
        "④ 왜 정의? — '범위 밖'을 느낌→수학 판정",
        "   min/max·감각 대신 공식 기준 → 3부 검증 가능",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "bartley", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · Bartley")
footer(s, P(), TOTAL, "1부")
note(s, "[1.5분]\n"
     "대본: 오른쪽 네 단계입니다. "
     "① 개념 — 훈련 row(●)들을 고무줄로 감싼 영역, 외곽+안. "
     "② 수식 — Conv(X_train) = Σ α_i x_i, α_i ≥ 0, Σ α_i = 1. "
     "③ 의미 — hull 안=보간, 밖=외삽. Σ=1은 섞기 100%라 고무줄 안만. "
     "④ 왜 정의? — '범위 밖'을 느낌이 아니라 수학으로 판정해야 3부 검증이 됩니다. "
     "min/max나 '비슷해 보여'로는 부족합니다. "
     "왼쪽 함정 — 변수별로는 범위 안인데 조합은 hull 밖일 수 있습니다.")
takeaway(s, "훈련 범위 = convex hull — 안이면 보간, 밖이면 외삽",
         "느낌이 아니라 판정 가능한 기준 (3부 검증의 토대)")

# ── S6 고전 사례: 다항식 blow-up ──
s = blank()
content_header(s, "1부 · 대표 사례", "훈련구간 적합도는 외삽 성능을 보증하지 않는다")
if (ASSETS / "fig_poly_extrap.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_poly_extrap.png"), Inches(0.30), Inches(1.20), width=Inches(7.3))
card_text(
    s, Inches(0.35), Inches(5.15), Inches(7.85), Inches(1.05),
    "예시 — 숫자로",
    [
        "설비 열화 곡선 100일치에 고차 다항식 적합: 훈련 오차 0.1% (완벽해 보임)",
        "130일째 예측: 물리적으로 불가능한 음수 출력 — 안의 성적표는 밖을 모른다",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.40), Inches(1.20), Inches(4.55), Inches(2.75),
    "관찰",
    [
        "4차 다항식 · 훈련구간 오차 ≈ 0",
        "경계를 넘는 즉시 발산",
        "차수 높일수록 안은 더 잘 맞고",
        "  밖은 더 크게 터진다",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.40), Inches(4.10), Inches(4.55), Inches(2.05),
    "여기서 고정할 교훈",
    [
        "'안에서 잘 맞는다'와",
        "  '밖에서 맞는다'는 다른 질문",
        "→ 왜 다른가? 2부에서 해부",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "1부")
note(s, "[1.5분]\n"
     "대본: 밖에서 무슨 일이 생기는지 대표 사례입니다. 4차 다항식이 훈련 구간 안에서는 사실상 완벽합니다 — 오차가 0에 가깝습니다. "
     "그런데 경계를 넘는 순간 폭발합니다. 차수를 높이면 안은 더 잘 맞는데 밖은 더 크게 터집니다. "
     "여기서 교훈 하나만 새기고 2부로 갑니다: '안에서 잘 맞는다'와 '밖에서 맞는다'는 서로 다른 질문이다. "
     "왜 다른 질문일 수밖에 없는지 — 그게 2부입니다.")
takeaway(s, "훈련구간 오차 ≈ 0 이어도 경계 밖에서 즉시 발산",
         "안의 성적은 밖을 보증하지 않는다 — 이유는 2부에서")

# ════════ 2부 · 외삽 실패 원인 ════════
section_slide(
    "2부 · 외삽 실패 원인 (4)",
    "훈련 범위 밖에서 예측이 깨지는 이유",
    "2", "15분", P(), TOTAL,
    points=[
        "이유 1·2: 데이터가 함수를 못 정하고, 고차원에선 거의 전부가 밖",
        "이유 3: 범위 밖에선 '모름'(불확실성)이 커진다",
        "이유 4: 신경망도 밖에서는 직선이 된다",
    ],
)
note(prs.slides[-1], "[20초]\n대본: 2부, **외삽 실패 원인**입니다. '한계'가 아니라 — 밖에서 왜 틀리는지 네 가지. "
     "식별 불가(데이터만으론 못 고름), 고차원에서 범위 밖, 밖에서 불확실성 폭발, ReLU 직선. 하나씩 보겠습니다.")

# ── S8 이유 1: 데이터가 함수를 못 정한다 ──
s = blank()
content_header(s, "2부 · 실패 원인 ①", "데이터가 함수를 정해 주지 않는다")
if (ASSETS / "fig_identifiability.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_identifiability.png"), Inches(0.30), Inches(1.20), width=Inches(7.3))
card_text(
    s, Inches(0.35), Inches(5.15), Inches(7.85), Inches(1.05),
    "예시 — 배터리 수명",
    [
        "100사이클까지의 용량 데이터: '선형 열화'와 '지수 열화' 둘 다 완벽 적합",
        "500사이클 수명 예측은 두 가설이 2배 차이 — 데이터는 어느 쪽인지 말해주지 않는다",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.40), Inches(1.20), Inches(4.55), Inches(2.75),
    "같은 점, 다른 함수",
    [
        "훈련점을 지나는 함수(가설)는 무한히 많다",
        "그림 — 직선·지수: 훈련 구간(OBSERVED) 안",
        "  에서는 둘 다 ● 통과 → 구분 불가",
        "밖(UNIDENTIFIED)으로 가면 전부 갈라진다",
        "데이터를 더 모아도 (안에서만 늘면)",
        "  구별 불가는 그대로 (Pfister 2024)",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.40), Inches(4.10), Inches(4.55), Inches(2.05),
    "핵심 결론 (S03과 연결)",
    [
        "밖을 지탱하는 것은",
        "  데이터가 아니라 가정이다",
        "가정 = 밖 거동에 대한 사전 제약",
        "→ 어느 함수족을 믿을지 선택해야 함",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "pfister", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · Pfister")
footer(s, P(), TOTAL, "2부")
note(s, "[2분] ★ 발표 전체에서 가장 중요한 장.\n"
     "대본: 첫 번째 이유이자 오늘의 심장입니다. 그림 — 직선 가설과 지수 가설, "
     "훈련 구간(OBSERVED) 안에서는 둘 다 ●를 통과합니다. 구분이 안 됩니다. "
     "밖(UNIDENTIFIED)으로 나가는 순간 전부 다른 답을 냅니다. "
     "아무 함수나 맞는 게 아니라, 훈련점을 지나는 가설이 (원칙적으로) 무한히 많다는 뜻입니다. "
     "그럼 데이터를 더 모으면 되지 않나? — 훈련 범위 안에서만 늘어나면 소용없습니다. "
     "Pfister 2024: 범위 밖 추론은 추가 가정 없이는 원리적으로 불가능. "
     "배터리 예 — 100사이클까지 선형·지수 열화 둘 다 완벽 적합, "
     "500사이클 수명은 2배 차이. 데이터는 어느 쪽인지 말해주지 않습니다. "
     "밖을 지탱하는 것은 데이터가 아니라 가정입니다.")
takeaway(s, "훈련 구간 안에서는 맞춰지지만, 밖에서 갈라진다 — 데이터로는 못 고른다",
         "밖을 지탱하는 것은 데이터가 아니라 가정이다 (핵심 문장)")

# ── S9 이유 2: 고차원에선 거의 전부가 밖 ──
s = blank()
content_header(s, "2부 · 실패 원인 ②", "고차원에선 새 데이터가 거의 전부 범위 밖")
if (ASSETS / "fig_hull_dimension.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_hull_dimension.png"), Inches(0.30), Inches(1.20), width=Inches(8.3))
card_text(
    s, Inches(0.35), Inches(4.45), Inches(8.25), Inches(1.75),
    "예시 — 센서 20채널이면",
    [
        "각 채널이 훈련 범위 안에 있을 확률 95%라 해도, 20채널 '조합'이 전부 안일 확률 ≈ 0.95²⁰ ≈ 36%",
        "실제 훈련 범위는 축 정렬 상자보다 훨씬 좁다 → 조합 기준으로는 대부분이 이미 밖",
        "이미지(수만 차원)·다변량 시계열 — 여러분이 다루는 데이터가 정확히 이 영역",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.80), Inches(1.20), Inches(4.15), Inches(2.75),
    "차원의 저주 (Bartley 2019)",
    [
        "2차원(왼쪽): 범위가 커 보인다",
        "차원↑ → 새 입력이 훈련 범위 안에",
        "  있을 확률이 지수적으로 0으로",
        "  (오른쪽 곡선)",
        "이미지·다채널 센서 = 이미 고차원",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.80), Inches(4.10), Inches(4.15), Inches(2.05),
    "함의",
    [
        "'우리 테스트는 안이겠지' 근거 없음",
        "외삽은 예외가 아니라 기본값",
        "→ 특수한 예외가 아니라 일반적 상황",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "bartley", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · Bartley")
footer(s, P(), TOTAL, "2부")
note(s, "[1.5분]\n"
     "대본: 두 번째 이유 — 규모의 문제입니다. '데이터가 충분히 많으니 테스트는 범위 안이겠지'라고 생각하기 쉽습니다. "
     "2차원에서는 맞는 직관입니다, 왼쪽 그림처럼 훈련 범위가 꽤 넓거든요. "
     "그런데 오른쪽 곡선 — 차원이 올라가면 **새 입력**이 훈련 범위 안일 확률이 지수적으로 0에 수렴합니다. "
     "100차원이면 사실상 0입니다. 체감 예시 — 센서 20채널에서 채널별로 95% 확률로 범위 안이라 해도, "
     "20개 조합이 전부 안일 확률은 0.95의 20제곱, 36%입니다. 실제 훈련 범위는 그보다 훨씬 좁으니 대부분이 이미 밖입니다. "
     "이미지, 다채널 센서 시계열 — 전부 고차원입니다. "
     "결론: 외삽은 예외 상황이 아니라 기본값입니다.")
takeaway(s, "차원 ↑ ⇒ 훈련 범위 안 확률 → 0 — 고차원 테스트는 대부분 이미 외삽",
         "외삽은 예외가 아니라 기본값이다")

# ── S10 이유 3: 밖에서 인식적 불확실성이 지배 ──
s = blank()
content_header(s, "2부 · 실패 원인 ③", "범위 밖에선 '모름'(불확실성)이 커진다")
if (ASSETS / "fig_error_decomp.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_error_decomp.png"), Inches(0.30), Inches(1.12), width=Inches(7.75))
card_text(
    s, Inches(0.35), Inches(4.52), Inches(7.70), Inches(1.68),
    "그림 읽기",
    [
        "음영 = 훈련 범위 · 안: total error 작음",
        "(a) 세 색 쌓임 = (b) 펼친 곡선 (같은 분해)",
        "밖으로 갈수록 주황(모름·Variance)이 total을 끌어올림",
        "회색(잡음)은 어디서나 비슷 · 파랑(편향)은 완만",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.35), Inches(1.12), Inches(4.60), Inches(2.85),
    "왜 세 가지로 쪼개 보나?",
    [
        "오차 원인이 다름 → 고치는 방법도 다름",
        "잡음(Noise): 측정·센서·환경 — 본질적 랜덤",
        "  → 과녁: 손 떨림 · 어디서나 비슷 · 못 줄임",
        "편향(Bias): 모델 가정·구조가 틀림 — 체계적 오차",
        "  → 과녁: 조준 틀림 · 가정·모델 바꾸면 줄임",
        "모름(Var/epistemic): 데이터 없어 예측이 흔들림",
        "  → 과녁: 위치 모름 · 그 구간 데이터 있으면 줄임",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(8.35), Inches(4.05), Inches(4.60), Inches(2.25),
    "범위 밖 · 분해 목적",
    [
        "훈련 범위 밖 = 정의상 데이터 없음",
        "→ '모름'만 구조적으로 커짐 (그림 주황 폭발)",
        "→ 잡음·편향만으로는 설명 안 됨",
        "분해하는 이유: 어디서·왜 틀리는지 진단",
        "  (가정 문제? 데이터 부족? 측정 한계?)",
        "→ 3부 UQ: '모름'을 신호로 쓰기",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "2부")
note(s, "[1.5분]\n"
     "대본: 밖에서 왜 불확실해지는지 — 오차를 세 덩어리로 쪼개 봅니다. "
     "왜 쪼개나? 원인마다 대처가 다르기 때문입니다. "
     "잡음 — 측정·센서·환경, 본질적 랜덤, 과녁 비유로 손 떨림, 줄일 수 없습니다. "
     "편향 — 모델 가정·구조가 틀린 체계적 오차, 조준 틀림, 가정을 고치면 줄입니다. "
     "모름 — 데이터가 없어서 예측이 흔들리는 것, 과녁 위치를 모르는 것, "
     "그 구간에 데이터를 모으면 줄일 수 있습니다. "
     "그림 — 음영 안에서는 total error가 작고, 밖으로 가면 주황 '모름'이 폭발합니다. "
     "훈련 범위 밖은 정의상 데이터가 없으니 '모름'만 구조적으로 못 줄입니다. "
     "분해 목적 — 어디서 왜 틀리는지 진단하는 것, 3부 UQ에서 이 신호를 씁니다.")
takeaway(s, "훈련 범위 밖 — '모름'(epistemic)이 커짐 · 데이터로는 못 줄임",
         "이 약점을 신호로 활용하는 방법이 3부 끝(UQ)에 나온다")

# ── S11 이유 4: 신경망도 직선 (Xu Thm.1) ──
s = blank()
content_header(s, "2부 · 실패 원인 ④", "ReLU NN — 밖에서는 직선 (Xu 2021)")
pic = add_paper_fig(s, "xu_fig1_relu_extrap.png", Inches(0.28), Inches(1.18), width=Inches(7.4),
                    caption="논문: Xu et al., How Neural Networks Extrapolate: From Feedforward to Graph Neural Networks (ICLR 2021), Fig.1")
if (ASSETS / "fig_xu_sin_tangent.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_xu_sin_tangent.png"), Inches(0.28), Inches(3.78), width=Inches(7.4))
card_text(
    s, Inches(0.35), Inches(2.92), Inches(7.35), Inches(0.78),
    "Fig.1 읽기 — φ = 활성화 함수 (activation)",
    [
        "파란 = 훈련 구간(안) · 바깥 격자 = ReLU MLP 예측",
        "3D 3장: 타깃만 다름 — 공통으로 밖은 평면(plane) · 2D: sin vs 직선",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.95), Inches(1.18), Inches(5.0), Inches(2.35),
    "Xu의 정리 1 (ReLU)",
    [
        "ReLU MLP — 입력 멀어지면 방향마다 직선",
        "깊이·너비 키워도 동일 (Thm.1, 증명)",
        "위: 논문 Fig.1 · 아래: sin 학습 MLP 예시",
        "→ 실무 기본(ReLU)의 밖 거동을 정식화",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.95), Inches(3.63), Inches(5.0), Inches(2.68),
    "예시 해석 · Q 다른 활성화 φ?",
    [
        "φ = 활성화 함수 — '밖은 이 모양' 선언 (S16)",
        "음영 안: 완벽 적합 · 밖: 접선 직선, 다시 안 굽음",
        "재해석: NN 외삽 불능 ❌ — '밖 가정'을 이미 함",
        "Q: ReLU만? — tanh/sig: 밖 포화 · sin/cos: 주기(맞으면)",
        "  GELU 등: ReLU 근사 → 직선 쪽 · φ 안 맞으면 10²–10³",
        "φ만으론 부족(S09) — 맞춰도 검증·UQ · 자세히 → S16",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "xu", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · Xu")
footer(s, P(), TOTAL, "2부")
note(s, "[2.5분] ★ 2부 클라이맥스.\n"
     "대본: 위 Fig.1부터 — φ는 활성화 함수(activation)입니다. 파란 부분이 훈련 구간, "
     "안에서는 복잡한 타깃을 ReLU MLP가 잘 맞춥니다. 바깥 회색 격자면을 보세요 — "
     "3D 세 장 모두 밖으로 나가면 평평한 plane으로 이어집니다. 타깃 함수만 다르고 "
     "밖 거동은 같습니다. 맨 오른쪽 2D는 단면 — sin을 배워도 밖은 직선. "
     "아래 큰 그림도 같은 메시지. Xu 정리 1: ReLU MLP는 멀어지면 방향마다 직선. "
     "외삽 불능이 아니라 '밖은 직선' 가정. Q ReLU만? tanh는 포화, sin/cos는 주기, "
     "φ=밖 모양 선언. φ 맞춰도 S09 한계·검증 필요 — S16. "
     "가정을 어차피 하고 있다면, 내가 아는 것으로 골라 넣자.")
takeaway(s, "ReLU NN — 밖에서 직선 (Xu Thm.1) · φ마다 밖 가정 다름",
         "NN은 외삽 불능이 아니라 '직선'이라는 가정을 이미 하고 있다")

# ── S12 이유 4 메커니즘 + 2부 정리 ──
s = blank()
content_header(s, "2부 · 실패 원인 ④", "왜 직선인가 — 3단 논법 · 2부 정리")
if (ASSETS / "fig_relu_affine.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_relu_affine.png"), Inches(1.62), Inches(1.15), width=Inches(10.1))
steps = [
    ("① 뉴런 고정 (그림 a)", ["ReLU 뉴런은 딱 한 번만 꺾인다", "입력이 커지면 마지막 꺾임을 지나", "  켜짐/꺼짐이 영원히 고정"], TEAL),
    ("② 층이 선형화", ["켜짐이 고정된 네트워크는", "그냥 행렬 곱셈의 연속", "비선형 굽힘이 전부 사라짐"], TEAL),
    ("③ 직선 수렴 (그림 b)", ["고정된 직선들의 합 = 직선 하나", "굽힘은 전부 데이터 근처에만 산다", "→ 멀리서 보면 평평 (Thm.1)"], CORAL),
]
for i, (title, bullets, accent) in enumerate(steps):
    card_text(s, Inches(0.30 + i * 4.28), Inches(4.72), Inches(4.15), Inches(1.55), title, bullets, accent, dense=True)
footer(s, P(), TOTAL, "2부")
note(s, "[1.5분]\n"
     "대본: 왜 직선이 되는지 3단으로. 첫째 — 입력이 커지면 각 ReLU 뉴런의 켜짐/꺼짐이 고정됩니다. "
     "둘째 — 켜짐이 고정된 네트워크는 그냥 행렬 곱셈입니다. 셋째 — 행렬 곱을 아무리 쌓아도 직선 하나입니다. "
     "직관적으로: 네트워크의 비선형성은 훈련 범위 근처에만 살고, 멀리서 보면 전부 평평합니다. "
     "2부 정리하겠습니다. 데이터는 함수를 못 정하고, 고차원 테스트는 어차피 대부분 밖이고, 밖에서는 '모름'(불확실성)이 커지고, "
     "신경망조차 밖에서는 직선입니다. 그러니 질문을 바꿔야 합니다 — '외삽이 되나?'가 아니라 "
     "'내 모델이 하고 있는 가정이 내 문제와 맞나?'. 가정을 직접 골라서 넣는 방법들, 3부입니다.")
takeaway(s, "2부 정리 — 질문을 바꿔라: '외삽 되나?' → '모델의 가정이 내 문제와 맞나?'",
         "가정을 직접 골라 넣으면? → 3부")

# ════════ 3부 · 외삽 대응법 + 성능 검증 ════════
section_slide(
    "3부 · 회귀 외삽 대응법 + 성능 검증",
    "가정 넣기(대응법) → '밖에서도 된다' 주장 확인(검증)",
    "3", "23분", P(), TOTAL,
    points=[
        "대응법: **회귀** — 사전 지식별 가정 4종",
        "성능 검증: hull 밖 holdout · baseline · UQ",
    ],
)
note(prs.slides[-1], "[20초]\n대본: 3부 — **회귀** 외삽 대응법과 검증. "
     "앞 — EQL·CMNN·Physics-ML·UQ. 뒤 — hull 밖 테스트·공정 baseline·UQ 동작 확인.")

# ── S14 방법 지도: 가정 스펙트럼 ──
s = blank()
content_header(s, "3부 · 대응법", "방법 지도 — **회귀** 외삽 · 사전 지식별 가정")
if (ASSETS / "fig_assumption_spectrum.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_assumption_spectrum.png"), Inches(1.55), Inches(1.12), height=Inches(3.20))
# 스펙트럼 아래: 방법 4종을 이름·이득·대가로 명시 (장표만으로 읽히게)
card_text(
    s, Inches(0.35), Inches(4.45), Inches(3.05), Inches(1.75),
    "① 함수족 (강)",
    [
        "EQL: 수식 유닛 내장",
        "NALU: 산술로 가중치 제약",
        "이득: 전역 유효 가능",
        "대가: 가정이 틀릴 때 전역 실패",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(3.55), Inches(4.55), Inches(3.05), Inches(1.65),
    "② 방향",
    [
        "CMNN: 단조를 구조로 강제",
        "실무 지식의 대부분 수준",
        "이득: 밖에서도 방향 보장",
        "대가: 구조 설계 비용",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(6.75), Inches(4.55), Inches(3.05), Inches(1.65),
    "③ 물리 (회귀)",
    [
        "**Physics-ML**: Aykol · Li",
        "L = L_data + λ·L_constraint",
        "센서→품질 **해당** (PINN ✗)",
        "대가: soft loss · **밖 검증**",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(9.95), Inches(4.55), Inches(3.0), Inches(1.65),
    "④ 무가정",
    [
        "UQ: '모름'을 경보로",
        "기권: 신뢰 낮으면 예측 안 함",
        "이득: 과신 오답 방지",
        "대가: 답변율 감소",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "3부")
note(s, "[1.5분]\n"
     "대본: 3부 지도입니다. 가로축 = **내가 아는 사전 지식의 강도**입니다. "
     "왼쪽 끝 — 타깃 함수의 형태를 안다: EQL과 NALU. 올바르면 전 영역에서 유효하고, 가정이 틀리면 전역 실패입니다. "
     "두 번째 — 방향성만 안다: 단조 신경망. 실무에서 가장 흔한 지식 수준입니다. "
     "세 번째 — **물리 지식**: 회귀면 **Physics-ML / physics loss** (Aykol·Li·Muckley — 다음 장). "
     "**PDE+장** u(x,t)만 Raissi식 PINN. "
     "센서→품질 = physics **feature·soft loss** 또는 CMNN·UQ. "
     "오른쪽 끝 — 사전 지식이 없다: 이때는 불확실성(UQ)과 기권(예측 안 함)으로 유보합니다. "
     "순서가 중요합니다. 방법을 먼저 고르는 것이 아니라, 가용 지식이 먼저이고 방법이 뒤따릅니다.")
takeaway(s, "지식 우선, 방법 후행 — 가용 사전 지식이 방법을 정한다",
         "강한 가정 = 높은 이득 + 높은 위험")

# ── S15 입구: 활성화 함수 교체 ──
s = blank()
content_header(s, "3부 · 입구", "최소 개입의 가정 주입 — 활성화 함수를 타깃에 맞추기 (Xu Fig.5)")
pic = add_paper_fig(s, "xu_fig5_activations.png", Inches(0.28), Inches(1.15), width=Inches(6.9),
                    caption="논문: Xu et al., How Neural Networks Extrapolate (ICLR 2021), Fig.5")
if (ASSETS / "fig_activation_match.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_activation_match.png"), Inches(0.28), Inches(3.75), width=Inches(6.9))
card_text(
    s, Inches(7.45), Inches(1.15), Inches(5.5), Inches(2.55),
    "실험 요지 (위: 논문 원본)",
    [
        "타깃이 주기함수 → 활성화를 cos으로",
        "  바꾸면 외삽이 된다",
        "타깃이 tanh형 → tanh 활성화가 성공",
        "짝이 틀리면 오차 100~1000배",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.45), Inches(3.85), Inches(5.5), Inches(2.30),
    "예시 — 같은 sin 타깃",
    [
        "(a) ReLU: 밖에서 직선 → 파동 놓침",
        "(b) cos 활성화: 밖에서도 파동 유지",
        "활성화 선택 = '밖은 이 모양' 선언",
        "가정에는 대가가 있다: 불일치 시 오차 수백 배",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "xu", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · Xu")
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: 본격적인 방법 전에, 가장 단순한 가정 주입부터 봅니다 — 활성화 함수 교체입니다. "
     "Xu의 Fig.5: 타깃이 주기함수일 때 활성화를 cos으로 바꾸면 외삽이 됩니다. tanh 타깃엔 tanh가 됩니다. "
     "반대로 짝이 틀리면 오차가 백 배, 천 배로 뜁니다. "
     "즉 활성화 선택이 곧 '밖은 이런 모양일 것'이라는 선언입니다. "
     "가정에는 대가가 있다 — 이 원리를 들고 방법 1로 갑니다.")
takeaway(s, "활성화와 타깃이 같은 족일 때만 외삽 성립 — 틀린 짝은 오차 수백 배",
         "활성화 선택 = 함수 모양 가정의 선언. 가정에는 대가가 있다")

# ── S16 방법 1a: EQL ──
s = blank()
content_header(s, "3부 · 방법 1", "함수 모양을 안다 (a) — EQL: 수식을 통째로 배우기")
pic = add_paper_fig(s, "eql_fig1_architecture.png", Inches(0.28), Inches(1.15), width=Inches(6.4),
                    caption="논문: Martius & Lampert, Extrapolation and Learning Equations (arXiv 2016), Fig.1")
if (ASSETS / "fig_eql_example.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_eql_example.png"), Inches(0.28), Inches(3.65), width=Inches(6.9))
card_text(
    s, Inches(7.45), Inches(1.15), Inches(5.5), Inches(2.85),
    "원리 (위: 논문 구조도)",
    [
        "뉴런 자리에 sin·cos·곱셈 같은",
        "  수식 유닛을 심는다",
        "학습이 끝나면 네트워크 자체가",
        "  하나의 닫힌 수식: ŷ ≈ sin(1.01x)",
        "배운 식이 참 식과 같은 족이면",
        "  → 훈련 구간 아닌 전 영역에서 유효",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.45), Inches(4.15), Inches(5.5), Inches(2.0),
    "예시 + 대가",
    [
        "초록(EQL): 식이라서 멀리서도 파동 유지",
        "주황(MLP): 같은 데이터인데 직선",
        "대가: 유닛 목록 밖 함수 표현 불가 ·",
        "  희소화 조절 까다로움 · 고차원 취약",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "eql", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · EQL")
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: 방법 1 — 타깃이 '닫힌 수식'이라는 확신이 있는 경우입니다. "
     "EQL은 sin, cos, 곱셈 같은 수식 유닛을 뉴런 자리에 심습니다. 학습이 끝나면 네트워크 자체가 하나의 식이 됩니다. "
     "배운 식이 참 식과 같은 족이면, 훈련 구간이 아니라 정의역 '전체'에서 유효합니다 — 아까 MLP가 직선이 되던 것과 정반대입니다. "
     "대가: 유닛 목록에 없는 함수는 아예 표현 못 하고, 희소화 조절이 까다롭고, 고차원·잡음에 약합니다.")
takeaway(s, "배운 식이 참 식과 같은 족이면 전 영역 유효 — MLP 직선화와 정반대",
         "대가: 유닛 밖 함수 표현 불가 · 희소화 · 고차원 취약")

# ── S17 방법 1b: NALU ──
s = blank()
content_header(s, "3부 · 방법 1", "함수 모양을 안다 (b) — NALU: 산술만 하도록 묶기")
pic = add_paper_fig(s, "nalu_fig2_architecture.png", Inches(0.28), Inches(1.15), width=Inches(6.2),
                    caption="논문: Trask et al., NALU (NeurIPS 2018), Fig.2 (a) NAC · (b) NALU · arXiv:1808.00508")
if (ASSETS / "fig_nalu_example.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_nalu_example.png"), Inches(0.28), Inches(4.00), width=Inches(6.6))
card_text(
    s, Inches(7.45), Inches(1.15), Inches(5.5), Inches(3.05),
    "Fig.2 — (a) NAC · (b) NALU",
    [
        "(a) NAC — +/− 부품",
        "  W = tanh(Ŵ)⊙σ(M̂) → w≈+1,0,−1",
        "  y = W·x (matmul) · w=+1 더함 · w=−1 뺌",
        "(b) NALU — NAC + × + 스위치",
        "  가운데: NAC · 보라: log→matmul→exp (×÷)",
        "  g=σ → y = g·(덧셈)+(1−g)·(곱셈)",
        "  (a)는 부품 · (b)는 +단·×단 완성",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.45), Inches(4.35), Inches(5.5), Inches(1.85),
    "예시 — 덧셈 학습",
    [
        "0~10으로 a+b 학습 → 10,000까지 테스트",
        "MLP 오차 3,100배 폭발 · NALU는 평탄",
        "대가: w·g ±1 수렴 보장 ❌ · 게이트 불안정",
        "  타깃이 산술 아니면 MLP로 붕괴 · y=a+b 알면 코드가 나음",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "nalu", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · NALU")
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: Fig.2 두 장입니다. (a) NAC — +/− 부품. 아래 W 공장에서 W = tanh(Ŵ)⊙σ(M̂)로 "
     "w를 ±1,0에 가깝게 만들고, 위에서 y = W·x로 덧셈·뺄셈만 합니다. ×는 선형합으로 불가합니다. "
     "(b) NALU — NAC에 log-exp 곱셈 경로(보라)와 게이트 g(빨강)를 더한 완성입니다. "
     "y = g·(덧셈)+(1−g)·(곱셈). a+b만이면 (a)로 충분, a×b면 (b)가 필요합니다. "
     "0~10으로 a+b 학습 후 10,000까지 테스트하면 MLP는 오차가 폭발하고 NALU는 평탄합니다. "
     "단, ±1·0 수렴은 목표이지 보장은 아니고, 게이트 불안정·비산술 타깃이면 MLP로 붕괴합니다. "
     "식 전체를 알면 NALU 없이 코드 한 줄이 낫습니다. "
     "방법 1 요약: 함수 모양 가정은 수익도 최대, 리스크도 최대입니다.")
takeaway(s, "가중치를 산술로 묶기 → 작은 수 훈련, 큰 수 일반화",
         "방법 1 요약: 함수 모양 가정 = 최대 수익 · 최대 리스크")

# ── S18 방법 2a: 단조 — 문제 ──
s = blank()
content_header(s, "3부 · 방법 2", "방향만 안다 (a) — 단조를 강제하면 표현력이 죽는다")
pic = add_paper_fig(s, "runje_fig1_cubic.png", Inches(0.28), Inches(1.15), width=Inches(7.4),
                    caption="논문: Runje & Shankaranarayana, Constrained Monotonic Neural Networks (ICML 2023), Fig.1 · arXiv:2205.11775")
card_text(
    s, Inches(0.30), Inches(3.85), Inches(7.35), Inches(2.35),
    "Fig.1 — y = x³ (단조 증가 함수)",
    [
        "(a) 무제약 NN: 적합은 되나 단조 보장 없음",
        "(b) 가중치 전부 양수: 단조는 되나 사실상 직선 — x³ 실패",
        "(c) CMNN: 단조 + 비선형 동시 (다음 장)",
        "실무: 보험료↑·전압↑·수요↓ — 방향은 아는데 수식은 모를 때",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(7.85), Inches(1.15), Inches(5.1), Inches(5.05),
    "단조(monotonic) 가정",
    [
        "입력↑ → 출력이 절대 안 줄어듦",
        "  (또는 반대 방향)",
        "실무 지식 대부분이 이 수준",
        "",
        "딜레마",
        "순진한 강제(가중치 자르기)는",
        "  표현력을 죽인다 — (b)",
        "단조 vs 표현력, 둘 다?",
        "→ 다음 장 CMNN",
    ],
    TEAL, dense=True,
)
pdf_btn(s, "runje", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · CMNN")
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: 방법 2 — 현실적으로 가장 중요한 방법입니다. 수식까진 몰라도 '가격이 오르면 수요는 준다', "
     "'충전량이 늘면 전압은 오른다' 같은 방향은 아는 경우가 대부분이죠. "
     "그런데 방향을 강제하는 순진한 방법 — 가중치를 전부 양수로 자르는 것 — 은 그림 (b)처럼 표현력을 죽입니다. "
     "단조는 얻었는데 사실상 직선이 돼서 x³조차 못 맞춥니다. "
     "단조성과 표현력이 상충하는 것처럼 보입니다. 이걸 푸는 게 다음 장입니다.")
takeaway(s, "가중치 자르기: 단조는 얻지만 x³조차 실패 — 표현력 붕괴",
         "단조 vs 표현력의 상충 — 해법은 다음 장")

# ── S19 방법 2b: 단조 — 해법 ──
s = blank()
content_header(s, "3부 · 방법 2", "방향만 안다 (b) — 해법: 구조 자체가 단조이게 만들기 (CMNN)")
pic = add_paper_fig(s, "runje_fig3_mono_unit.png", Inches(0.28), Inches(1.15), height=Inches(3.1),
                    caption="논문: Runje & Shankaranarayana, Constrained Monotonic Neural Networks (ICML 2023), Fig.3")
card_text(
    s, Inches(0.30), Inches(4.90), Inches(6.9), Inches(1.30),
    "실무 예시 — 배터리 제어기",
    [
        "충전량(SOC) → 전압(OCV) 곡선을 CMNN으로 학습",
        "훈련에 없던 고충전 조건 예측에도 '전압 급락'은 구조상 불가",
        "→ 이 보장 하나로 제어기 안전 인증이 가능해진다",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.55), Inches(1.15), Inches(5.4), Inches(2.85),
    "Fig.3 — 한 층 흐름",
    [
        "x → |W|_t·x+b → split(s) → concat → W_out → ŷ",
        "t: 입력 feature마다 단조 방향 (+1 / −1 / 0)",
        "  +1 → x_j↑일 때 출력↑ (예: SOC↑ → OCV↑)",
        "  −1 → x_j↑일 때 출력↓",
        "  0 → 그 feature는 단조 제약 없음",
        "  |W|_t: t=±1일 때 w 부호만 고정 (전체 |W| 아님)",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.55), Inches(4.15), Inches(5.4), Inches(2.0),
    "s — 활성화 배치 (곡선 모양)",
    [
        "s = (볼록 ρ̊, 오목 ρ̂, 포화 ρ̃) — 합 = hidden m",
        "  ρ̊: h>0 구간에서 증가 (ReLU 계열)",
        "  ρ̂: h<0 구간에서 증가 (거울 ReLU 계열)",
        "  ρ̃: 위·아래 평평 (x³·배터리는 보통 0)",
        "예 m=10, s=(4,4,0) → ρ̊ 4개 + ρ̂ 4개",
        "뉴런 번호로 φ 고정 — x가 +/−냐로 바뀌지 않음",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "runje", Inches(11.4), Inches(0.28), w=Inches(1.55), label="PDF · CMNN")
footer(s, P(), TOTAL, "3부")
note(s, "[2분] ★ 구조 내장 단조.\n"
     "대본: Fig.3 한 층. t는 입력 feature마다 +1, −1, 0 중 하나입니다. "
     "+1이면 그 입력이 올라갈 때 출력도 올라가야 하고, −1이면 반대, 0이면 그 feature는 단조 제약 없습니다. "
     "배터리 예는 SOC에 +1 — 충전량 올라가면 전압도 올라가야 합니다. "
     "|W|_t는 t=±1일 때 가중치 부호만 고정하는 것이지, 가중치 전체 절댓값이 아닙니다. "
     "s는 hidden m개를 볼록 ρ̊, 오목 ρ̂, 포화 ρ̃에 나누는 인원표입니다. "
     "예 (4,4,0)이면 볼록 4개, 오목 4개 — x³나 OCV 곡선처럼 양쪽으로 휘는 단조 함수에 씁니다. "
     "ReLU만 쓰면 clip+ReLU와 같아서 x³ 왼쪽이 깨집니다. "
     "몇 번째 뉴런이냐로 φ가 정해지지, x 부호로 바뀌지 않습니다. "
     "앞 장 Fig.1 (c)가 이 구조의 결과입니다 — (b) clip+ReLU는 x³ 실패, CMNN은 t+s로 둘 다 잡습니다. "
     "외삽 보장은 t 쪽 **방향**뿐 — 값·곡선은 틀릴 수 있습니다. "
     "다음 장 ③ — CMNN보다 physics **한 단계 더**: dQ/dV feature, thermo/단조 **loss**. "
     "같은 ‘밖에서 이상한 예측’ 문제지만 **넣는 지식·방식**이 다릅니다.")
takeaway(s, "② CMNN — 방향만 **구조**로 고정",
         "다음 ③ — feature·physics loss (mechanism)")

# ── S22 방법 3: Physics-ML / Physics loss (외삽) ──
s = blank()
content_header(s, "3부 · 방법 3", "③ **외삽** · Physics-ML — **방향(② CMNN)보다 더 아는** 경우")
if (ASSETS / "fig_physics_loss_example.png").exists():
    s.shapes.add_picture(
        str(ASSETS / "fig_physics_loss_example.png"),
        Inches(0.18), Inches(1.04), width=Inches(7.28),
    )
card_text(
    s, Inches(7.52), Inches(1.04), Inches(5.42), Inches(1.82),
    "외삽 문제 · 방법 선택 (그림 a)",
    [
        "고 DoD 학습 → 저 DoD 시험 — **train hull 밖**",
        "**방향만** (cycle↑→용량↓) → **② CMNN** · **법칙·열화 지표** → **③ Physics-ML**",
        "순수 ML 외삽은 비선형 실제와 **괴리** — 물리는 OOD에서도 성립",
    ],
    NAVY, dense=True,
)
card_text(
    s, Inches(7.52), Inches(2.94), Inches(5.42), Inches(1.72),
    "③-A Feature — 입력에 물리 (그림 b · Li)",
    [
        "충방전 곡선에서 **열화 peak(dQ/dV)** 뽑을 수 있을 때",
        "Q(V) → **dQ/dV peak** → 회귀 — **loss 아님**",
        "DoD 바뀌어도 peak 유효 · OOD MAPE **~22%**",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(7.52), Inches(4.74), Inches(5.42), Inches(1.72),
    "③-B Loss — 학습에 법칙 (그림 c · Aykol)",
    [
        "**단조만**이면 CMNN으로 충분 — **loss 불필요**",
        "**Arrhenius·열역학** 등 **관계 형태**를 알 때",
        "L = L_data + **λ·L_physics** · soft · **holdout** 필수",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "li2023", Inches(9.55), Inches(0.28), w=Inches(1.55), label="PDF · Li")
pdf_btn(s, "aykol", Inches(11.15), Inches(0.28), w=Inches(1.55), label="PDF · Aykol")
footer(s, P(), TOTAL, "3부")
note(s, "[2.5분]\n"
     "Q: 방향만 알면 CMNN 아닌가? — 맞다. ②가 그 경우.\n"
     "③ Physics-ML은 방향보다 더 안다 — 열화 peak(Li) 또는 법칙 형태(Aykol).\n"
     "Li: 고 DoD 학습 → 저 DoD 시험. ML 외삽 실패 → dQ/dV peak를 입력에.\n"
     "Aykol loss: 단조는 CMNN이 더 적합. Arrhenius·열역학 등 형태를 알 때 loss.\n"
     "Fesser: soft constraint — OOD holdout 재검증 필수.")
takeaway(s, "**방향만 → CMNN(②)** · **더 알면 Physics-ML(③)**",
         "Feature(Li) · Loss(Aykol) · **holdout**")

# ── S21 방법 4: UQ + 기권 ──
s = blank()
content_header(s, "3부 · 방법 4", "사전 지식 없음 — UQ(경보) + 기권(예측 안 함)")
fig_panel(s, Inches(0.30), Inches(1.14), Inches(6.20), Inches(3.05))
if (ASSETS / "fig_uq_aleatoric_epistemic.png").exists():
    s.shapes.add_picture(
        str(ASSETS / "fig_uq_aleatoric_epistemic.png"),
        Inches(0.38), Inches(1.22), width=Inches(6.05),
    )
pdf_btn(s, "ghahramani", Inches(9.55), Inches(0.28), w=Inches(1.55), label="PDF · Ghahramani")
pdf_btn(s, "zhu", Inches(11.15), Inches(0.28), w=Inches(1.55), label="PDF · Zhu")
card_text(
    s, Inches(0.32), Inches(4.28), Inches(6.18), Inches(1.35),
    "UQ 그림",
    [
        "훈련 범위 밖 → 불확실성 띠 벌어짐 = **경보**",
        "예측값은 **나옴** — ‘믿지 말 것’ 신호",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.32), Inches(5.72), Inches(6.18), Inches(1.28),
    "실무 예 — 센서 3개 → 품질",
    [
        "학습: 센서 각 0~10 · 추론: 센서=**20** → **범위 밖(외삽)**",
        "UQ: ‘품질 5.2 (불확실↑)’ · 기권: ‘**예측 안 함** → 재측정’",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(6.65), Inches(1.14), Inches(6.35), Inches(1.45),
    "UQ — 경보등",
    [
        "예측값 + 불확실성(띠·분산) **둘 다** 출력",
        "밖/낯선 입력 → 불확실성↑ — **숫자는 나옴**",
        "예: 앙상블 5개 갈리면 경보 · GP는 멀수록 띠↑",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(6.65), Inches(2.68), Inches(6.35), Inches(1.55),
    "기권 — 답변 거부",
    [
        "신뢰·불확실성이 임계 미만 → **예측값 자체 안 냄**",
        "규칙: if 센서∉[0,10]: 거부 도 동일 철학",
        "성능: **정확도 × coverage** — Zhu (operator 문헌)",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(6.65), Inches(4.32), Inches(6.35), Inches(1.35),
    "UQ vs 기권",
    [
        "**UQ** = 경보 (숫자 O) · **기권** = 차단 (숫자 X)",
        "불확실 작음 → 사용 · 큼 → 경보 또는 **거부**",
        "고위험: **틀린 답 1번 > 기권 1번**(실험·재측정)",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(6.65), Inches(5.72), Inches(6.35), Inches(1.28),
    "DeepONet(Zhu) — 부록",
    [
        "PDE **연산자** extrap + 기권 — **이 장 그림 아님**",
        "센서→품질 = **표 회귀** — 위 UQ·기권으로 충분",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "3부")
note(s, "[2.5분]\n"
     "대본: 방법 ④ — 사전 지식 없을 때. UQ는 ‘예측값 + 얼마나 확신?’을 같이 말합니다. "
     "그림 — 훈련 범위 밖으로 가면 불확실성 띠가 벌어집니다. 숫자는 나오지만 경보입니다. "
     "기권은 한 단계 더 — 신뢰가 낮거나 입력이 범위 밖이면 예측값 자체를 안 냅니다. "
     "실무 예 — 센서 0~10으로 학습했는데 20이 들어오면 외삽입니다. UQ는 ‘품질 5.2 불확실’, "
     "기권은 ‘예측 안 함, 재측정’입니다. 고위험이면 틀린 숫자보다 안 하는 게 낫습니다. "
     "Zhu·DeepONet은 operator 문헌 예 — 부록. 방법 네 가지 끝 — 종합합니다.")
takeaway(s, "UQ=경보(숫자 O) · 기권=차단(숫자 X)",
         "센서 범위 밖(20) → 재측정 · DeepONet=부록")

# ── S22 종합 1: 같은 데이터, 다른 밖 ──
s = blank()
content_header(s, "3부 · 종합 (1)", "같은 훈련 데이터 — 가정이 다르면 밖 거동이 전부 다르다")
if (ASSETS / "fig_method_cases.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_method_cases.png"), Inches(0.90), Inches(1.20), width=Inches(9.2))
card_text(
    s, Inches(10.35), Inches(1.20), Inches(2.6), Inches(4.9),
    "패널 = 넣은 가정",
    [
        "훈련점은 6패널 동일",
        "MLP: '밖=직선' (암묵)",
        "EQL: '닫힌 수식'",
        "Mono: '단조 방향'",
        "Physics-ML: soft loss·feature",
        "UQ: '점추정 포기'",
        "→ 데이터는 동일,",
        "  밖 거동만 다름",
        "차이의 원인 = 가정",
        "(S03 정의의 시각화)",
    ],
    TEAL, dense=True,
)
footer(s, P(), TOTAL, "3부")
note(s, "[1.5분]\n"
     "대본: 종합 첫 장. 여섯 패널 모두 '완전히 같은' 훈련 데이터입니다. 넣은 가정만 다릅니다. "
     "MLP는 직선으로, EQL은 식으로, 단조망은 방향을 지키며, Physics-ML은 soft loss로, UQ는 폭을 벌리며 나갑니다. "
     "2부에서 말한 '밖에서의 행동 = 모델에 들어 있는 가정의 발현'을 눈으로 확인하는 장입니다. "
     "데이터는 아무것도 결정하지 않았습니다 — 전부 가정이 결정했습니다.")
takeaway(s, "같은 데이터, 다른 가정 → 완전히 다른 밖 거동",
         "차이는 데이터가 아니라 전부 가정에서 온다")

# ── S25 종합 2: 선택 가이드 ──
s = blank()
content_header(s, "3부 · 종합 (2)", "선택 가이드 — 무엇을 아는가 → 그만큼의 가정 → 나머지는 UQ")
if (ASSETS / "fig_method_decision.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_method_decision.png"), Inches(0.30), Inches(1.15), height=Inches(2.95))
headers = ["가용 지식", "방법", "이득", "대가"]
rows = [
    ["함수 모양·산술", "EQL / NALU", "전 영역 닫힌 식", "가정이 틀릴 때 전역 실패·불안정"],
    ["방향(단조)", "CMNN", "방향 역전 차단", "구조 설계 비용"],
    ["물리 (회귀)", "Physics-ML / soft loss", "Aykol·Li·Muckley", "λ val 튜닝 · ≠ PINN"],
    ["PDE·장 u(x,t)", "PINN (Raissi)", "시뮬·희소 data", "잔차≠보증 · tabular✗"],
    ["없음", "UQ + 기권", "예측 안 함", "답변율 감소"],
]
col_w = [2.4, 2.3, 2.6, 3.2]
x0 = 0.35
y0 = 4.30
for j, htxt in enumerate(headers):
    x = Inches(x0 + sum(col_w[:j]))
    rect(s, x, Inches(y0), Inches(col_w[j] - 0.06), Inches(0.36), NAVY, LINE)
    tb = s.shapes.add_textbox(x + Inches(0.06), Inches(y0 + 0.045), Inches(col_w[j] - 0.16), Inches(0.28))
    add_text(tb, htxt, size=11, bold=True, color=TEAL)
for i, row in enumerate(rows):
    y = Inches(y0 + 0.40 + i * 0.38)
    for j, cell in enumerate(row):
        x = Inches(x0 + sum(col_w[:j]))
        rect(s, x, y, Inches(col_w[j] - 0.06), Inches(0.34), CARD if i % 2 == 0 else PANEL, LINE)
        tb = s.shapes.add_textbox(x + Inches(0.06), y + Inches(0.04), Inches(col_w[j] - 0.16), Inches(0.26))
        add_text(tb, cell, size=10.5, bold=(j == 0), color=INK if j == 0 else SLATE)
card_text(
    s, Inches(10.95), Inches(4.30), Inches(2.0), Inches(1.90),
    "규칙",
    ["지식 우선", "방법 후행", "불확실하면 예측 안 함"],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: 선택을 표 하나로 정리합니다. 왼쪽 열이 여러분이 실제로 아는 것이고, 나머지가 따라옵니다. "
     "함수 모양을 알면 EQL·NALU — 전 영역 식을 사지만 틀리면 전역 실패. 방향만 알면 CMNN — 방향 보장을 사고 설계 비용을 냅니다. "
     "물리·soft loss(Aykol·Li·Muckley) — **센서→품질 해당**. PDE·장 u(x,t)만 Raissi PINN — tabular ✗. 검증 없이는 보증 아님(Fesser). 아무것도 모르면 UQ+기권 — 정직함을 사고 답변율을 냅니다. "
     "그리고 마지막 규칙이 어디에나 적용됩니다: 어떤 방법을 쓰든, 모르는 부분에는 UQ를 함께 태우세요. "
     "그런데 — physics loss 행이 남긴 숙제가 있었죠. 넣은 가정이 작동하는지 어떻게 확인하는가. 이제 검증입니다.")
takeaway(s, "지식 우선, 방법 후행, 불확실하면 예측 안 함(UQ)",
         "남은 숙제: 가정이 작동하는지 → 검증 2장")

# ── S26 검증 동기: silent failure ──
s = blank()
content_header(s, "3부 · 성능 검증 ①", "왜 **회귀 extrap** 검증이 별도인가 — 밖에서도 확신")
if (ASSETS / "fig_fail_cases.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_fail_cases.png"), Inches(0.45), Inches(1.18), width=Inches(8.0))
card_text(
    s, Inches(8.65), Inches(1.18), Inches(4.3), Inches(3.0),
    "왜 가정이 있어도 검증이 필요한가",
    [
        "가정 주입(3부 방법) ≠ 자동 보증",
        "Physics-ML: soft loss·feature",
        "  시간 외삽 실패 (Fesser 2023)",
        "평가 설계가 틀리면 '외삽 성능'",
        "  숫자가 보간 점수가 됨",
        "→ 가정 + 프로토콜이 한 세트",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.45), Inches(4.55), Inches(12.4), Inches(1.55),
    "세 가지 함정 (장표만으로도 판정 가능해야 함)",
    [
        "함정 1 — 테스트가 훈련 범위 안만인가: '외삽 성능' 숫자가 실은 보간 점수",
        "함정 2 — 높은 확신 + 낮은 정확도: 밖에서도 **분산·확신** 안 낮춤 (그림 c)",
        "함정 3 — 내부 지표 착시: 훈련 잔차·훈련 RMSE가 좋아도 밖 성능과 무관할 수 있음",
    ],
    CORAL, dense=True,
)
footer(s, P(), TOTAL, "3부")
note(s, "[1.5분]\n"
     "대본: 가정 주입의 마지막 조각 — **검증**입니다. 왜 별도로 필요한가, 함정이 셋입니다. "
     "첫째, 테스트가 훈련 범위 안이면 '외삽 성능'이라 부르는 숫자가 사실 보간 점수입니다. "
     "둘째, 그림 (c) — 모델은 범위 밖에서도 확신을 낮추지 않습니다. 자신 있게 틀립니다. "
     "셋째, 방금 본 physics loss 교훈 — 내부 지표가 좋아도 밖 성능과는 무관할 수 있습니다. "
     "그래서 확인 절차가 필요합니다. **회귀 extrap** 체크리스트로 정리합니다.")
takeaway(s, "테스트가 범위 안이면 외삽 점수는 오인 · 모델은 밖에서도 높은 확신",
         "→ hull 밖 holdout · baseline · UQ 확인")

# ── S27 검증: 회귀 extrap 체크리스트 ──
s = blank()
content_header(s, "3부 · 성능 검증 ②", "**회귀** 외삽 성능 주장 전 — 체크 4개")
if (ASSETS / "fig_checklist.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_checklist.png"), Inches(2.0), Inches(1.10), height=Inches(2.15))
card_text(
    s, Inches(0.35), Inches(3.35), Inches(6.2), Inches(1.75),
    "① 진짜 밖인가? (Bartley)",
    [
        "test ∩ Conv(X_train) = ∅ 인가",
        "  (= 테스트가 훈련 범위와 겹치지 않는가)",
        "min–max box만으로 '밖' 판정 **금지**",
        "아니면 '보간 성능'으로 보고",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(6.75), Inches(3.35), Inches(6.2), Inches(1.75),
    "② 회귀 holdout 설계",
    [
        "입력 **범위 밖** (예: 센서 0~10→20)",
        "목표 y도 범위 밖이면 **별도** holdout",
        "시간 extrap이면 **train window 뒤** 테스트",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(0.35), Inches(5.25), Inches(6.2), Inches(1.0),
    "③ baseline 공정 비교",
    [
        "선형·RF vs fancy NN — **같은 튜닝 예산·용량**",
        "튜닝 안 맞추면 '향상' = **오인**",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(6.75), Inches(5.25), Inches(6.2), Inches(1.0),
    "④ UQ/기권 + 민감도",
    [
        "밖에서 불확실↑·거부 되는지",
        "분할·시드·가정 약화 — 'SOTA' 자제",
    ],
    CORAL, dense=True,
)
pdf_btn(s, "bartley", Inches(9.55), Inches(0.28), w=Inches(1.55), label="PDF · Bartley")
pdf_btn(s, "pfister", Inches(11.15), Inches(0.28), w=Inches(1.55), label="PDF · Pfister")
footer(s, P(), TOTAL, "3부")
note(s, "[2분]\n"
     "대본: **회귀 extrap** 검증 체크 4개. "
     "① Bartley hull. ② holdout — 센서·y·시간 **범위 밖**. "
     "③ baseline **공정 비교** — IRM·DomainBed는 **분류 DG**, 오늘 범위 밖. "
     "④ UQ/기권 — 밖에서 경보·거부, 민감도.")
takeaway(s, "hull 밖 holdout · 공정 baseline · UQ 동작",
         "IRM/DomainBed=분류 DG · 회귀엔 hull+holdout")

# ════════ 엔딩 ════════

# ── S28 전체 요약 ──
s = blank()
content_header(s, "요약", "**회귀** 외삽 — 정의 · 실패 · 대응 · 검증")
card_text(
    s, Inches(0.35), Inches(1.25), Inches(4.1), Inches(3.35),
    "1. 회귀 외삽 정의",
    [
        "외삽 = 새 입력이 훈련 hull 밖",
        "연속 y · tabular/시계열",
        "분류·DG(OOD) → **범위 밖**",
        "가정 = 밖 거동 제약",
    ],
    TEAL, dense=True,
)
card_text(
    s, Inches(4.60), Inches(1.25), Inches(4.1), Inches(3.35),
    "2. 실패 원인 (4)",
    [
        "① 함수 못 고름 (Pfister 2024)",
        "② 고차원·대부분 밖 (Bartley)",
        "③ 범위 밖에선 모름↑",
        "④ ReLU → 멀리서 직선",
        "     (Xu Thm.1, ICLR 2021)",
        "→ '외삽 되나?'가 아니라",
        "  '가정이 문제에 맞나?'",
    ],
    CORAL, dense=True,
)
card_text(
    s, Inches(8.85), Inches(1.25), Inches(4.1), Inches(3.35),
    "3. 대응 + 검증",
    [
        "EQL/NALU · CMNN · Physics-ML · UQ/기권",
        "검증: hull 밖 holdout",
        "     · baseline 공정 비교",
        "     · UQ/기권 동작 · 민감도",
        "지식 → 대응 → 검증",
    ],
    TEAL, dense=True,
)
box = rect(s, Inches(0.45), Inches(4.85), Inches(12.4), Inches(1.15), NAVY)
t = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12.0), Inches(0.9))
tf = add_text(t, "밖을 지탱하는 것은 데이터가 아니라 가정이다", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "가정 = 훈련에 없는 영역의 함수 거동을 미리 제약하는 지식  ·  올바르게 넣고  ·  프로토콜로 검증한다", size=13, color=SOFT, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL)
note(s, "[1.5분]\n"
     "대본: **회귀 extrap** 세 질문 답으로 닫습니다. "
     "정의 — hull 밖 + 연속 y. 실패 — 데이터·차원·UQ·ReLU. "
     "대응 — 가정 넣기. 검증 — hull holdout·baseline·UQ. "
     "한 문장: 밖을 지탱하는 것은 데이터가 아니라 가정.")
takeaway(s, "밖을 지탱하는 것은 데이터가 아니라 가정이다",
         "회귀 extrap: 정의 → 원인 4 → 대응 + hull 검증")

# ── S29 필독 3편 ──
s = blank()
content_header(s, "필독", "딱 3편 — **회귀 extrap** 진단 · 처방 · 검증")
trio = [
    ("진단", "Xu et al. 2021 (ICLR)", ["How NNs Extrapolate", "ReLU→직선 (회귀)", "2부 근거"], CORAL, "xu"),
    ("처방", "Runje et al. 2023 (ICML)", ["Constrained Monotonic NN", "방향 가정 → 구조", "3부 근거"], TEAL, "runje"),
    ("검증", "Bartley et al. 2019", ["Characterizing Extrapolation", "hull·고차원 — 범위 판정", "검증 근거"], TEAL, "bartley"),
]
for i, (role, cite, bl, c, key) in enumerate(trio):
    x = Inches(0.45 + i * 4.2)
    card_text(s, x, Inches(1.45), Inches(4.0), Inches(3.6), f"{role} — {cite}", bl, c)
    pdf_btn(s, key, x + Inches(0.15), Inches(5.25), w=Inches(1.6), label="PDF 열기")
card_text(
    s, Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.85),
    "더 읽기",
    ["Pfister 2024 · Fesser 2023 · Zhu 2022 (UQ/기권) — OoD-Bench·DomainBed는 **분류 DG** (INDEX.md)"],
    TEAL, dense=True,
)
footer(s, P(), TOTAL)
note(s, "[1분]\n"
     "대본: **회귀 extrap** 필독 3편 — Xu(진단), Runje(CMNN), Bartley(hull 검증). "
     "OoD-Bench·DomainBed·IRM은 **분류 DG** — INDEX에 있으나 오늘 범위 밖.")
takeaway(s, "Xu(진단) · Runje(처방) · Bartley(검증)",
         "분류 DG(OoD-Bench·DomainBed) → INDEX 참고")

# ── S30 Q&A / 참고문헌 ──
s = blank()
content_header(s, "Q&A · 참고문헌", "**회귀 extrap** 본문  ·  분류 DG는 INDEX")
refs_l = [
    ("xu", "Xu 2021 — NN Extrapolate (2부·3부)"),
    ("pfister", "Pfister 2024 — 함수 못 정함 (2부)"),
    ("bartley", "Bartley 2019 — 고차원·훈련 범위 (1부·2부)"),
    ("eql", "Martius 2016 — EQL (3부)"),
    ("nalu", "Trask 2018 — NALU (3부)"),
    ("runje", "Runje 2023 — CMNN (3부)"),
]
refs_r = [
    ("aykol", "Aykol 2021 — Physics-ML (3부)"),
    ("li2023", "Li 2023 — Physics RUL (3부)"),
    ("raissi", "Raissi 2019 — PINN PDE only (3부)"),
    ("fesser", "Fesser 2023 — PINN 반례 (3부)"),
    ("ghahramani", "Ghahramani 2013 — UQ (3부)"),
    ("zhu", "Zhu 2022 — 기권 (3부)"),
    ("note", "학습노트 PDF (전체 정리)"),
]
for i, (key, title) in enumerate(refs_l):
    y = Inches(1.45 + i * 0.62)
    pdf_btn(s, key, Inches(0.45), y, w=Inches(1.45), h=Inches(0.42), label="PDF")
    tb = s.shapes.add_textbox(Inches(2.05), y + Inches(0.06), Inches(4.4), Inches(0.35))
    add_text(tb, title, size=12, color=INK)
for i, (key, title) in enumerate(refs_r):
    y = Inches(1.45 + i * 0.62)
    pdf_btn(s, key, Inches(6.75), y, w=Inches(1.45), h=Inches(0.42), label="PDF")
    tb = s.shapes.add_textbox(Inches(8.35), y + Inches(0.06), Inches(4.5), Inches(0.35))
    add_text(tb, title, size=12, color=INK)
footer(s, P(), TOTAL)
note(s, "[Q&A]\n"
     "예상 질문 대비:\n"
     "Q. '데이터를 엄청 키우면 (foundation model) 해결 안 되나?' → 범위가 커질 뿐 범위 밖은 여전히 존재. "
     "함수를 못 정하는 건 데이터 양이 아니라 범위의 문제.\n"
     "Q. 'ReLU 말고 다른 활성화면?' → Xu Fig.5 — 밖의 모양이 바뀔 뿐 '가정이 생긴다'는 구조는 동일.\n"
     "Q. '단조 가정이 틀린 데이터면?' → 방법 1과 같은 리스크. 그래서 검증 + 민감도 확인이 필수.\n"
     "Q. 'IRM·DomainBed는?' → **분류 domain generalization**. 회귀 extrap 본문 **범위 밖**. baseline 공정 비교 **교훈만** 차용.\n"
     "Q. '분류 extrap은?' → reject/OOD — **별도** 주제 (INDEX 참고).\n"
     "Q. 'PINN 실패는 해결됐나?' → 완화책은 있으나 '검증 별도' 원칙은 유지.")
actual = len(prs.slides)
print(f"slides={actual} (footer TOTAL was {TOTAL})")
prs.save(OUT)
print(f"saved: {OUT}")
print(f"size_mb={OUT.stat().st_size/1e6:.2f}")
from pptx import Presentation as _P
_prs = _P(str(OUT))
n = 0
for sl in _prs.slides:
    for sh in sl.shapes:
        try:
            a = sh.click_action.hyperlink.address
            if a and a.startswith("paper_pdfs/"):
                n += 1
                assert (ROOT / a).exists(), a
        except Exception:
            pass
print(f"relative pdf links ok: {n}")
