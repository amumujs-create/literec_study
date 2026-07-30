#!/usr/bin/env python3
"""Export each PPTX slide to PNG preview for the speaker guide MD."""
from __future__ import annotations

import io
import re
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
_V = sys.argv[1] if len(sys.argv) > 1 else "v5"
if _V == "v6":
    PPTX = ROOT / "외삽_50분_발표자료_v6.pptx"
    OUT = ROOT / "_assets" / "slide_previews_v6"
else:
    PPTX = ROOT / "외삽_50분_발표자료_v5_심화.pptx"
    OUT = ROOT / "_assets" / "slide_previews"

# Render at ~half native (13.333"×7.5" @ 96dpi ≈ 1280×720)
SCALE = 1280 / 13.333333


def emu_to_px(v: int) -> int:
    return max(0, int(Emu(v).inches * SCALE))


def rgb_of(color) -> tuple[int, int, int] | None:
    if color is None:
        return None
    try:
        if color.type is None:
            return None
        if hasattr(color, "rgb") and color.rgb is not None:
            c = color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None


def fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        return rgb_of(fill.fore_color)
    except Exception:
        return None


def font_for(size_pt: float, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    size = max(10, int(size_pt * SCALE / 96 * 1.15))
    candidates = [
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size=size, index=0)
            except Exception:
                continue
    return ImageFont.load_default()


def draw_text_frame(draw: ImageDraw.ImageDraw, shape, canvas_w: int, canvas_h: int) -> None:
    if not shape.has_text_frame:
        return
    x = emu_to_px(shape.left)
    y = emu_to_px(shape.top)
    w = emu_to_px(shape.width)
    h = emu_to_px(shape.height)
    cur_y = y + 2
    for para in shape.text_frame.paragraphs:
        text = para.text or ""
        if not text.strip() and not para.runs:
            cur_y += 8
            continue
        # style from first run
        size_pt = 14
        bold = False
        color = (236, 241, 247)
        if para.runs:
            r0 = para.runs[0]
            if r0.font.size:
                size_pt = r0.font.size.pt
            bold = bool(r0.font.bold)
            c = rgb_of(r0.font.color)
            if c:
                color = c
            # rebuild full paragraph text from runs (markdown ** stripped)
            text = "".join(r.text or "" for r in para.runs)
        text = text.replace("**", "")
        font = font_for(size_pt, bold)
        # word wrap
        words = text.split(" ")
        lines: list[str] = []
        line = ""
        for word in words:
            trial = (line + " " + word).strip() if line else word
            bbox = draw.textbbox((0, 0), trial, font=font)
            if bbox[2] - bbox[0] > w - 8 and line:
                lines.append(line)
                line = word
            else:
                line = trial
        if line:
            lines.append(line)
        if not lines and text:
            lines = [text]
        for ln in lines:
            if cur_y > y + h:
                return
            draw.text((x + 4, cur_y), ln, fill=color, font=font)
            bbox = draw.textbbox((0, 0), ln, font=font)
            cur_y += max(12, bbox[3] - bbox[1] + 3)


def render_slide(slide, out_path: Path) -> None:
    W = emu_to_px(slide.part.package.presentation_part.presentation.slide_width)
    H = emu_to_px(slide.part.package.presentation_part.presentation.slide_height)
    # fallback size
    if W < 100:
        W, H = 1280, 720
    img = Image.new("RGB", (W, H), (11, 16, 22))
    draw = ImageDraw.Draw(img)

    # sort by z-order approx: background first (large full-bleed), then others
    shapes = list(slide.shapes)
    shapes.sort(key=lambda s: (0 if (s.width > slide.part.package.presentation_part.presentation.slide_width * 0.9 and s.height > slide.part.package.presentation_part.presentation.slide_height * 0.9) else 1,
                               s.top or 0, s.left or 0))

    sw = slide.part.package.presentation_part.presentation.slide_width
    sh = slide.part.package.presentation_part.presentation.slide_height

    for shape in shapes:
        left, top = emu_to_px(shape.left), emu_to_px(shape.top)
        width, height = max(1, emu_to_px(shape.width)), max(1, emu_to_px(shape.height))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                pic = Image.open(io.BytesIO(blob)).convert("RGBA")
                pic = pic.resize((width, height), Image.Resampling.LANCZOS)
                img.paste(pic, (left, top), pic if pic.mode == "RGBA" else None)
            except Exception:
                pass
            continue

        fill = fill_rgb(shape)
        if fill and shape.shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.RECTANGLE,
            MSO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        ):
            # skip fully transparent / line-only thin bars still useful
            r = 8 if height > 40 and width > 40 else 0
            if r:
                draw.rounded_rectangle([left, top, left + width, top + height], radius=r, fill=fill)
            else:
                draw.rectangle([left, top, left + width, top + height], fill=fill)

        if shape.has_text_frame and shape.text_frame.text.strip():
            draw_text_frame(draw, shape, W, H)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, "PNG", optimize=True)
    print("wrote", out_path.name, f"{W}x{H}")


def main() -> None:
    prs = Presentation(str(PPTX))
    OUT.mkdir(parents=True, exist_ok=True)
    # fix size using presentation
    global SCALE
    # keep SCALE from top; recompute W from inches
    for i, slide in enumerate(prs.slides, 1):
        # monkey: set slide dims via closure by attaching
        slide.part.package.presentation_part.presentation  # noqa: B018
        out = OUT / f"S{i:02d}.png"
        # patch dimensions onto slide for helper
        render_slide_fixed(prs, slide, out, i)
    print(f"done: {len(prs.slides)} previews → {OUT}")


def render_slide_fixed(prs: Presentation, slide, out_path: Path, idx: int) -> None:
    W = max(640, emu_to_px(prs.slide_width))
    H = max(360, emu_to_px(prs.slide_height))
    img = Image.new("RGB", (W, H), (11, 16, 22))
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        left, top = emu_to_px(shape.left or 0), emu_to_px(shape.top or 0)
        width = max(1, emu_to_px(shape.width or 1))
        height = max(1, emu_to_px(shape.height or 1))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                pic = pic.resize((width, height), Image.Resampling.LANCZOS)
                img.paste(pic, (left, top), pic)
            except Exception:
                pass
            continue

        fill = fill_rgb(shape)
        if fill is not None:
            # tiny 1pt lines
            if height <= 3:
                draw.rectangle([left, top, left + width, top + max(1, height)], fill=fill)
            elif width <= 8:
                draw.rectangle([left, top, left + width, top + height], fill=fill)
            else:
                rad = 10 if min(width, height) > 36 else 0
                box = [left, top, left + width, top + height]
                if rad:
                    draw.rounded_rectangle(box, radius=rad, fill=fill)
                else:
                    draw.rectangle(box, fill=fill)

        if shape.has_text_frame and (shape.text_frame.text or "").strip():
            # local draw using fixed W/H
            x, y, w, h = left, top, width, height
            cur_y = y + 2
            for para in shape.text_frame.paragraphs:
                text = "".join((r.text or "") for r in para.runs) if para.runs else (para.text or "")
                text = text.replace("**", "")
                if not text.strip():
                    cur_y += 6
                    continue
                size_pt = 14.0
                color = (236, 241, 247)
                if para.runs:
                    r0 = para.runs[0]
                    if r0.font.size:
                        size_pt = r0.font.size.pt
                    c = rgb_of(r0.font.color)
                    if c:
                        color = c
                font = font_for(size_pt)
                # wrap by character for Korean
                lines: list[str] = []
                line = ""
                for ch in text:
                    trial = line + ch
                    bb = draw.textbbox((0, 0), trial, font=font)
                    if bb[2] - bb[0] > w - 10 and line:
                        lines.append(line)
                        line = ch
                    else:
                        line = trial
                if line:
                    lines.append(line)
                for ln in lines:
                    if cur_y > y + h - 2:
                        break
                    draw.text((x + 4, cur_y), ln, fill=color, font=font)
                    bb = draw.textbbox((0, 0), ln, font=font)
                    cur_y += max(11, bb[3] - bb[1] + 2)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, "PNG", optimize=True)
    print(f"wrote {out_path.name} {W}x{H}")


if __name__ == "__main__":
    main()
