#!/usr/bin/env python3
"""v4·v5·v6·CA-CSS 부록을 한글 섹션으로 묶은 문헌조사 통합 PPT."""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "외삽_문헌조사_통합.pptx"

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(0x0B, 0x10, 0x16)
CARD = RGBColor(0x14, 0x1B, 0x24)
INK = RGBColor(0xEC, 0xF1, 0xF7)
SLATE = RGBColor(0xA0, 0xAD, 0xBC)
TEAL = RGBColor(0x3D, 0xD6, 0xC6)
CORAL = RGBColor(0xE8, 0x9A, 0x5C)
GREEN = RGBColor(0x6B, 0xC9, 0x8A)
FONT = "Apple SD Gothic Neo"

SOURCES = [
    ("PART 1 · 기초 문헌", "외삽_50분_발표자료_v4.pptx", "Hull · OOD · 신경망 외삽 · 필독 논문", TEAL),
    ("PART 2 · 심화", "외삽_50분_발표자료_v5_심화.pptx", "가정 스펙트럼 · UQ · 포기(abstention)", CORAL),
    ("PART 3 · 실전 정리", "외삽_50분_발표자료_v6.pptx", "검증 가능한 가정 · 평가 체크리스트", GREEN),
    ("PART 4 · 사례", "외삽_50분_사례_CA-CSS_v4_부록.pptx", "N-CMAPSS / CA-CSS v4 부록", TEAL),
]


def _run(p, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", FONT)


def _bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _box(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    _run(tf.paragraphs[0], text, size, bold=bold, color=color, align=align)


def _panel(slide, x, y, w, h, title, body, accent):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD
    sh.line.fill.background()
    _box(slide, x + 0.18, y + 0.12, w - 0.36, 0.38, title, 14, True, accent)
    _box(slide, x + 0.18, y + 0.52, w - 0.36, h - 0.7, body, 14, False, SLATE)


def _blank(prs):
    return prs.slide_layouts[6]


def _clear(slide) -> None:
    tree = slide.shapes._spTree
    for child in list(tree):
        if child.tag.endswith("}sp") or child.tag.endswith("}pic") or child.tag.endswith("}grpSp"):
            tree.remove(child)


def _copy_slide(dst: Presentation, src_slide) -> None:
    new_slide = dst.slides.add_slide(_blank(dst))
    _clear(new_slide)
    for shape in src_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            stream = io.BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(
                stream, shape.left, shape.top, width=shape.width, height=shape.height
            )
            continue
        new_slide.shapes._spTree.append(deepcopy(shape.element))


def _cover(prs: Presentation, counts: list[tuple[str, int]]) -> None:
    slide = prs.slides.add_slide(_blank(prs))
    _bg(slide)
    _box(slide, 0.8, 1.3, 11.7, 0.5, "문헌조사 통합본", 18, True, TEAL, PP_ALIGN.CENTER)
    _box(slide, 0.8, 1.9, 11.7, 0.9, "외삽(Extrapolation) 리터레처 스터디", 32, True, INK, PP_ALIGN.CENTER)
    _box(
        slide,
        1.4,
        2.85,
        10.5,
        0.5,
        "v4 기초 · v5 심화 · v6 실전 · CA-CSS 사례를 한 파일로",
        16,
        False,
        SLATE,
        PP_ALIGN.CENTER,
    )
    y = 3.55
    for name, n in counts:
        _panel(slide, 2.2, y, 8.9, 0.62, name, f"{n}장", TEAL)
        y += 0.72
    total = 1 + len(counts) + sum(n for _, n in counts)
    _box(slide, 0.8, 6.7, 11.7, 0.35, f"총 {total}장  ·  https://github.com/amumujs-create/literec_study", 13, False, SLATE, PP_ALIGN.CENTER)


def _section(prs: Presentation, title: str, subtitle: str, accent: RGBColor) -> None:
    slide = prs.slides.add_slide(_blank(prs))
    _bg(slide)
    _box(slide, 0.8, 2.5, 11.7, 0.45, title.split("·")[0].strip(), 16, True, accent, PP_ALIGN.CENTER)
    _box(slide, 0.8, 3.05, 11.7, 0.8, title, 32, True, INK, PP_ALIGN.CENTER)
    _box(slide, 1.5, 4.05, 10.3, 0.6, subtitle, 18, False, SLATE, PP_ALIGN.CENTER)


def main() -> None:
    counts: list[tuple[str, int]] = []
    decks: list[tuple[str, Path, str, RGBColor, Presentation]] = []
    for title, fname, subtitle, accent in SOURCES:
        path = ROOT / fname
        if not path.exists():
            raise FileNotFoundError(path)
        src = Presentation(str(path))
        counts.append((title, len(src.slides)))
        decks.append((title, path, subtitle, accent, src))

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    _cover(prs, counts)
    for title, _path, subtitle, accent, src in decks:
        _section(prs, title, f"{subtitle}  ·  {len(src.slides)}장", accent)
        for slide in src.slides:
            _copy_slide(prs, slide)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
