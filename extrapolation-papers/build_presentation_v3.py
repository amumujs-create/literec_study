#!/usr/bin/env python3
"""Build 외삽 50분 발표자료 v3 (widescreen, speaker notes, figures)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "_assets"
OUT = ROOT / "외삽_50분_발표자료_v3.pptx"

# 16:9
W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x1A, 0x7A, 0x6D)
CORAL = RGBColor(0xC4, 0x5C, 0x26)
SLATE = RGBColor(0x3D, 0x4A, 0x56)
MUTED = RGBColor(0x6B, 0x76, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE6, 0xEB)


def _set_run(run, size=18, bold=False, color=NAVY, font="Apple SD Gothic Neo"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(shape, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return p


def rect(slide, x, y, w, h, fill=CARD, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    # soft corners
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def bar(slide, color=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def footer(slide, page, total, part=""):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(10), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"외삽 완전 정복  ·  Kookmin Univ. IE Lab  ·  2026  {('·  ' + part) if part else ''}"
    _set_run(run, size=10, color=MUTED)
    num = slide.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.35))
    tf2 = num.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    _set_run(r2, size=10, color=MUTED)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def section_slide(title, subtitle, part_no, mins, page, total):
    s = blank()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(0.5))
    add_text(t, f"PART {part_no}  ·  {mins}", size=16, bold=True, color=TEAL)
    t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.85), Inches(11), Inches(1.2))
    add_text(t2, title, size=40, bold=True, color=WHITE)
    t3 = s.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11), Inches(0.6))
    add_text(t3, subtitle, size=18, color=RGBColor(0xB8, 0xC4, 0xD0))
    footer_box = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.4), Inches(0.35))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page} / {total}"
    _set_run(r, size=10, color=RGBColor(0x8A, 0x9A, 0xAA))
    return s


def content_header(slide, kicker, title):
    bar(slide, TEAL)
    k = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12), Inches(0.35))
    add_text(k, kicker, size=12, bold=True, color=TEAL)
    t = slide.shapes.add_textbox(Inches(0.45), Inches(0.55), Inches(12.3), Inches(0.55))
    add_text(t, title, size=26, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.15), Inches(12.4), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def card_text(slide, x, y, w, h, title, bullets, accent=TEAL):
    sh = rect(slide, x, y, w, h, CARD, LINE)
    # left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.18), w - Inches(0.4), h - Inches(0.3))
    tf = add_text(box, title, size=15, bold=True, color=NAVY)
    for b in bullets:
        add_para(tf, "•  " + b, size=13, color=SLATE, space_before=6)
    return sh


# ── build ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Estimate total slides for footers (update if you add/remove)
TOTAL = 40
page = 0


def P():
    global page
    page += 1
    return page


# 1 Title
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), W, Inches(1.6))
acc.fill.solid()
acc.fill.fore_color.rgb = RGBColor(0x0E, 0x2A, 0x4A)
acc.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.82), W, Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = TEAL
stripe.line.fill.background()
t = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(0.4))
add_text(t, "논문 기반 심화 세미나  ·  50분", size=16, bold=True, color=TEAL)
t2 = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.5), Inches(1.4))
add_text(t2, "외삽 (Extrapolation)\n완전 정복", size=44, bold=True, color=WHITE)
t3 = s.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.5), Inches(0.8))
add_text(
    t3,
    "기초 이론  ·  OOD 일반화  ·  신경망 외삽 방법론  ·  N-CMAPSS APEX-Guard",
    size=18,
    color=RGBColor(0xB8, 0xC4, 0xD0),
)
t4 = s.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11), Inches(0.8))
tf = add_text(t4, "Kookmin Univ. IE Lab", size=16, bold=True, color=WHITE)
add_para(tf, "literec_study · 36편 논문  ·  2026.07", size=13, color=RGBColor(0x9A, 0xAA, 0xBA))
note(s, "오프닝 30초. 오늘 목표: 외삽이 왜 어려운지 → 어떻게 방어하는지 → 우리 N-CMAPSS 결과로 연결.")
P()

# 2 Why
s = blank()
content_header(s, "들어가며", "왜 외삽을 공부해야 하는가")
cards = [
    ("01  데이터가 항상 불완전", ["실험 범위 밖 조건 예측이 일상", "N-CMAPSS: 미경험 TRA에서 RUL", "미래 cycle · 극한 운용점"]),
    ("02  ML의 근본 한계", ["대부분 모델은 in-distribution 전제", "범위 밖 오차는 급격히 증가", "테스트셋 ≠ 외삽 평가"]),
    ("03  산업 리스크", ["잘못된 외삽 → 안전·비용 사고", "예측값보다 UQ가 중요", "물리 제약·불변성이 방어막"]),
]
for i, (title, bullets) in enumerate(cards):
    card_text(s, Inches(0.45 + i * 4.2), Inches(1.5), Inches(4.0), Inches(4.4), title, bullets, TEAL if i < 2 else CORAL)
footer(s, P(), TOTAL, "Intro")
note(s, "2분. 핵심 메시지: '잘 맞는 테스트'가 보간일 수 있다. 외삽은 별도 프로토콜이 필요.")

# 3 Agenda
s = blank()
content_header(s, "AGENDA", "50분 구성 (Q&A 별도)")
if (ASSETS / "fig_timeline.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_timeline.png"), Inches(0.6), Inches(1.4), width=Inches(12.1))
rows = [
    ("Part 1", "외삽 기초 이론", "보간 vs 외삽 · Convex Hull · Richardson · UQ", "11분"),
    ("Part 2", "OOD 일반화", "분포 이동 · IRM · DomainBed · 인과 · 시계열", "11분"),
    ("Part 3", "신경망 외삽", "Xu · EQL · NALU · Monotonic · PINN · DeepONet", "14분"),
    ("Part 4", "N-CMAPSS 응용", "Strict TRA · APEX-Guard · CA-CSS · 결과", "10분"),
    ("Part 5", "동향 & 가이드", "타임라인 · 필독 논문 · 연구 적용", "4분"),
]
y0 = Inches(3.35)
for i, (p, t, d, m) in enumerate(rows):
    y = y0 + Inches(i * 0.58)
    box = rect(s, Inches(0.5), y, Inches(12.3), Inches(0.52), WHITE, LINE)
    a = s.shapes.add_textbox(Inches(0.7), y + Inches(0.08), Inches(1.3), Inches(0.35))
    add_text(a, p, size=14, bold=True, color=TEAL)
    b = s.shapes.add_textbox(Inches(2.1), y + Inches(0.08), Inches(2.6), Inches(0.35))
    add_text(b, t, size=14, bold=True, color=NAVY)
    c = s.shapes.add_textbox(Inches(4.8), y + Inches(0.08), Inches(6.2), Inches(0.35))
    add_text(c, d, size=13, color=SLATE)
    dbox = s.shapes.add_textbox(Inches(11.3), y + Inches(0.08), Inches(1.3), Inches(0.35))
    add_text(dbox, m, size=14, bold=True, color=CORAL, align=PP_ALIGN.RIGHT)
footer(s, P(), TOTAL)
note(s, "30초. 시간 배분 고지. Part3가 가장 김. Part4에서 우리 실험 숫자.")

# ════════════════ PART 1 ════════════════
section_slide("외삽이란 무엇인가", "기초 이론 · Convex Hull · Richardson · 불확실성", "1", "11분", P(), TOTAL)

# P1 interp vs extrap
s = blank()
content_header(s, "PART 1 · 기초 이론", "보간(Interpolation) vs 외삽(Extrapolation)")
card_text(
    s,
    Inches(0.45),
    Inches(1.45),
    Inches(6.0),
    Inches(4.6),
    "보간 — Interpolation",
    [
        "훈련 데이터 범위(지지 집합) 내부",
        "신뢰도 일반적으로 높음",
        "오차 제어·교차검증 가능",
        "대표: 표준 ML 테스트셋 평가",
        "예: 같은 운용 조건의 중간값",
    ],
    TEAL,
)
card_text(
    s,
    Inches(6.8),
    Inches(1.45),
    Inches(6.0),
    Inches(4.6),
    "외삽 — Extrapolation",
    [
        "훈련 지지 집합 밖",
        "신뢰도 급감 · 오차 거리 비례 폭발",
        "추가 가정(물리·단조성) 없으면 식별 불가",
        "대표: OOD, 미래 예측, RUL",
        "예: 고TRA holdout unit의 RUL",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. Bartley 2019 인용. '테스트 잘 나옴'만으로는 외삽 성공 주장 불가.")

# poly fig
s = blank()
content_header(s, "PART 1 · 직전 예시", "다항식 피팅의 외삽 실패")
if (ASSETS / "fig_poly_extrap.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_poly_extrap.png"), Inches(0.5), Inches(1.35), width=Inches(8.2))
card_text(
    s,
    Inches(8.9),
    Inches(1.45),
    Inches(3.9),
    Inches(4.6),
    "교훈",
    [
        "범위 [0,5] 안: 잘 맞음",
        "x>5: 오차 급발산",
        "고차 다항 = 외삽에 취약",
        "test∈hull → 보간 평가",
        "N-CMAPSS도 동일 구조",
        "Bartley 2019 · Tsai 2024",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 그림 가리키며 훈련 범위/외삽 영역 구분. 차수 올릴수록 밖이 더 위험.")

# convex hull definition
s = blank()
content_header(s, "PART 1 · 수학적 정의", "볼록 껍질(Convex Hull)과 외삽 영역")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(2.5),
    "정의",
    [
        "x ∈ Conv(X_train) → 보간 영역",
        "x ∉ Conv(X_train) → 외삽 영역",
        "Pfister 2024: support 밖은 추가 가정 없이 추론 불가",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(2.5),
    "차원의 저주",
    [
        "d↑ → hull 부피 비율 급감",
        "고차원: 테스트 대부분 외삽",
        "예: (T, C-rate, SOC, DOD) 4D",
    ],
    CORAL,
)
card_text(
    s,
    Inches(0.45),
    Inches(4.15),
    Inches(6.1),
    Inches(2.2),
    "불확실성 폭발",
    [
        "인근 관측 없을수록 bounds 확장",
        "기존 CI는 coverage 보장 상실",
        "→ extrapolation-aware CI 필요",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(4.15),
    Inches(6.0),
    Inches(2.2),
    "물리 지식 = 외삽 가정",
    [
        "단조성 · PDE · smoothness",
        "가정이 hull을 '논리적으로' 확장",
        "도메인 지식이 최고의 외삽 도구",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "2.5분. Pfister: 가정 없으면 밖은 식별 불가. 물리 제약이 왜 등장하는지 연결.")

# hull fig
s = blank()
content_header(s, "PART 1 · 시각화", "2D Convex Hull로 보간/외삽 구분")
if (ASSETS / "fig_convex_hull.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_convex_hull.png"), Inches(0.5), Inches(1.35), width=Inches(8.3))
card_text(
    s,
    Inches(9.0),
    Inches(1.45),
    Inches(3.8),
    Inches(4.6),
    "실무 체크",
    [
        "hull 안: 비교적 안전",
        "hull 밖: epistemic UQ↑",
        "d>10: 대부분 외삽",
        "대책: 물리 제약 주입",
        "또는 외삽축만 따로 설계",
        "Bartley · Pfister",
    ],
)
footer(s, P(), TOTAL, "Part 1")
note(s, "1.5분. X 표시가 외삽. N-CMAPSS의 TRA 축을 이 그림으로 비유.")

# error decomposition
s = blank()
content_header(s, "PART 1 · 오차 분해", "외삽 오차는 왜 폭발하는가")
headers = ["오차 요인", "보간 영역", "외삽 영역", "대책"]
rows = [
    ["편향 (Bias)", "복잡도 제어", "가정 위반 시 누적", "물리·smooth 가정"],
    ["분산 (Variance)", "인근 데이터 많음", "관측 없음 → 급증", "extrap-aware CI"],
    ["노이즈 σ²", "분리 측정 가능", "epistemic에 혼재", "Bayesian UQ"],
]
# table-like cards
for j, h in enumerate(headers):
    box = rect(s, Inches(0.45 + j * 3.15), Inches(1.4), Inches(3.05), Inches(0.55), NAVY)
    t = s.shapes.add_textbox(Inches(0.55 + j * 3.15), Inches(1.48), Inches(2.85), Inches(0.4))
    add_text(t, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        y = Inches(2.1 + i * 1.15)
        fill = WHITE if i % 2 == 0 else BG
        rect(s, Inches(0.45 + j * 3.15), y, Inches(3.05), Inches(1.05), fill, LINE)
        t = s.shapes.add_textbox(Inches(0.55 + j * 3.15), y + Inches(0.3), Inches(2.85), Inches(0.5))
        add_text(t, cell, size=13, color=NAVY if j == 0 else SLATE, align=PP_ALIGN.CENTER, bold=(j == 0))
footer(s, P(), TOTAL, "Part 1")
note(s, "2분. 외삽에서 variance(epistemic)가 지배. UQ 없이 점추정만 내는 건 위험.")

# Richardson
s = blank()
content_header(s, "PART 1 · 수치해석", "Richardson 외삽 → Probabilistic Richardson")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "고전 Richardson",
    [
        "A(h)=A* + c·h^p + O(h^(p+1))",
        "두 step h1,h2로 A* 추정",
        "수치 미분·적분·PDE 수렴값",
        "점추정만 — 불확실성 없음",
        "공학에서 오랜 표준",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "Probabilistic (Teckentrup 2024)",
    [
        "Richardson + GP prior",
        "사후분포 · 신뢰구간 제공",
        "수치해석 × Bayesian 융합",
        "공학 외삽의 안전장치",
        "JRSS-B · 필독 ★★☆",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "1.5분. 점추정→분포 추정으로 업그레이드. Part5 필독과 연결.")

# UQ
s = blank()
content_header(s, "PART 1 · UQ", "불확실성 정량화 — 외삽에서 왜 필수인가")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "Aleatoric",
    [
        "데이터 자체의 노이즈",
        "데이터↑로도 안 사라짐",
        "예: 센서 오차, 공정 변동",
    ],
    MUTED,
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "Epistemic",
    [
        "모델 지식 부족",
        "데이터↑로 줄일 수 있음",
        "외삽 영역에서 급증",
        "핵심: 얼마나 틀릴 수 있나",
    ],
    CORAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "외삽 UQ",
    [
        "epistemic 지배 + aleatoric 혼재",
        "기존 CI coverage 붕괴",
        "bounds + CI 결합",
        "Ghahramani 2013",
        "Pfister 2024",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 1")
note(s, "1.5분. Part1 마무리: 예측값 < 불확실성. Q로 넘기기 전 한 줄 요약.")

# ════════════════ PART 2 ════════════════
section_slide("OOD 일반화", "분포 이동 · IRM · DomainBed · 인과 · 시계열", "2", "11분", P(), TOTAL)

# TRA = OOD hook
s = blank()
content_header(s, "PART 2 · 실전 훅", "미경험 TRA 조건 = OOD")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "설정",
    [
        "Train: 저TRA (≤q70)",
        "Test: 고TRA holdout (>q90)",
        "Unit holdout: 10, 20",
        "Late band: RUL≤50",
    ],
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "실패 모드",
    [
        "ERM → 저TRA overfit",
        "고TRA에서 급락",
        "전이 window가 RMSE 오염",
        "→ strict 프로토콜 필요",
    ],
    CORAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "방어",
    [
        "CA-CSS 제약 학습",
        "Isotonic 단조 후처리",
        "ERM baseline 항상 비교",
        "DomainBed 교훈 적용",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "1.5분. Part4 예고. OOD를 추상 개념이 아닌 TRA 숫자로 체감.")

# shift types
s = blank()
content_header(s, "PART 2 · 분류", "분포 이동(Distribution Shift) 유형")
items = [
    ("Covariate Shift", "P(X)↑  P(Y|X) 동일", "다른 장비 의료영상\n고TRA 센서 분포"),
    ("Label Shift", "P(Y)↑  P(X|Y) 동일", "클래스 비율 변화"),
    ("Concept Drift", "P(Y|X) 자체 변화", "시간 경과 열화 법칙"),
    ("Domain Shift", "생성과정 전체 변화", "시뮬→실기 전이"),
]
for i, (title, defn, ex) in enumerate(items):
    x = Inches(0.45 + (i % 4) * 3.2)
    card_text(s, x, Inches(1.45), Inches(3.05), Inches(4.5), title, [defn, ex.replace("\n", " · ")], TEAL if i % 2 == 0 else NAVY)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. N-CMAPSS는 covariate(+concept) 혼합. Liu 2023 survey.")

# OOD formal
s = blank()
content_header(s, "PART 2 · 정의", "OOD 일반화의 공식 목표 (Ye et al., 2021)")
box = rect(s, Inches(0.45), Inches(1.5), Inches(12.4), Inches(1.8), NAVY)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.8), Inches(1.3))
tf = add_text(t, "min_h   max_{e ∈ E_all}   R^e(h)", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "훈련 환경 E_train에서 학습 → 미지 e_test에서도 낮은 리스크", size=14, color=RGBColor(0xB8, 0xC4, 0xD0), align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(3.6),
    Inches(6.0),
    Inches(2.7),
    "표준 ERM",
    ["min 평균_{e∈E_train} R^e(h)", "훈련 분포 평균에 overfit", "OOD에서 깨지기 쉬움"],
)
card_text(
    s,
    Inches(6.8),
    Inches(3.6),
    Inches(6.0),
    Inches(2.7),
    "OOD 목표",
    ["min 최악 환경 리스크", "robust / invariant 학습", "Ye 2021 · Arjovsky 2021"],
    TEAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. ERM vs worst-case. 수식은 직관만.")

# IRM
s = blank()
content_header(s, "PART 2 · 알고리즘", "Invariant Risk Minimization (IRM)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "핵심 아이디어 (Arjovsky)",
    [
        "모든 환경에서 동시에 최적인 w",
        "그런 w를 유도하는 표현 Φ 학습",
        "IRMv1: R + λ‖∇_{w=1} R‖²",
        "Colored MNIST: 색 함정 vs 숫자",
        "환경 정의가 애매하면 약해짐",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "비교",
    [
        "ERM — 평균 리스크, 단순",
        "IRM — 불변 특징",
        "GroupDRO — 최악 그룹",
        "Causal ML — 진짜 원인",
        "실무: ERM baseline 먼저",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. IRM 직관만. DomainBed로 바로 연결.")

# DomainBed
s = blank()
content_header(s, "PART 2 · 벤치마크", "DomainBed — 충격적인 재평가")
box = rect(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(1.6), CORAL)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.1))
tf = add_text(t, "공정한 HP 탐색 하에서 ERM ≥ 대부분의 OOD 알고리즘", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "Gulrajani & Lopez-Paz (2020) · ICML 2021", size=14, color=WHITE, align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "모델 선택이 핵심",
    ["성능 차의 상당 부분 = HP·선택", "알고리즘 자체보다 튜닝"],
)
card_text(
    s,
    Inches(4.65),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "공정 비교 체크",
    ["① ERM baseline", "② 동일 HP budget", "③ 동일 capacity", "④ target val 전략 명시"],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(3.3),
    Inches(4.0),
    Inches(3.0),
    "우리 실험에 적용",
    ["strict_late에서", "TabPFN·XGB·TF·TCN 동시 비교", "유닛 바꾸면 순위 뒤집힘"],
    CORAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "2분. 가장 중요한 슬라이드 중 하나. '새 방법' 주장 전 ERM 재검증.")

# Causal + time series
s = blank()
content_header(s, "PART 2 · 관점", "인과적 OOD · 시계열 OOD")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "통계적",
    ["상관 최대화", "데이터↑ · 정규화", "허위상관에 취약"],
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "불변성 / 인과",
    ["IRM · GroupDRO", "SCM · 진짜 원인", "개입에도 강건"],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "시계열 OOD (Wu 2025)",
    [
        "① 시간적 분포 이동",
        "② 주기적 분포 이동",
        "RUL = 둘 다 해당",
        "충방전/비행 주기 구조",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 2")
note(s, "1.5분. Part2 마무리. 시계열 OOD → Part4 RUL.")

# ════════════════ PART 3 ════════════════
section_slide("신경망 외삽 방법론", "왜 실패하는가 · EQL · NALU · Monotonic · PINN", "3", "14분", P(), TOTAL)

# Xu
s = blank()
content_header(s, "PART 3 · Xu et al. (2021)", "신경망은 왜 외삽에 실패하는가")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "관측",
    [
        "훈련 밖: 평면/직선",
        "ReLU MLP의 전형 패턴",
        "ICLR 2021 바이블",
    ],
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "Theorem 1",
    [
        "방향 v로 O(1/t) affine 수렴",
        "마지막 활성 선형구간 연장",
        "Appendix: R²>0.99",
    ],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "함의",
    [
        "비선형 타깃 외삽 실패",
        "활성화 = 외삽 가정",
        "구조 없이 '학습만'으론 한계",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "2.5분. 발표 핵심 논문 #1. 그림/정리만 강조.")

# activation fig
s = blank()
content_header(s, "PART 3 · 활성화 함수", "활성화 선택이 곧 외삽 가정")
if (ASSETS / "fig_activation.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_activation.png"), Inches(0.5), Inches(1.35), width=Inches(12.2))
box = rect(s, Inches(0.45), Inches(5.0), Inches(12.4), Inches(1.4), BG, LINE)
t = s.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(1.0))
tf = add_text(t, "ReLU: 방향별 조각 선형  ·  Tanh/Sigmoid: 상수 포화(회귀 치명)  ·  Sin: 주기 반복(물리엔 유용/과적합 위험)", size=15, color=NAVY)
add_para(tf, "Xu et al. §3.3 — 활성화 함수를 고르는 순간, 범위 밖 거동을 이미 선택한 것", size=13, color=SLATE)
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. 회색 띠=훈련 범위.")

# EQL
s = blank()
content_header(s, "PART 3 · 아키텍처", "EQL — Equation Learner (Martius & Lampert, 2016)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "아이디어",
    [
        "뉴런 = {sin, cos, ×, ÷, id, σ}",
        "학습 결과 = 해석 가능한 수식",
        "진짜 sin을 배우면 밖에서도 sin",
        "MLP는 조각선형 근사 → 발산",
        "소량 데이터 · 물리 관계 발견",
    ],
    TEAL,
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "적용",
    [
        "Symbolic regression",
        "전기화학 관계 모델링",
        "센서 곡선 외삽",
        "해석 가능성이 최우선일 때",
        "arxiv:1610.02995",
    ],
)
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. 수식을 배우면 외삽이 '공짜'로 따라옴.")

# NALU
s = blank()
content_header(s, "PART 3 · 아키텍처", "NALU — Neural Arithmetic Logic Units (Trask 2018)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "구조",
    [
        "NAC: W=tanh(Ŵ)⊙σ(M̂) → +/−",
        "NALU: g·NAC + (1−g)·exp(NAC(log|x|))",
        "× ÷ 까지 명시적 학습",
        "카운팅·누적합 외삽 특화",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "벤치마크 요지",
    [
        "덧셈/뺄셈: MLP 대비 대폭 향상",
        "훈련범위 5× 밖에서도 유지",
        "체계적 수치 외삽 가능",
        "arxiv:1808.00508",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. 숫자 연산이 필요한 과제에.")

# Monotonic
s = blank()
content_header(s, "PART 3 · 제약", "단조 제약 신경망 (Monotonic NN)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "왜?",
    [
        "공학 관계에 단조성 흔함",
        "제약을 강제 → 밖에서도 타당",
        "SOC↑ → OCV↑",
        "cycle↑ → RUL↓",
    ],
    CORAL,
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "방법",
    [
        "Weight clipping",
        "Certified MNN (Liu 2022)",
        "CMNN (Runje 2023)",
        "Min-Max 구조",
    ],
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "우리 연결",
    [
        "APEX-Guard isotonic",
        "cycle_mono hinge",
        "H1 위반 28%→0%",
        "물리 제약 = 외삽 가드",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. Part4 isotonic과 직접 연결. 강조.")

# PINN
s = blank()
content_header(s, "PART 3 · 물리", "PINN — 가능성과 실패 (Fesser 2023)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(6.1),
    Inches(4.7),
    "PINN (Raissi 2019)",
    [
        "L = L_data + λ L_physics",
        "L_physics = ‖N[u]‖² (PDE 잔차)",
        "데이터 희소 영역도 물리 일관",
        "전기화학·열전달에 매력적",
    ],
)
card_text(
    s,
    Inches(6.8),
    Inches(1.4),
    Inches(6.0),
    Inches(4.7),
    "왜 그래도 실패하나",
    [
        "Spectral bias (저주파 우선)",
        "경계 손실 가중 불균형",
        "외삽 방향 훈련점 희박",
        "DeepONet+UQ (Zhu 2022)",
        "신뢰도 낮으면 경고",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 3")
note(s, "2분. PINN 만능 아님. UQ 결합이 현실적.")

# method compare
s = blank()
content_header(s, "PART 3 · 종합", "외삽 방법론 한눈에 비교")
headers = ["방법", "원리", "데이터", "해석", "공학"]
rows = [
    ["MLP", "경험 근사", "대량", "낮음", "기준선"],
    ["EQL", "수식 발견", "소량", "매우 높음", "물리 관계"],
    ["NALU", "수치 연산", "중간", "중간", "카운팅"],
    ["Mono NN", "단조 제약", "중간", "높음", "SOC/RUL"],
    ["PINN", "PDE 내재", "소량", "높음", "전기화학"],
    ["IRM", "불변 특징", "다환경", "낮음", "도메인이동"],
    ["DeepONet+UQ", "연산자+UQ", "많음", "중간", "다물리"],
]
for j, h in enumerate(headers):
    rect(s, Inches(0.4 + j * 2.5), Inches(1.35), Inches(2.4), Inches(0.45), NAVY)
    t = s.shapes.add_textbox(Inches(0.45 + j * 2.5), Inches(1.4), Inches(2.3), Inches(0.35))
    add_text(t, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        y = Inches(1.85 + i * 0.65)
        fill = TEAL if (i == 3 and j == 0) else (WHITE if i % 2 == 0 else BG)
        fc = WHITE if (i == 3 and j == 0) else NAVY
        rect(s, Inches(0.4 + j * 2.5), y, Inches(2.4), Inches(0.6), fill, LINE)
        t = s.shapes.add_textbox(Inches(0.45 + j * 2.5), y + Inches(0.12), Inches(2.3), Inches(0.4))
        add_text(t, cell, size=12, bold=(j == 0), color=fc, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL, "Part 3")
note(s, "1.5분. Mono NN 하이라이트 → Part4로.")

# ════════════════ PART 4 ════════════════
section_slide("N-CMAPSS 응용", "Strict TRA · APEX-Guard · CA-CSS · 실험 결과", "4", "10분", P(), TOTAL)

# APEX intro
s = blank()
content_header(s, "PART 4 · 모델", "APEX-Guard — 외삽을 설계한 RUL 모델")
box = rect(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(1.2), NAVY)
t = s.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.8), Inches(0.8))
add_text(t, "Adaptive Physics EXtrapolation Guard  ·  외삽은 피할 수 없다 → 설계로 방어한다", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
card_text(
    s,
    Inches(0.45),
    Inches(2.9),
    Inches(4.0),
    Inches(3.4),
    "L1 평가",
    ["strict_late 프로토콜", "순수 외삽만 측정", "전이 window 제거"],
)
card_text(
    s,
    Inches(4.65),
    Inches(2.9),
    Inches(4.0),
    Inches(3.4),
    "L2 학습",
    ["CA-CSS 제약 발견·컴파일", "TRA(−) + cycle_mono", "hull을 논리적으로 확장"],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(2.9),
    Inches(4.0),
    Inches(3.4),
    "L3 추론",
    ["unit별 isotonic", "H1 위반 → 0%", "단조 RUL 궤적"],
    CORAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. 3층 방어 프레임 먼저.")

# strict protocol
s = blank()
content_header(s, "PART 4 · 프로토콜", "Strict TRA Extrap — 3겹으로 강제")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "① TRA shift",
    ["저TRA train → 고TRA test", "covariate shift 직접 검증", "tra_end>q90 · high_frac≥0.7"],
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "② Unit holdout",
    ["test units: 10, 20", "엔진 간 일반화 분리", "tra_rich 선정"],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "③ Late-life band",
    ["RUL≤50만 측정", "전이 window RMSE 오염 차단", "strict_late n=201"],
    CORAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2분. '어디까지를 test로 볼 것인가'가 외삽 연구의 출발점.")

# architecture
s = blank()
content_header(s, "PART 4 · 구조", "v4_disentangled — Dual Encoder")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(3.05),
    Inches(4.7),
    "입력",
    ["OC-Norm", "운용점 편차 분리", "cycle channel"],
)
card_text(
    s,
    Inches(3.7),
    Inches(1.4),
    Inches(3.05),
    Inches(4.7),
    "Health Enc",
    ["TRA=0 masking", "열화 패턴", "TRA overfit 방지"],
    TEAL,
)
card_text(
    s,
    Inches(6.95),
    Inches(1.4),
    Inches(3.05),
    Inches(4.7),
    "Load Enc",
    ["TRA, φ, T30, cycle", "외삽 축 전담"],
    CORAL,
)
card_text(
    s,
    Inches(10.2),
    Inches(1.4),
    Inches(2.7),
    Inches(4.7),
    "RUL 분해",
    ["RUL=health−damage", "w≥0 제약", "d=64, L=2, ep=15"],
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. Dual path가 핵심.")

# CA-CSS
s = blank()
content_header(s, "PART 4 · 학습", "CA-CSS + Isotonic")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "CA-CSS 파이프라인",
    [
        "Discovery → Pseudo-extrap",
        "→ Compile → Train",
        "train-band만 (leakage 방지)",
        "TRA(−) λ=0.05 adopt",
        "cycle_mono adopt",
    ],
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "손실",
    [
        "L = L_RUL + L_physics + L_TRA",
        "TRA(−) hinge",
        "cycle_mono hinge",
        "물리 감독 = 외삽 가드",
    ],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "Isotonic 후처리",
    [
        "unit별 cycle 정렬",
        "decreasing isotonic",
        "H1: 28% → 0%",
        "단조 궤적 보장",
    ],
    CORAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "1.5분. Part3 Mono NN과 동일 철학.")

# results
s = blank()
content_header(s, "PART 4 · 결과", "strict_late 검증 (tra_rich · seed=42)")
if (ASSETS / "fig_strict_rmse.png").exists():
    s.shapes.add_picture(str(ASSETS / "fig_strict_rmse.png"), Inches(0.4), Inches(1.3), width=Inches(7.8))
card_text(
    s,
    Inches(8.4),
    Inches(1.4),
    Inches(4.4),
    Inches(4.7),
    "숫자",
    [
        "APEX-Guard RMSE 3.26",
        "TabPFN 3.80 · R² 0.927",
        "Transformer 7.31",
        "TCN 14.09",
        "주의: nasa_drop14면 순위 뒤집힘",
        "→ multi-seed · unit 규칙 명시",
    ],
    TEAL,
)
footer(s, P(), TOTAL, "Part 4")
note(s, "2.5분. 메인 결과. 한계(유닛 민감)도 솔직히.")

# takeaway part4
s = blank()
content_header(s, "PART 4 · 교훈", "설계 → 성능으로 이어지는 한 줄")
box = rect(s, Inches(0.45), Inches(1.6), Inches(12.4), Inches(4.4), NAVY)
t = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.4), Inches(3.2))
tf = add_text(t, "외삽을 '우연히' 잘하길 바라지 말고\n평가(L1) · 제약(L2) · 단조(L3)를 설계하라", size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "", size=10)
add_para(tf, "이론(Xu · DomainBed · Monotonic) → 프로토콜(strict_late) → 모델(APEX-Guard)", size=16, color=RGBColor(0xB8, 0xC4, 0xD0), align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL, "Part 4")
note(s, "1분. Part4 클로징.")

# ════════════════ PART 5 ════════════════
section_slide("연구 동향 & 논문 가이드", "30년 타임라인 · 필독 · 연구 적용", "5", "4분", P(), TOTAL)

# timeline
s = blank()
content_header(s, "PART 5 · 타임라인", "외삽 연구 30년")
items = [
    ("1990s", "Vapnik · VC\n보간 이론"),
    ("2000s", "GP · Covariate\nShift"),
    ("2016-19", "EQL · NALU\nIRM"),
    ("2020-21", "DomainBed\nXu · PINN"),
    ("2022-23", "CMNN\nDeepONet+UQ"),
    ("2024-25", "Pfister\nProb.Richardson\nWu 시계열"),
]
for i, (yr, txt) in enumerate(items):
    x = Inches(0.4 + i * 2.15)
    rect(s, x, Inches(2.0), Inches(2.0), Inches(3.5), WHITE, LINE)
    top = slide_shapes_top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.0), Inches(0.55))
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL if i % 2 == 0 else NAVY
    top.line.fill.background()
    ty = s.shapes.add_textbox(x, Inches(2.08), Inches(2.0), Inches(0.4))
    add_text(ty, yr, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(x + Inches(0.1), Inches(2.8), Inches(1.8), Inches(2.4))
    add_text(tb, txt, size=13, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL, "Part 5")
note(s, "1분. 흐름만.")

# must read
s = blank()
content_header(s, "PART 5 · 필독", "우선순위 가이드 (레포 PDF 동봉)")
card_text(
    s,
    Inches(0.45),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★★★ 필수",
    [
        "Xu 2021 — NN 외삽",
        "Gulrajani 2020 — DomainBed",
        "Liu 2023 — OOD Survey",
    ],
    CORAL,
)
card_text(
    s,
    Inches(4.65),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★★☆ 방법",
    [
        "Trask 2018 NALU",
        "Runje 2023 CMNN",
        "Fesser 2023 PINN 실패",
        "Teckentrup 2024",
    ],
    TEAL,
)
card_text(
    s,
    Inches(8.85),
    Inches(1.4),
    Inches(4.0),
    Inches(4.7),
    "★☆☆ 응용",
    [
        "Aykol 2021 Battery+Phys",
        "Li 2023 Lifetime OOD",
        "Wu 2025 시계열 OOD",
        "총 36편 / extrapolation-papers/",
    ],
)
footer(s, P(), TOTAL, "Part 5")
note(s, "1.5분. 읽을 순서만 찍어주기.")

# apply
s = blank()
content_header(s, "PART 5 · 적용", "연구에 바로 붙이는 매핑")
rows = [
    ("문제", "추천 방법", "논문"),
    ("단조 관계 외삽 (SOC-OCV, RUL)", "Monotonic NN + isotonic", "Runje 2023"),
    ("조건 OOD (고TRA)", "제약 학습 + ERM 비교", "Liu · Gulrajani"),
    ("물리 PDE 영역", "PINN + UQ (DeepONet)", "Fesser · Zhu"),
    ("소량·해석 필요", "EQL + Bayesian UQ", "Martius · Teckentrup"),
]
for j, h in enumerate(rows[0]):
    rect(s, Inches(0.45 + j * 4.15), Inches(1.4), Inches(4.05), Inches(0.5), NAVY)
    t = s.shapes.add_textbox(Inches(0.55 + j * 4.15), Inches(1.48), Inches(3.85), Inches(0.35))
    add_text(t, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows[1:]):
    for j, cell in enumerate(row):
        y = Inches(2.0 + i * 1.05)
        rect(s, Inches(0.45 + j * 4.15), y, Inches(4.05), Inches(0.95), WHITE if i % 2 == 0 else BG, LINE)
        t = s.shapes.add_textbox(Inches(0.55 + j * 4.15), y + Inches(0.28), Inches(3.85), Inches(0.45))
        add_text(t, cell, size=13, color=NAVY, align=PP_ALIGN.CENTER, bold=(j == 0))
footer(s, P(), TOTAL, "Part 5")
note(s, "1분.")

# summary
s = blank()
content_header(s, "핵심 요약", "외삽을 정복하는 4가지 열쇠")
keys = [
    ("01  이론", "보간 ≠ 외삽\nsupport 밖은 가정 없으면 식별 불가"),
    ("02  OOD", "ERM 실패 가능\nDomainBed로 재검증 필수"),
    ("03  구조", "물리·단조 제약이\n최고의 외삽 도구"),
    ("04  UQ", "예측값보다\n얼마나 틀릴 수 있는가"),
]
for i, (t, b) in enumerate(keys):
    x = Inches(0.45 + i * 3.2)
    rect(s, x, Inches(1.5), Inches(3.05), Inches(4.4), WHITE, LINE)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(3.05), Inches(0.7))
    top.fill.solid()
    top.fill.fore_color.rgb = [NAVY, TEAL, CORAL, RGBColor(0x4A, 0x6F, 0xA5)][i]
    top.line.fill.background()
    tt = s.shapes.add_textbox(x, Inches(1.62), Inches(3.05), Inches(0.5))
    add_text(tt, t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bb = s.shapes.add_textbox(x + Inches(0.2), Inches(2.6), Inches(2.65), Inches(2.8))
    add_text(bb, b, size=15, color=SLATE, align=PP_ALIGN.CENTER)
footer(s, P(), TOTAL)
note(s, "1.5분. 클로징 메시지.")

# final advice
s = blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), W, Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = TEAL
stripe.line.fill.background()
t = s.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(0.7))
add_text(t, "최종 조언", size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
t2 = s.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(2.0))
add_text(
    t2,
    "외삽을 풀 때 가장 할 일:\n어떤 물리적·수학적 불변성이 존재하는가?",
    size=28,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
footer_box = s.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5))
add_text(footer_box, "단조성 · PDE · 보존법칙 — 있다면 그것이 최고의 외삽 도구다", size=16, color=RGBColor(0xB8, 0xC4, 0xD0), align=PP_ALIGN.CENTER)
note(s, "30초.")
P()

# refs
s = blank()
content_header(s, "참고문헌 · Q&A", "주요 논문 (전체 PDF는 literec_study/extrapolation-papers)")
refs = [
    "Xu et al. (2021) How Neural Networks Extrapolate — arxiv:2009.11848",
    "Gulrajani & Lopez-Paz (2020) In Search of Lost Domain Generalization — arxiv:2007.01434",
    "Liu et al. (2023) Towards OOD Generalization: A Survey — arxiv:2108.13624",
    "Arjovsky (2021) OOD Generalization in ML — arxiv:2103.02667",
    "Martius & Lampert (2016) Extrapolation and Learning Equations — arxiv:1610.02995",
    "Trask et al. (2018) Neural Arithmetic Logic Units — arxiv:1808.00508",
    "Runje & Shankaranarayana (2023) Constrained Monotonic NN — arxiv:2205.11775",
    "Fesser et al. (2023) Extrapolation Failures in PINNs — arxiv:2306.09478",
    "Teckentrup et al. (2024) Probabilistic Richardson Extrapolation — JRSS-B",
    "Wu et al. (2025) OOD Generalization in Time Series — arxiv:2503.13868",
    "Zhu et al. (2022) Reliable Extrapolation with DeepONet — arxiv:2212.06347",
    "Pfister & Bühlmann (2024) Extrapolation-Aware Nonparametric Inference — arxiv:2402.09758",
]
box = s.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.2), Inches(5.3))
tf = box.text_frame
tf.clear()
p = tf.paragraphs[0]
r = p.add_run()
r.text = refs[0]
_set_run(r, size=13, color=SLATE)
for ref in refs[1:]:
    add_para(tf, ref, size=13, color=SLATE, space_before=3, space_after=1)
footer(s, P(), TOTAL)
note(s, "Q&A. github.com/amumujs-create/literec_study")

# fix TOTAL in footers if page count differs
actual = len(prs.slides)
print(f"slides={actual} (footer TOTAL was {TOTAL})")
prs.save(OUT)
print(f"saved: {OUT}")
print(f"size_mb={OUT.stat().st_size/1e6:.2f}")
