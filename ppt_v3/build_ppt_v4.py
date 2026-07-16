#!/usr/bin/env python3
"""외삽 50분 세미나 v4 — outline_v4.md 기준 (다듬은 판).

원칙:
- 장표마다 논문 출처 배지(SOURCE) 명시
- 핵심 메시지 + 불릿 + Figure (가능하면 세 가지 모두)
- 논문으로 확인 가능한 내용만 (미검증 Extra R² 결합 우위 주장 금지)
- EQL/NALU/Monotonic은 필독 보조로 강등
"""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_ppt_v3 import (  # noqa: E402
    BASE,
    C_ACCENT,
    C_MUTED,
    C_PRIMARY,
    C_TEXT,
    C_WHITE,
    CONTENT_TOP,
    FIG_BOTTOM,
    MARGIN,
    MAX_FIG_H,
    ROOT,
    add_bullets,
    add_footer,
    add_header,
    add_image,
    add_key_message,
    add_pdf_ole,
    blank_slide,
    fig,
    set_run,
    slide_section,
    slide_table,
    sync_paper_copies,
)

OUT = BASE / "외삽_50분_발표자료_v4.pptx"

# GitHub 100MB 제한: --no-ole 시 PDF OLE 미첨부 (원문은 papers/ · papers_to_add/ 참고)
EMBED_PDF = True

C_BADGE = RGBColor(15, 118, 110)  # teal badge
C_BADGE_BG = RGBColor(204, 251, 241)

# 단축키 → PDF 상대경로 (ROOT 기준)
PDF = {
    "pfister": "extrapolation-papers/01_foundations/Pfister2024_Extrapolation-Aware_Nonparametric_Inference.pdf",
    "bartley": "extrapolation-papers/01_foundations/Bartley2019_Characterizing_Extrapolation_Multivariate.pdf",
    "xu": "extrapolation-papers/03_neural_network_extrapolation/Xu2021_How_Neural_Networks_Extrapolate.pdf",
    "bonnasse": "extrapolation-papers/03_neural_network_extrapolation/Bonnasse-Gahot2022_Interpolation_Extrapolation_NN.pdf",
    "ye": "extrapolation-papers/02_OOD_generalization/Ye2021_Theoretical_Framework_OOD.pdf",
    "liu": "extrapolation-papers/02_OOD_generalization/Liu2023_OOD_Generalization_Survey.pdf",
    "yuan": "extrapolation-papers/02_OOD_generalization/Yuan2022_OOD_Mechanics.pdf",
    "krueger": "extrapolation-papers/02_OOD_generalization/Krueger2021_Risk_Extrapolation_REx.pdf",
    "wu": "extrapolation-papers/02_OOD_generalization/Wu2025_OOD_Time_Series_Survey.pdf",
    "gulrajani": "extrapolation-papers/02_OOD_generalization/Gulrajani2020_In_Search_of_Lost_Domain_Generalization.pdf",
    "irm": "papers_to_add/02_OOD/Arjovsky2019_Invariant_Risk_Minimization.pdf",
    "scholkopf": "papers_to_add/02_OOD/Scholkopf2021_Toward_Causal_Representation_Learning.pdf",
    "moe": "papers_to_add/03_NN/Shazeer2017_Sparsely_Gated_MoE.pdf",
    "tent": "papers_to_add/03_NN/Wang2021_Tent_Test_Time_Adaptation.pdf",
    "adapter": "papers_to_add/03_NN/Rebuffi2017_Residual_Adapters.pdf",
    "fesser": "extrapolation-papers/03_neural_network_extrapolation/Fesser2023_Extrapolation_Failures_PINNs.pdf",
    "raissi": "papers_to_add/03_NN/Raissi2019_Physics_Informed_Neural_Networks.pdf",
}

V4_PAPERS = [
    ("Pfister 2024", "Hull·외삽 정의", PDF["pfister"]),
    ("Xu 2021", "NN 외삽 한계", PDF["xu"]),
    ("Ye 2021", "OOD·ERM 실패", PDF["ye"]),
    ("Arjovsky 2019", "IRM / spurious", PDF["irm"]),
    ("Shazeer 2017", "Sparse MoE", PDF["moe"]),
    ("Wang 2021", "Tent TTA", PDF["tent"]),
    ("Rebuffi 2017", "Residual Adapter", PDF["adapter"]),
    ("Gulrajani 2020", "DomainBed / ERM", PDF["gulrajani"]),
]


def pdf_path(*keys: str) -> list[Path]:
    """단축키 → 존재하는 Path 목록."""
    out = []
    for k in keys:
        rel = PDF.get(k, k)
        p = ROOT / rel if not Path(rel).is_absolute() else Path(rel)
        if p.exists():
            out.append(p)
        else:
            print(f"WARN missing PDF: {k} -> {p}")
    return out


def sync_all_v4_pdfs():
    """인용 논문 PDF를 ppt_v3/papers/ 에 복사."""
    papers_dir = BASE / "papers"
    papers_dir.mkdir(parents=True, exist_ok=True)
    for rel in PDF.values():
        src = ROOT / rel
        if not src.exists():
            continue
        dst = papers_dir / src.name
        if not dst.exists() or dst.stat().st_size != src.stat().st_size:
            shutil.copy2(src, dst)
            print(f"COPY {dst.name}")


def run_prep():
    subprocess.run([sys.executable, str(BASE / "generate_v4_diagrams.py")], check=False)
    venv = BASE / ".venv-figcrop" / "bin" / "python"
    layout = BASE / "extract_figures_layout.py"
    if venv.exists() and layout.exists():
        subprocess.run([str(venv), str(layout)], check=False)
    sync_paper_copies()
    sync_all_v4_pdfs()


def add_source(slide, text: str, top: float = 0.88):
    """장표 상단 논문 출처 배지 — 모든 콘텐츠 슬라이드에 사용."""
    box = slide.shapes.add_shape(1, Inches(MARGIN), Inches(top), Inches(13.333 - 2 * MARGIN), Inches(0.30))
    box.fill.solid()
    box.fill.fore_color.rgb = C_BADGE_BG
    box.line.color.rgb = C_BADGE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(2)
    r = tf.paragraphs[0].add_run()
    r.text = f"SOURCE  {text}"
    set_run(r, size=11, bold=True, color=C_BADGE)


def attach_pdfs(slide, pdfs: list[Path] | None, label: bool = True):
    """우하단 PDF OLE 아이콘 — 더블클릭으로 원문 열기. (--no-ole 시 경로 안내만)"""
    if not pdfs:
        return
    n = min(len(pdfs), 4)
    names = " · ".join(p.stem.split("_")[0] for p in pdfs[:n])
    if not EMBED_PDF:
        box = slide.shapes.add_textbox(Inches(9.2), Inches(6.55), Inches(3.8), Inches(0.35))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = f"PDF: ppt_v3/papers/ · {names}"
        set_run(r, size=8, bold=True, color=C_ACCENT)
        return
    # 최대 4개까지 가로로 배치 (우하단)
    icon_w, icon_h = 0.55, 0.68
    gap = 0.06
    total_w = n * icon_w + (n - 1) * gap
    left0 = 13.333 - MARGIN - total_w
    top = 6.22
    for i, p in enumerate(pdfs[:n]):
        x = left0 + i * (icon_w + gap)
        add_pdf_ole(slide, p, x, top, icon_w, icon_h)
    if label:
        box = slide.shapes.add_textbox(Inches(left0 - 0.05), Inches(top + icon_h - 0.02), Inches(max(total_w, 2.2)), Inches(0.22))
        tf = box.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = f"PDF 더블클릭  {names}"
        set_run(r, size=8, bold=True, color=C_ACCENT)


def add_key(slide, text: str, top: float = 1.22):
    add_key_message(slide, text, top=top)


def content_top_with_source() -> float:
    """출처 배지 사용 시 본문 시작 y."""
    return 1.72


def slide_rich_text(prs, section, title, subtitle, source, key, bullets, footer="", pdfs=None):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_source(slide, source)
    add_key(slide, key)
    # PDF 아이콘 공간 확보: 불릿 폭을 살짝 줄임
    add_bullets(slide, bullets, MARGIN, content_top_with_source(), 12.2, FIG_BOTTOM - content_top_with_source() - 0.85, size=15)
    attach_pdfs(slide, pdfs)
    add_footer(slide, footer or source)


def slide_fig_bullets(prs, section, title, subtitle, source, key, image, bullets, footer="", fig_ratio=0.58, caption="", pdfs=None):
    """Figure(좌) + 핵심 불릿(우) + SOURCE 배지 + PDF."""
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_source(slide, source)
    add_key(slide, key)
    top = content_top_with_source()
    usable = 13.333 - 2 * MARGIN
    fig_w = usable * fig_ratio
    text_w = usable - fig_w - 0.18
    # PDF 아이콘이 우측 하단에 오므로 불릿 높이만 살짝 줄임
    add_image(slide, image, MARGIN, top, fig_w, caption, max_h=MAX_FIG_H - 0.15, valign="center")
    add_bullets(slide, bullets, MARGIN + fig_w + 0.18, top, text_w, MAX_FIG_H - 0.95, size=13, compact=True)
    attach_pdfs(slide, pdfs)
    add_footer(slide, footer or source)


def slide_fig_bottom(prs, section, title, subtitle, source, key, image, bullets, footer="", fig_h_ratio=0.55, pdfs=None):
    """상단 불릿 + 하단 Figure + PDF."""
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_source(slide, source)
    add_key(slide, key)
    top = content_top_with_source()
    usable = 13.333 - 2 * MARGIN
    bullet_h = 1.55
    add_bullets(slide, bullets, MARGIN, top, usable - 0.9, bullet_h, size=13, compact=True)
    fig_top = top + bullet_h + 0.08
    add_image(slide, image, MARGIN, fig_top, usable - 0.85, "", max_h=FIG_BOTTOM - fig_top - 0.05, valign="center")
    attach_pdfs(slide, pdfs)
    add_footer(slide, footer or source)


def slide_dual_fig(prs, section, title, subtitle, source, key, img_l, img_r, bullets, footer="", pdfs=None):
    """두 Figure + 하단 요약 불릿 + PDF."""
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_source(slide, source)
    add_key(slide, key)
    top = content_top_with_source()
    usable = 13.333 - 2 * MARGIN
    half = (usable - 0.15) / 2
    fig_h = MAX_FIG_H * 0.68
    add_image(slide, img_l, MARGIN, top, half, "", max_h=fig_h, valign="center")
    add_image(slide, img_r, MARGIN + half + 0.15, top, half, "", max_h=fig_h, valign="center")
    add_bullets(slide, bullets, MARGIN, top + fig_h + 0.06, usable - 2.2, 1.0, size=12, compact=True)
    attach_pdfs(slide, pdfs)
    add_footer(slide, footer or source)


def slide_title(prs):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()
    for y, t, sz in [
        (1.6, "외삽 완전 정복", 42),
        (2.55, "이론 → 현상 → 적응 → 실전", 22),
        (3.35, "엄밀한 측정(Hull·Extra R²) · 구조(MoE·Adapter) · 적응(TTA)", 15),
        (4.3, "핵심 메시지: 외삽은 피할 수 없다. 측정 + 구조 + 적응으로 방어한다.", 14),
        (5.4, "v4 | 50분 · 국민대 IE Lab · 2026.07", 13),
    ]:
        b = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.8), Inches(0.7))
        tf = b.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = t
        set_run(r, size=sz, bold=sz >= 20, color=C_WHITE)


def slide_pdf_index(prs):
    """필수 8편 + 보조 인용 PDF 전부."""
    extras = [
        ("Bonnasse 2022", "Convex Hull", PDF["bonnasse"]),
        ("Fesser 2023", "PINN 외삽 실패", PDF["fesser"]),
        ("Yuan 2022", "OOD Mechanics", PDF["yuan"]),
        ("Liu 2023", "OOD Survey", PDF["liu"]),
        ("Krueger 2021", "REx", PDF["krueger"]),
        ("Bartley 2019", "다변량 CI", PDF["bartley"]),
        ("Raissi 2019", "PINN", PDF["raissi"]),
        ("Wu 2025", "OOD Time Series", PDF["wu"]),
        ("Schölkopf 2021", "인과 표현", PDF["scholkopf"]),
    ]
    all_papers = list(V4_PAPERS) + extras

    # page 1: top 8
    slide = blank_slide(prs)
    add_header(slide, "필독 논문 PDF (필수 8)", "아이콘 더블클릭 → 원문 열기", "APPENDIX")
    add_source(slide, "필수 8편 — 각 콘텐츠 장표에도 동일 PDF가 첨부되어 있음")
    add_key(slide, "발표 중에도 우하단 PDF 아이콘으로 바로 열 수 있음")
    for i, (short, role, rel) in enumerate(V4_PAPERS):
        col, row = i % 4, i // 4
        x = 0.45 + col * 3.2
        y = 1.95 + row * 2.25
        pdf = ROOT / rel
        if EMBED_PDF and pdf.exists():
            add_pdf_ole(slide, pdf, x + 1.0, y, 0.75, 0.9)
        box = slide.shapes.add_textbox(Inches(x), Inches(y + (1.0 if EMBED_PDF else 0.15)), Inches(2.9), Inches(0.9 if EMBED_PDF else 1.6))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        r = tf.paragraphs[0].add_run()
        fname = Path(rel).name
        r.text = f"{i + 1}. {short}\n{role}" + ("" if EMBED_PDF else f"\n→ papers/{fname}")
        set_run(r, size=12 if EMBED_PDF else 11, bold=True, color=C_PRIMARY)
    add_footer(slide, "복사본: ppt_v3/papers/" + ("" if EMBED_PDF else "  (GitHub용: OLE 미첨부)"))

    # page 2: extras
    slide2 = blank_slide(prs)
    add_header(slide2, "보조 인용 논문 PDF", "아이콘 더블클릭 → 원문 열기" if EMBED_PDF else "경로: ppt_v3/papers/", "APPENDIX")
    add_source(slide2, "Bonnasse · Fesser · Yuan · Liu · REx · Bartley · Raissi · Wu · Schölkopf")
    add_key(slide2, "콘텐츠 장표 SOURCE에 나온 보조 논문도 여기에 모음")
    for i, (short, role, rel) in enumerate(extras):
        col, row = i % 5, i // 5
        x = 0.35 + col * 2.55
        y = 1.95 + row * 2.35
        pdf = ROOT / rel
        if EMBED_PDF and pdf.exists():
            add_pdf_ole(slide2, pdf, x + 0.75, y, 0.7, 0.85)
        box = slide2.shapes.add_textbox(Inches(x), Inches(y + (0.95 if EMBED_PDF else 0.1)), Inches(2.4), Inches(0.9 if EMBED_PDF else 1.8))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = f"{short}\n{role}" + ("" if EMBED_PDF else f"\n→ {Path(rel).name}")
        set_run(r, size=11, bold=True, color=C_PRIMARY)
    add_footer(slide2, f"총 {len(all_papers)}편 · ppt_v3/papers/")


def build(skip_extract: bool = False):
    if not skip_extract:
        run_prep()
    else:
        subprocess.run([sys.executable, str(BASE / "generate_v4_diagrams.py")], check=False)
        sync_all_v4_pdfs()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== Part 0 · 도입 ==========
    slide_title(prs)

    slide_rich_text(
        prs, "PART 0", "왜 지금 외삽인가", "산업·센서·극한 조건",
        "동기 프레임 (학술 정의는 Pfister 2024 · OOD는 Ye/Liu)",
        "공정·센서·극한 조건은 항상 train 범위를 벗어난다",
        [
            "운영 조건(부하·온도·수명 구간)은 실험실 train 분포와 다름",
            "일반 ML/DL은 train 분포 안에서 최적화된 모델 → 범위 밖은 ‘보너스’가 아님",
            "범위 밖 예측을 그냥 연장하면 오차·위험 급증 (다변량 CI 붕괴: Bartley 2019 보조)",
            "오늘 흐름: 측정(Hull) → 현상(OOD/spurious) → 한계(ReLU·PINN) → 해법(MoE·Adapter·TTA) → 실전 평가",
        ],
        "도입 — 외삽은 선택이 아니라 운영 현실",
        pdfs=pdf_path("pfister", "ye", "bartley"),
    )

    if fig("storyline_roadmap.png").exists():
        slide_fig_bullets(
            prs, "PART 0", "발표 로드맵", "50분 스토리라인",
            "outline_v4 — 하이브리드(Hull·OOD·구조·적응)",
            "측정 → 현상 → 한계 → MoE/Adapter/TTA → Extra R² / strict holdout",
            fig("storyline_roadmap.png"),
            [
                "Part 0 (8′): Hull · 보간/외삽 · UQ",
                "Part 1 (10′): OOD · ERM 실패 · Flip/Regime · IRM",
                "Part 2 (5′): ReLU · PINN 한계",
                "Part 3 (15′): MoE · Adapter · TTA ← 핵심",
                "Part 4–5 (12′): Extra R² · N-CMAPSS · Takeaways",
            ],
            "시간 배분: 해법(Part 3)에 가장 많이 투자",
            fig_ratio=0.55,
        )
    else:
        slide_rich_text(
            prs, "PART 0", "발표 로드맵", "50분 흐름",
            "outline_v4",
            "측정 → 현상 → 한계 → 해법 → 실전",
            [
                "Part 0: Hull · 보간/외삽 · UQ (8′)",
                "Part 1: OOD · ERM 실패 · spurious / regime (10′)",
                "Part 2: ReLU · PINN 한계 (5′)",
                "Part 3: MoE · Adapter · TTA (15′) ← 핵심",
                "Part 4–5: 평가 철학 · N-CMAPSS · Takeaways (12′)",
            ],
        )

    if fig("pfister_extrapolation_fig2.png").exists():
        slide_fig_bullets(
            prs, "PART 0", "보간 ≠ 외삽", "Convex Hull 기준",
            "Pfister & Bühlmann (2024) · Fig.2 — Extrapolation-aware nonparametric inference",
            "train convex hull 안 = 보간, 밖 = 외삽",
            fig("pfister_extrapolation_fig2.png"),
            [
                "보간: 추가 가정 없이 조건부기댓값을 논할 수 있는 영역",
                "외삽: hull 밖 — 추정·예측이 본질적으로 어려워짐",
                "test가 hull 안이면 ‘외삽 성공’이라고 부르면 안 됨",
                "실무: 먼저 ‘어디가 밖인가’를 정의해야 평가가 성립",
            ],
            "Pfister & Bühlmann (2024) Fig.2",
            fig_ratio=0.60,
            pdfs=pdf_path("pfister"),
        )

    if fig("bonnasse_convex_hull_fig1.png").exists() and fig("pfister_rmse_extrap_fig3.png").exists():
        slide_dual_fig(
            prs, "PART 0", "Convex Hull + UQ", "고차원 · 외삽 오차",
            "Bonnasse-Gahot et al. (2022) Fig.1 · Pfister & Bühlmann (2024) Fig.3 · (보조) Bartley 2019",
            "고차원일수록 hull 밖이 다수 · 외삽 구간에서 오차·불확실성 폭증",
            fig("bonnasse_convex_hull_fig1.png"),
            fig("pfister_rmse_extrap_fig3.png"),
            [
                "Bonnasse: 내재 공간 hull ≠ 임베딩 공간 hull (A–D)",
                "Pfister: 외삽 거리↑ → RMSE / UQ 폭↑",
                "다변량 CI는 쉽게 붕괴 (Bartley 보조 인용)",
            ],
            pdfs=pdf_path("bonnasse", "pfister", "bartley"),
        )
    elif fig("bonnasse_convex_hull_fig1.png").exists():
        slide_fig_bullets(
            prs, "PART 0", "Convex Hull", "표현 공간의 함정",
            "Bonnasse-Gahot et al. (2022) Fig.1",
            "내재 공간 hull과 신경 표현 공간 hull은 다를 수 있다",
            fig("bonnasse_convex_hull_fig1.png"),
            [
                "파란점: 관측/학습 영역 · 빨간점: hull 밖/다른 지지",
                "고차원·비선형 매핑 후 ‘안쪽처럼 보이는’ 점이 실제로는 외삽일 수 있음",
                "→ 입력 공간 hull만으로도 부족, 평가 프로토콜이 필요",
            ],
            pdfs=pdf_path("bonnasse"),
        )

    # ========== Part 1 · OOD ==========
    slide_section(prs, "PART 1", "OOD 물리/구조 현상", "분포 이동 · ERM 실패 · spurious · regime")

    if fig("ood_shift_diagram.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "OOD란?", "Out-of-Distribution",
            "Liu et al. (2023) — Out-of-Distribution Generalization survey (정의·분류)",
            "train과 다른 test 분포 = OOD",
            fig("ood_shift_diagram.png"),
            [
                "P_test ≠ P_train 이면 OOD (세부 유형은 아래 장표)",
                "외삽과 직결: covariate shift — 입력 지지가 이동",
                "산업 데이터: 센서/운전점 이동이 곧 OOD",
                "‘정확도 높음’만으로는 OOD 안전성을 보장하지 않음",
            ],
            "Liu et al. (2023) OOD Survey",
            pdfs=pdf_path("liu"),
        )

    if fig("distribution_shift_types.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "Covariate Shift vs Concept Drift", "분포 이동 유형",
            "Liu et al. (2023) survey 정리 · (서적) Quinonero-Candela 2009는 비공개 → survey로 대체",
            "P(X) 변화 vs P(Y|X) 변화 — 대응 전략이 달라진다",
            fig("distribution_shift_types.png"),
            [
                "Covariate shift: 입력 분포만 이동 — Hull/외삽과 가장 직접 연결",
                "Label / concept shift: 라벨 규칙 자체가 변화",
                "운영: 두 유형이 섞여 나타남 → 단일 ‘보정’으로 해결하기 어려움",
                "오늘 해법(MoE·Adapter·TTA)은 주로 covariate/regime 이동을 겨냥",
            ],
            pdfs=pdf_path("liu"),
        )

    if fig("ye_ood_failure_fig2.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "ERM이 OOD에서 실패", "평균 최적 ≠ 환경 일반화",
            "Ye et al. (2021) Fig.2 · 관련: Krueger et al. (2021) Risk Extrapolation (REx)",
            "train 평균 손실 최소화 ≠ 다른 도메인/환경에서의 성공",
            fig("ye_ood_failure_fig2.png"),
            [
                "ERM은 관측된 train 환경의 평균에 맞춰짐",
                "도메인/환경이 바뀌면 같은 feature라도 성능이 붕괴할 수 있음",
                "REx: 환경 간 risk를 맞추려는 방향 (보조 인용)",
                "실무 함의: in-distribution CV 점수만 보고 배포하면 위험",
            ],
            "Ye et al. (2021) · Krueger et al. (2021)",
            pdfs=pdf_path("ye", "krueger"),
        )

    if fig("v4_correlation_flip.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "Partial Correlation Flip", "환경이 바뀌면 부호가 뒤집힐 수 있다",
            "개념 정리 — Yuan et al. (2022) mechanism/environment · Arjovsky IRM spurious · Schölkopf 2021",
            "정상 regime에서 보인 상관이 극한 regime에서 역전될 수 있다",
            fig("v4_correlation_flip.png"),
            [
                "단일 전역 회귀/NN은 ‘평균 상관’을 학습하기 쉬움",
                "Flip이 있는 구간을 섞어 학습하면 상쇄·왜곡",
                "해결 방향: regime 분할(MoE) 또는 불변 표현(IRM 계열)",
                "※ 수치 예시는 도메인별로 별도 검증 필요 (이 장표는 개념)",
            ],
            "Yuan 2022 · IRM 2019 (개념도)",
            pdfs=pdf_path("yuan", "irm", "scholkopf"),
        )

    if fig("irm_fig3_colored_mnist_setup.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "Spurious correlation", "IRM · Colored MNIST",
            "Arjovsky et al. (2019) — Invariant Risk Minimization · Fig.3 (Colored MNIST 설정)",
            "우연 상관(색–라벨)을 배우면 환경이 바뀌는 순간 붕괴한다",
            fig("irm_fig3_colored_mnist_setup.png"),
            [
                "Train: 색이 라벨과 강하게 상관 → ERM이 색을 암기",
                "Test: 색–라벨 상관이 바뀌면 정확도 급락",
                "IRM 메시지: 환경마다 변하지 않는 예측기를 찾자",
                "산업 비유: 센서 drift·운전점이 ‘색’ 역할 → spurious feature",
            ],
            "Arjovsky et al. (2019) Fig.3",
            pdfs=pdf_path("irm"),
        )

    if fig("yuan_ood_mechanics_fig1.png").exists() and fig("v4_regime_shift.png").exists():
        slide_dual_fig(
            prs, "PART 1", "Regime / Mechanism Shift", "환경 파라미터 s가 바뀌면",
            "Yuan et al. (2022) Fig.1 · (시계열) Wu (2025) OOD Time Series survey · 우: 개념도",
            "동일 과제라도 환경(s)·regime이 바뀌면 입·출력 통계가 함께 이동",
            fig("yuan_ood_mechanics_fig1.png"),
            fig("v4_regime_shift.png"),
            [
                "Train env ≠ Test env → mechanism shift",
                "시계열/공정: 저부하→고부하 전이가 비선형일 수 있음",
                "단일 전역 모델 취약 → 분할·적응 해법으로 연결",
            ],
            pdfs=pdf_path("yuan", "wu"),
        )
    elif fig("yuan_ood_mechanics_fig1.png").exists():
        slide_fig_bullets(
            prs, "PART 1", "Regime / Mechanism Shift", "Yuan 2022",
            "Yuan et al. (2022) Fig.1 — OOD Mechanics",
            "환경 파라미터 s가 바뀌면 입·출력 통계가 함께 이동",
            fig("yuan_ood_mechanics_fig1.png"),
            [
                "동일 ‘과제’라도 환경(s)에 따라 분포 변화",
                "상관·역학이 regime에 따라 달라질 수 있음",
                "→ 단일 ERM 모델의 취약점 · MoE/TTA로 연결",
            ],
            pdfs=pdf_path("yuan"),
        )

    # ========== Part 2 · 한계 ==========
    slide_section(prs, "PART 2", "왜 단순 모델로는 부족한가", "ReLU · PINN — 구조 없이는 외삽이 깨진다")

    if fig("xu_relu_extrapolation_fig1.png").exists():
        slide_fig_bottom(
            prs, "PART 2", "ReLU MLP 외삽 실패", "조각별 직선",
            "Xu et al. (2021) ICLR — How Neural Networks Extrapolate · Fig.1",
            "훈련 영역(파란 영역) 밖에서는 ReLU MLP가 조각별 affine으로 직선 연장된다",
            fig("xu_relu_extrapolation_fig1.png"),
            [
                "학습 구간에서는 복잡한 곡면도 근사 가능",
                "hull 밖: 마지막 선형 조각이 무한히 연장 → 목표 함수와 어긋남",
                "‘더 깊게’만으로는 외삽 형태가 자동으로 좋아지지 않음",
            ],
            pdfs=pdf_path("xu"),
        )

    if fig("fesser_pinn_failure_fig1.png").exists():
        slide_fig_bullets(
            prs, "PART 2", "PINN도 외삽에서 깨짐", "물리 제약이 만능은 아님",
            "Fesser et al. (2023) Fig.1 · PINN 원형: Raissi et al. (2019)",
            "물리 손실을 넣어도 외삽 구간 잔차·오차가 급증할 수 있다",
            fig("fesser_pinn_failure_fig1.png"),
            [
                "PINN: PDE/물리 잔차를 loss에 포함",
                "그래도 학습 지지 밖에서는 해가 불안정해질 수 있음",
                "함의: ‘제약만’이 아니라 평가·분할·적응이 함께 필요",
                "다음: MoE(분할) · Adapter(잔차) · TTA(적응)",
            ],
            "Fesser 2023 · Raissi 2019",
            pdfs=pdf_path("fesser", "raissi"),
        )

    if fig("v4_solution_preview.png").exists():
        slide_fig_bullets(
            prs, "PART 2", "해법 예고", "세 축으로 방어",
            "모듈별 근거 — MoE: Shazeer 2017 · Adapter: Rebuffi 2017 · TTA: Wang 2021 (Tent)",
            "분할(MoE) · 잔차(Adapter) · 적응(TTA) — Part 3에서 논문 Figure로 확인",
            fig("v4_solution_preview.png"),
            [
                "MoE: regime별로 다른 Expert",
                "Adapter: 공유 backbone + 작은 잔차만 학습",
                "TTA: 배포 후 unlabeled stream에서 소수 파라미터 적응",
                "결합은 설계 제안 — 결합 Extra R² 우위는 이 자료에서 주장하지 않음",
            ],
            "Part 3 예고",
            fig_ratio=0.55,
            pdfs=pdf_path("moe", "adapter", "tent"),
        )

    # ========== Part 3 · 해법 ==========
    slide_section(prs, "PART 3", "해법 딥다이브", "MoE · Adapter · TTA — 역할 분담")

    if fig("v4_combine_pipeline.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "결합 아키텍처 (제안)", "역할 분담 파이프라인",
            "제안 구조 — 개별 모듈만 논문 근거 (Shazeer / Rebuffi / Wang Tent)",
            "Gating → Expert + Residual Adapter → (고엔트로피 시) Tent-style TTA",
            fig("v4_combine_pipeline.png"),
            [
                "입력 x → Gating이 regime/expert 선택",
                "선택된 Expert + 해당 Adapter로 예측",
                "불확실/OOD면 TTA로 BN 등 소수 파라미터만 업데이트",
                "비교 시 항상 ERM baseline 포함 (DomainBed)",
            ],
            "설계 제안 · 모듈별 논문 근거",
            fig_ratio=0.62,
            pdfs=pdf_path("moe", "adapter", "tent"),
        )

    if fig("shazeer_moe_fig1.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "MoE 원리", "Sparsely-Gated MoE",
            "Shazeer et al. (2017) — Outrageously Large Neural Networks · Fig.1",
            "Gating이 Expert를 선택 · Sparse 조건부 계산으로 용량을 나눈다",
            fig("shazeer_moe_fig1.png"),
            [
                "입력마다 소수 Expert만 활성화 (sparse)",
                "regime/패턴별로 다른 함수를 학습하기 유리",
                "외삽 관점: ‘전 구간 단일 모델’ 부담을 분할",
                "원 논문 맥락은 대규모 언어/번역 — 아이디어를 공정 외삽에 이식",
            ],
            "Shazeer et al. (2017) Fig.1",
            pdfs=pdf_path("moe"),
        )

    # MoE 함정 — figure if available
    if fig("shazeer_moe_fig2.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "MoE 함정", "Gating이 틀리면 끝",
            "Shazeer et al. (2017) — load balancing / expert utilization 이슈 · Fig.2(관련)",
            "Gating이 OOD에서 잘못 고르면 Expert 선택 자체가 위험하다",
            fig("shazeer_moe_fig2.png"),
            [
                "Expert collapse / load imbalance — 학습 자체가 깨질 수 있음",
                "train에 없던 regime → gating 오분류 가능",
                "MoE만으로 ‘외삽 해결’이라 말하기 어려움",
                "→ Adapter(미세 보정) · TTA(배포 후 적응)와 역할 분담",
            ],
            "Shazeer et al. (2017)",
            pdfs=pdf_path("moe"),
        )
    else:
        slide_rich_text(
            prs, "PART 3", "MoE 함정", "논문이 지적하는 실패 모드",
            "Shazeer et al. (2017)",
            "Gating이 OOD에서 잘못 고르면 Expert 선택 자체가 위험",
            [
                "Expert collapse / load imbalance (Shazeer 2017에서 다루는 학습 이슈)",
                "train에 없던 regime → gating 오분류 가능",
                "따라서 MoE만으로 ‘외삽 해결’이라 말하기 어려움",
                "→ Adapter(미세 보정) · TTA(배포 후 적응)와 역할 분담",
            ],
            pdfs=pdf_path("moe"),
        )

    if fig("rebuffi_adapter_fig2.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "Residual Adapter", "작은 잔차만 학습",
            "Rebuffi, Bilen, Vedaldi (2017) — Learning multiple visual domains with residual adapters · Fig.2",
            "공유 backbone(w)은 고정하고, 파란 adapter 파라미터만 도메인별로 학습",
            fig("rebuffi_adapter_fig2.png"),
            [
                "Visual Decathlon: 다중 도메인에서 residual adapter로 조향",
                "장점: full fine-tune 대비 파라미터·안정성",
                "외삽/새 regime: 작은 adapter 추가로 보정하는 설계에 대응",
                "한계: regime·도메인 구분이 애매하면 부착 전략이 어려움",
            ],
            "Rebuffi et al. (2017) Fig.2",
            fig_ratio=0.62,
            pdfs=pdf_path("adapter"),
        )
    else:
        slide_rich_text(
            prs, "PART 3", "Residual Adapter", "Rebuffi et al. 2017",
            "Rebuffi, Bilen, Vedaldi (2017)",
            "공유 backbone은 고정하고, 작은 잔차 모듈만 학습",
            [
                "Visual Decathlon: 다중 도메인에서 adapter residual로 조향",
                "장점: full fine-tune 대비 파라미터·안정성",
                "외삽/새 regime: 작은 adapter 추가로 보정",
                "한계: regime 라벨이 애매하면 부착이 어려움",
            ],
            pdfs=pdf_path("adapter"),
        )

    slide_rich_text(
        prs, "PART 3", "Adapter × Regime", "MoE expert와 연결",
        "설계 연결 — Rebuffi 2017(잔차) · Shazeer 2017(분할) · 결합은 제안",
        "새 regime ≈ 작은 adapter (또는 expert 헤드) 추가",
        [
            "MoE: 큰 용량을 expert로 분할 — ‘어느 함수’를 쓸지 선택",
            "Adapter: 공유 표현 위에 가벼운 잔차 보정 — ‘얼마나 조향’할지",
            "실무 설계: gating/regime 추정 → 해당 expert·adapter 활성화",
            "둘 다 전 모델 재학습을 피하려는 방향 (각 논문 설정에서 검증)",
            "주의: ‘MoE+Adapter가 항상 이긴다’는 실험 결과는 아직 이 자료의 증거가 아님",
        ],
        pdfs=pdf_path("adapter", "moe"),
    )

    if fig("wang_tent_fig1.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "TTA — Tent", "Test-time Entropy Minimization",
            "Wang et al. (2021) — Tent: Fully Test-Time Adaptation by Entropy Minimization · Fig.1",
            "라벨 없는 test stream에서 entropy를 낮추도록 소수 파라미터를 적응",
            fig("wang_tent_fig1.png"),
            [
                "고엔트로피 ≈ 모델이 불확실한 영역 (오차와 상관)",
                "주로 BN 통계/affine 등 소수 파라미터만 업데이트",
                "배포 후 분포 이동에 대응하는 축",
                "외삽 운영: Hull/OOD 경보와 결합해 on/off",
            ],
            "Wang et al. (2021) Tent Fig.1",
            pdfs=pdf_path("tent"),
        )

    if fig("wang_tent_fig3.png").exists():
        slide_fig_bullets(
            prs, "PART 3", "TTA 한계", "적응은 공짜가 아니다",
            "Wang et al. (2021) Tent · Fig.3 및 논문 논의 · 실무 제약",
            "latency · forgetting · noisy stream — 켜는 조건을 명시해야 한다",
            fig("wang_tent_fig3.png"),
            [
                "매 배치 업데이트 → 지연·불안정 가능",
                "긴 스트림에서 이전 적응을 잊을 수 있음",
                "노이즈·이상치에 entropy 신호가 오염되면 오적응",
                "실무 규칙: 고엔트로피/OOD 경보 시에만 on",
            ],
            "Wang et al. (2021) Tent",
            pdfs=pdf_path("tent"),
        )
    else:
        slide_rich_text(
            prs, "PART 3", "TTA 한계", "Tent 논문·실무 제약",
            "Wang et al. (2021) Tent",
            "적응은 공짜가 아니다 — latency · forgetting · noisy stream",
            [
                "매 배치 업데이트 → 지연·불안정 가능",
                "긴 스트림에서 이전 적응을 잊을 수 있음",
                "노이즈·이상치에 entropy 신호가 오염되면 오적응",
                "실무: 고엔트로피/OOD 경보 시에만 on",
            ],
            pdfs=pdf_path("tent"),
        )

    # 역할 분담 표
    slide = blank_slide(prs)
    add_header(slide, "왜 셋이 필요한가", "역할 분담 표", "PART 3")
    add_source(slide, "MoE: Shazeer 2017 · Adapter: Rebuffi 2017 · TTA: Wang 2021 · ERM 비교: Gulrajani & Lopez-Paz 2020 DomainBed")
    add_key(slide, "비교 실험에는 항상 ERM baseline을 넣는다 (DomainBed 원칙)")
    # reuse table helper pattern inline for source-aware layout
    headers = ["방법", "푸는 문제", "실패 모드", "근거 논문"]
    rows = [
        ["MoE", "단일 모델 전 regime 과적합", "gating 오분류 · collapse", "Shazeer 2017"],
        ["Adapter", "full fine-tune 비용·불안정", "regime 라벨 애매", "Rebuffi 2017"],
        ["TTA", "배포 후 분포 이동", "latency · drift · 오적응", "Wang 2021"],
        ["ERM", "공정 비교 baseline", "spurious 암기", "Gulrajani 2020"],
    ]
    rs, cs = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(
        rs, cs, Inches(MARGIN), Inches(content_top_with_source()),
        Inches(13.333 - 2 * MARGIN), Inches(min(0.48 * rs, 4.2)),
    ).table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = C_PRIMARY
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=12, bold=True, color=C_WHITE)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=12)
    attach_pdfs(slide, pdf_path("moe", "adapter", "tent", "gulrajani"))
    add_footer(slide, "Gulrajani & Lopez-Paz (2020) DomainBed — ERM을 빼면 주장을 믿기 어렵다")

    # ========== Part 4 · 실험·검증 ==========
    slide_section(prs, "PART 4", "실험·검증", "평가 철학 · Extra R² · N-CMAPSS")

    if fig("v4_eval_philosophy.png").exists():
        slide_fig_bullets(
            prs, "PART 4", "평가 철학", "random split의 함정",
            "Ye et al. (2021) · Pfister & Bühlmann (2024) · Gulrajani & Lopez-Paz (2020) DomainBed",
            "random split만으로는 보간 평가일 수 있다 — holdout 축과 hull을 명시",
            fig("v4_eval_philosophy.png"),
            [
                "Holdout 축: 고부하 · late cycle · 미관측 unit 등",
                "Train hull 밖 test를 따로 리포트 (Pfister)",
                "OOD 주장 시 ERM과 동일 HP·모델로 비교 (DomainBed)",
                "Ablation: 모듈을 하나씩 켠 결과만 ‘기여’로 인정",
            ],
            "평가 설계 원칙",
            fig_ratio=0.58,
            pdfs=pdf_path("ye", "pfister", "gulrajani"),
        )
    else:
        slide_rich_text(
            prs, "PART 4", "평가 철학", "random split의 함정",
            "Ye 2021 · Pfister 2024 · Gulrajani 2020",
            "random split만으로는 보간 평가일 수 있다",
            [
                "Holdout 축을 명시: 고부하 · late cycle · 미관측 unit",
                "Train hull 밖 test를 따로 리포트",
                "OOD 주장 시 ERM과 동일 HP·모델로 비교",
                "Ablation으로 모듈 기여를 분리",
            ],
            pdfs=pdf_path("ye", "pfister", "gulrajani"),
        )

    slide_rich_text(
        prs, "PART 4", "Extra R²", "순수 외삽 구간 지표",
        "운영 지표 정의 — 정신은 Pfister(hull) · Ye(OOD 실험) · 수식은 랩 프로토콜",
        "전체 R²와 별도로, hull 밖(또는 holdout 축)만의 R²를 본다",
        [
            "운영 정의: Extra R² = R²( { (x,y) | x ∉ Conv(X_train) } )",
            "또는 프로토콜 holdout 축(예: 고 TRA · strict_late)만 사용",
            "전체 R²↑ 만으로 외삽 성공을 주장하지 않음",
            "체크리스트: Holdout 축 · Hull 밖 · ERM baseline · Ablation · Flip/Regime 구간별 리포트",
            "※ 아래 N-CMAPSS는 이 철학을 적용한 랩 실험 예시 (공개 논문 수치 아님)",
        ],
        pdfs=pdf_path("pfister", "ye"),
    )

    slide_rich_text(
        prs, "PART 4", "데이터셋 · Holdout 축", "무엇을 막을지 먼저 정한다",
        "일반론: Ye / Pfister · 적용 예: N-CMAPSS (랩 프로토콜)",
        "벤치마크마다 holdout 축이 있어야 ‘외삽 실험’이 된다",
        [
            "공개 벤치 예: 부하/온도 구간 holdout (UCI Naval / CCPP 등 — 축만 정의하면 됨)",
            "본 랩: N-CMAPSS — TRA(고부하) · unit · strict_late holdout",
            "Train: 저 TRA / Test: 고 TRA → covariate shift",
            "random unit split만 하면 보간에 가까운 평가가 될 수 있음",
            "Naval/CCPP 수치 비교표는 이 자료에서 확정 결과로 제시하지 않음",
        ],
        pdfs=pdf_path("ye", "pfister"),
    )

    if fig("ncmapss_results_chart.png").exists():
        slide_fig_bullets(
            prs, "PART 4", "N-CMAPSS 결과 예시", "랩 실험 (APEX-Guard)",
            "랩 내부 strict_extrap — 공개 논문 수치 아님 · MoE+TTA 결합 승리가 아님",
            "strict holdout에서 구조·제약 모델 vs 강한 in-dist 모델",
            fig("ncmapss_results_chart.png"),
            [
                "목적: ‘엄밀한 holdout이 순위를 바꾼다’는 교훈",
                "in-dist에 강한 모델이 OOD에서 질 수 있음",
                "MoE+Adapter+TTA 결합 Extra R² 우위는 추가 실험 필요",
            ],
        )
    else:
        slide = blank_slide(prs)
        add_header(slide, "N-CMAPSS strict_late 예시", "랩 실험", "PART 4")
        add_source(slide, "랩 내부 결과 — 공개 논문 수치가 아님 · MoE+TTA 결합 승리가 아님")
        add_key(slide, "제약+단조 구조가 TabPFN보다 OOD(고 TRA)에서 유리했던 사례")
        headers = ["모델", "RMSE", "R²", "비고"]
        rows = [
            ["APEX-Guard", "3.26", "0.947", "제약+구조"],
            ["TabPFN", "3.80", "0.927", "강한 in-dist"],
            ["Transformer", "7.31", "—", "OOD 취약"],
        ]
        table = slide.shapes.add_table(
            4, 4, Inches(MARGIN), Inches(content_top_with_source()),
            Inches(13.333 - 2 * MARGIN), Inches(2.2),
        ).table
        for j, h in enumerate(headers):
            c = table.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = C_PRIMARY
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=13, bold=True, color=C_WHITE)
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                c = table.cell(i, j)
                c.text = val
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        set_run(r, size=13)
        add_bullets(
            slide,
            [
                "교훈: strict split이 없으면 ‘승리’는 환상일 수 있음",
                "이 표는 MoE+Adapter+TTA 결합 결과가 아님",
            ],
            MARGIN, 4.3, 12.8, 1.5, size=14,
        )
        add_footer(slide, "랩 내부 strict_extrap 결과")

    slide_rich_text(
        prs, "PART 4", "실전 교훈", "측정이 모델보다 먼저",
        "평가 철학 + 랩 사례 정리 · 결합 Extra R² 우위는 미주장",
        "strict split 없으면 승리는 환상이다",
        [
            "L1: strict holdout으로 ‘순수 외삽’ 구간을 정의",
            "L2: 도메인 지식 제약(단조·물리)으로 구조를 박음",
            "L3: 배포 후 적응(TTA)은 옵션 — 켜는 조건을 명시",
            "다음 실험: Naval/CCPP·N-CMAPSS에서 MoE / Adapter / TTA ablation + Extra R²",
            "오늘 주장: 측정 프로토콜 + 역할 분담 설계 — ‘결합 완승’ 아님",
        ],
        pdfs=pdf_path("pfister", "ye", "gulrajani"),
    )

    # ========== Part 5 · 결론 ==========
    slide_section(prs, "PART 5", "결론", "Takeaways · 필독 · Q&A")

    slide = blank_slide(prs)
    add_header(slide, "Key Takeaways", "네 줄로 요약", "PART 5")
    add_source(slide, "Pfister · Ye/IRM/Yuan · Shazeer/Rebuffi/Wang · DomainBed/Extra R²")
    add_key(slide, "측정 + 현상 이해 + 구조/적응 + 엄밀한 평가")
    points = [
        "① Hull: 보간≠외삽 — hull 밖을 따로 측정 (Pfister 2024)",
        "② Spurious/Regime: 환경이 바뀌면 상관·역학이 흔들림 (IRM · Yuan)",
        "③ MoE·Adapter·TTA: 분할·잔차·적응 — 각각 논문으로 검증된 역할",
        "④ Extra R² / strict holdout: 전체 R²만 보고 외삽 성공을 말하지 말 것",
    ]
    add_bullets(slide, points, MARGIN, content_top_with_source(), 12.0, 4.0, size=16)
    attach_pdfs(slide, pdf_path("pfister", "irm", "moe", "xu"))
    add_footer(slide, "한계·Future는 다음 장표")

    slide_rich_text(
        prs, "PART 5", "한계 · Future", "오늘 주장하지 않은 것",
        "열린 문제 — Tent latency · MoE gating OOD · Adapter 라벨 · 인과(Schölkopf)",
        "결합 파이프라인의 대규모 Extra R² 우위는 아직 이 자료의 증거가 아님",
        [
            "TTA: latency · forgetting · noisy stream",
            "MoE: Gating OOD 오분류",
            "Adapter: 부착을 위한 regime 라벨 문제",
            "인과 표현(Schölkopf 2021)과의 결합은 열린 문제",
            "다음: Naval/CCPP·N-CMAPSS ablation + Extra R² + Flip 구간 리포트",
        ],
        pdfs=pdf_path("tent", "moe", "adapter", "scholkopf"),
    )

    slide_rich_text(
        prs, "PART 5", "필독 논문 Top 8", "슬라이드 주장의 근거",
        "필수 8편 + 보조 · EQL/NALU/Monotonic은 구조 외삽 보조 읽기",
        "발표에서 쓴 주장은 아래 논문으로만 뒷받침한다",
        [
            "1 Pfister 2024 — Hull·외삽 정의",
            "2 Xu 2021 — ReLU MLP 외삽 한계",
            "3 Ye 2021 — OOD framework / ERM 실패",
            "4 Arjovsky 2019 — IRM / spurious",
            "5 Shazeer 2017 — Sparse MoE",
            "6 Wang 2021 — Tent (TTA)",
            "7 Rebuffi 2017 — Residual Adapters",
            "8 Gulrajani 2020 — DomainBed / ERM baseline",
            "보조: Bonnasse-Gahot · Fesser · Yuan · Bartley · Krueger REx · Raissi · Liu survey",
            "구조 NN 보조(필독 강등): Martius EQL · Trask NALU · Runje Monotonic",
            "→ 다음 슬라이드에서 필수 8편 PDF 전부 더블클릭 가능",
        ],
        "다음 슬라이드: PDF 더블클릭",
        pdfs=pdf_path("pfister", "xu", "ye", "irm"),
    )

    slide_pdf_index(prs)

    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()
    for y, t, sz in [
        (2.6, "Q & A", 44),
        (3.7, "질문 · 토론 · 다음 실험 아이디어", 16),
    ]:
        b = s.shapes.add_textbox(Inches(2.5), Inches(y), Inches(8.5), Inches(0.8))
        tf = b.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = t
        set_run(r, size=sz, bold=sz >= 30, color=C_WHITE)

    prs.save(OUT)
    shutil.copy(OUT, ROOT / "extrapolation-papers" / OUT.name)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    skip = "--skip-extract" in sys.argv
    if "--no-ole" in sys.argv:
        EMBED_PDF = False
        print("NOTE: building without PDF OLE embeds (GitHub-sized)")
    build(skip_extract=skip)
