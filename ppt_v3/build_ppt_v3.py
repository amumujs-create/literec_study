#!/usr/bin/env python3
"""외삽 완전정복 v3.1 PPT — 스토리라인 + 섹션 정리 + Figure 재점검."""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from paper_catalog import PAPERS

ROOT = Path(__file__).resolve().parent.parent
BASE = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = BASE / "외삽_50분_발표자료_v3.pptx"
PDF_ICON = BASE / "assets" / "pdf_icon.png"
PAPERS_DIR = BASE / "papers"

# Palette
C_BG = RGBColor(248, 250, 252)
C_PRIMARY = RGBColor(30, 64, 110)
C_ACCENT = RGBColor(14, 116, 144)
C_RECAP = RGBColor(236, 245, 250)
C_TEXT = RGBColor(30, 41, 59)
C_MUTED = RGBColor(100, 116, 139)
C_WHITE = RGBColor(255, 255, 255)
C_KEY_BG = RGBColor(224, 242, 254)
C_ORANGE = RGBColor(194, 65, 12)

# Layout (inches) — figure-first, compact chrome
MARGIN = 0.22
KEY_TOP = 1.22
KEY_H = 0.40
CONTENT_TOP = 1.68
FOOTER_TOP = 7.02
FIG_BOTTOM = 6.92
MAX_FIG_H = FIG_BOTTOM - CONTENT_TOP


def fig(name: str) -> Path:
    return FIG / name


def run_prep():
    layout_py = BASE / "extract_figures_layout.py"
    venv_python = BASE / ".venv-figcrop" / "bin" / "python"
    model_dir = BASE / "models" / "PP-DocLayoutV2_infer"
    # 레이아웃 모델(PP-DocLayoutV2)이 있으면 우선 사용, 없으면 caption 기반 fallback
    if layout_py.exists() and venv_python.exists() and model_dir.exists():
        r = subprocess.run([str(venv_python), str(layout_py)], check=False)
        if r.returncode == 0:
            sync_paper_copies()
            return
    subprocess.run([sys.executable, str(BASE / "extract_figures_smart.py")], check=False)
    sync_paper_copies()


def sync_paper_copies():
    """PPT 옆에 papers/ 폴더로 PDF 복사 (하이퍼링크·백업용)."""
    import shutil

    PAPERS_DIR.mkdir(parents=True, exist_ok=True)
    for p in PAPERS:
        src = ROOT / p["pdf"]
        if not src.exists():
            continue
        dst = PAPERS_DIR / src.name
        if not dst.exists() or dst.stat().st_size != src.stat().st_size:
            shutil.copy2(src, dst)
            print(f"COPY {dst.name}")


def add_pdf_ole(slide, pdf_path: Path, left, top, width=0.95, height=1.15):
    """PDF를 OLE로 임베드 — 더블클릭하면 열림."""
    if not pdf_path.exists():
        return None
    icon = str(PDF_ICON) if PDF_ICON.exists() else None
    return slide.shapes.add_ole_object(
        str(pdf_path),
        "AcroExch.Document",
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        icon_file=icon,
        icon_width=Inches(width),
        icon_height=Inches(height),
    )


def slide_paper_index(prs):
    """부록 인덱스 — 모든 논문 PDF OLE 아이콘."""
    slide = blank_slide(prs)
    add_header(slide, "부록 · 논문 PDF 첨부", "아이콘 더블클릭 → PDF 열기", "APPENDIX")
    add_key_message(slide, "핵심 논문 13편 · PPT에 임베드됨 · papers/ 폴더에도 복사본 있음")

    cols, rows = 5, 3
    cell_w, cell_h = 2.4, 1.55
    origin_x, origin_y = 0.35, 1.85
    for i, p in enumerate(PAPERS):
        r, c = divmod(i, cols)
        if r >= rows:
            break
        x = origin_x + c * cell_w
        y = origin_y + r * cell_h
        pdf = ROOT / p["pdf"]
        add_pdf_ole(slide, pdf, x + 0.65, y, 0.7, 0.85)
        box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.9), Inches(cell_w - 0.1), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = f"{p['short']}\n{p['act']}"
        set_run(run, size=11, bold=True, color=C_PRIMARY)
    add_footer(slide, f"papers/ 경로: {PAPERS_DIR}")


def slide_paper_detail(prs, paper: dict):
    """논문별 핵심 Figure + PDF 첨부 슬라이드."""
    slide = blank_slide(prs)
    add_header(slide, paper["short"], paper["title"][:48], paper["act"])
    add_key_message(slide, paper["why"])

    # PDF OLE (right)
    pdf = ROOT / paper["pdf"]
    add_pdf_ole(slide, pdf, 12.0, 1.85, 0.95, 1.15)
    note = slide.shapes.add_textbox(Inches(11.7), Inches(3.05), Inches(1.5), Inches(0.7))
    ntf = note.text_frame
    ntf.word_wrap = True
    ntf.clear()
    nr = ntf.paragraphs[0].add_run()
    nr.text = "PDF\n더블클릭"
    set_run(nr, size=11, bold=True, color=C_ACCENT)

    # figures
    fig_names = list(paper["figs"].values())
    existing = [fig(n) for n in fig_names if fig(n).exists()]
    if existing:
        n = len(existing)
        if n == 1:
            add_image(slide, existing[0], MARGIN, CONTENT_TOP, 11.3, max_h=MAX_FIG_H * 0.95)
        else:
            w = min(5.5, 11.2 / n)
            for i, path in enumerate(existing[:3]):
                add_image(slide, path, MARGIN + i * (w + 0.12), CONTENT_TOP, w, max_h=MAX_FIG_H * 0.9)
    else:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(10), Inches(1.5))
        tf = box.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = "핵심 Figure 자동 추출 없음 — 오른쪽 PDF를 열어 확인하세요."
        set_run(r, size=16, color=C_MUTED)

    add_footer(slide, pdf.name if pdf.exists() else paper["pdf"])


def set_run(run, size=16, bold=False, color=C_TEXT, name="맑은 고딕"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()
    return slide


def add_header(slide, title: str, subtitle: str = "", section: str = ""):
    stripe = slide.shapes.add_shape(1, 0, 0, Inches(0.10), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = C_ACCENT
    stripe.line.fill.background()
    bar = slide.shapes.add_shape(1, Inches(0.10), 0, Inches(13.23), Inches(0.78))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    if section:
        sbox = slide.shapes.add_textbox(Inches(0.28), Inches(0.06), Inches(2.6), Inches(0.26))
        stf = sbox.text_frame
        stf.clear()
        sr = stf.paragraphs[0].add_run()
        sr.text = section
        set_run(sr, size=10, color=RGBColor(180, 210, 230))
    headline = title if not subtitle else f"{title}  ·  {subtitle}"
    tbox = slide.shapes.add_textbox(Inches(0.28), Inches(0.24), Inches(12.7), Inches(0.48))
    tf = tbox.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = headline
    set_run(r, size=22, bold=True, color=C_WHITE)


def add_key_message(slide, text: str, top=KEY_TOP):
    box = slide.shapes.add_shape(1, Inches(MARGIN), Inches(top), Inches(13.333 - 2 * MARGIN), Inches(KEY_H))
    box.fill.solid()
    box.fill.fore_color.rgb = C_KEY_BG
    box.line.color.rgb = C_ACCENT
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"핵심  {text}"
    set_run(r, size=14, bold=True, color=C_PRIMARY)


def add_bullets(slide, items, left, top, width, height, size=15, compact=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    gap = Pt(2 if compact else 4)
    for i, item in enumerate(items):
        text, level = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = gap
        p.line_spacing = 1.05 if compact else 1.15
        r = p.add_run()
        r.text = text
        set_run(r, size=size - level * 1)


def add_caption(slide, text: str, left, top, width=5.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=9, color=C_MUTED)


def _fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with PILImage.open(path) as im:
        w, h = im.size
    ar = w / h
    width = max_w
    height = width / ar
    if height > max_h:
        height = max_h
        width = height * ar
    return width, height


def add_footer(slide, text: str):
    if not text:
        return
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(FOOTER_TOP), Inches(12.9), Inches(0.28))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    set_run(r, size=9, color=C_MUTED)


def add_image(
    slide,
    path: Path,
    left,
    top,
    max_w,
    caption: str = "",
    max_h: float = MAX_FIG_H,
    valign: str = "top",
):
    if not path.exists():
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(max_w), Inches(0.5))
        box.text_frame.text = f"[Figure 없음: {path.name}]"
        return 0.0

    w_in, h_in = _fit_image(path, max_w, max_h)
    y = top + (max_h - h_in) / 2 if valign == "center" else top
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(y), width=Inches(w_in), height=Inches(h_in))
    frame = slide.shapes.add_shape(
        1, pic.left - Pt(1), pic.top - Pt(1), pic.width + Pt(2), pic.height + Pt(2)
    )
    frame.fill.background()
    frame.line.color.rgb = RGBColor(203, 213, 225)
    frame.line.width = Pt(0.75)
    slide.shapes._spTree.remove(frame._element)
    slide.shapes._spTree.insert(2, frame._element)
    if caption:
        cap_y = y + h_in + 0.03
        add_caption(slide, caption, left, cap_y, w_in)
    return y + h_in


def slide_section(prs, num: str, title: str, subtitle: str):
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()
    for y, text, sz, bold in [(2.3, num, 18, False), (2.8, title, 38, True), (4.0, subtitle, 18, False)]:
        b = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(1))
        tf = b.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = text
        set_run(r, size=sz, bold=bold, color=C_WHITE if bold else RGBColor(190, 220, 240))


def slide_recap(prs, section: str, title: str, points: list[str], next_hint: str = ""):
    slide = blank_slide(prs)
    add_header(slide, title, "섹션 정리", section)
    panel = slide.shapes.add_shape(1, Inches(0.35), Inches(1.55), Inches(12.55), Inches(4.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = C_RECAP
    panel.line.color.rgb = C_ACCENT
    for i, pt in enumerate(points):
        y = 1.75 + i * 0.85
        circ = slide.shapes.add_shape(9, Inches(0.6), Inches(y), Inches(0.38), Inches(0.38))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()
        nbox = slide.shapes.add_textbox(Inches(0.67), Inches(y + 0.02), Inches(0.3), Inches(0.35))
        ntf = nbox.text_frame
        ntf.clear()
        nr = ntf.paragraphs[0].add_run()
        nr.text = str(i + 1)
        set_run(nr, size=13, bold=True, color=C_WHITE)
        tbox = slide.shapes.add_textbox(Inches(1.15), Inches(y - 0.02), Inches(11.5), Inches(0.7))
        ttf = tbox.text_frame
        ttf.clear()
        tr = ttf.paragraphs[0].add_run()
        tr.text = pt
        set_run(tr, size=17)
    if next_hint:
        add_key_message(slide, f"다음 → {next_hint}", top=6.45)


def slide_text(prs, section, title, subtitle, key, bullets, footer=""):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    add_bullets(slide, bullets, MARGIN, CONTENT_TOP, 13.333 - 2 * MARGIN, FIG_BOTTOM - CONTENT_TOP, size=15)
    add_footer(slide, footer)


def slide_fig_right(prs, section, title, subtitle, key, image, caption, bullets, footer="", fig_ratio=0.62):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    usable = 13.333 - 2 * MARGIN
    fig_w = usable * fig_ratio
    text_w = usable - fig_w - 0.18
    text_left = MARGIN + fig_w + 0.18
    add_image(slide, image, MARGIN, CONTENT_TOP, fig_w, caption, valign="center")
    add_bullets(slide, bullets, text_left, CONTENT_TOP, text_w, MAX_FIG_H, size=13, compact=True)
    add_footer(slide, footer)


def slide_fig_top(prs, section, title, subtitle, key, image, caption, bullets, footer="", fig_h_ratio=0.72):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    usable = 13.333 - 2 * MARGIN
    fig_h = MAX_FIG_H * fig_h_ratio
    img_bottom = add_image(slide, image, MARGIN, CONTENT_TOP, usable, caption, max_h=fig_h, valign="center")
    bullet_top = img_bottom + 0.10
    bullet_h = FIG_BOTTOM - bullet_top
    mid = MARGIN + usable / 2
    half = (usable - 0.2) / 2
    if len(bullets) <= 2:
        add_bullets(slide, bullets, MARGIN, bullet_top, usable, bullet_h, size=13, compact=True)
    else:
        split = (len(bullets) + 1) // 2
        add_bullets(slide, bullets[:split], MARGIN, bullet_top, half, bullet_h, size=12, compact=True)
        add_bullets(slide, bullets[split:], mid + 0.1, bullet_top, half, bullet_h, size=12, compact=True)
    add_footer(slide, footer)


def slide_figure_full(prs, section, title, subtitle, key, image, footer=""):
    """그림 전용 — 설명은 Figure 안에, 슬라이드는 핵심+출처만."""
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    usable = 13.333 - 2 * MARGIN
    add_image(slide, image, MARGIN, CONTENT_TOP, usable, "", max_h=FIG_BOTTOM - CONTENT_TOP, valign="center")
    add_footer(slide, footer)


def slide_fig_hero(prs, section, title, subtitle, key, image, caption, bullets, footer=""):
    """Figure 거의 전폭 + 하단 2열 요약."""
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    usable = 13.333 - 2 * MARGIN
    fig_h = MAX_FIG_H * 0.82
    img_bottom = add_image(slide, image, MARGIN, CONTENT_TOP, usable, caption, max_h=fig_h, valign="center")
    bullet_top = img_bottom + 0.08
    bullet_h = FIG_BOTTOM - bullet_top
    mid = MARGIN + usable / 2
    half = (usable - 0.24) / 2
    split = (len(bullets) + 1) // 2
    add_bullets(slide, bullets[:split], MARGIN, bullet_top, half, bullet_h, size=12, compact=True)
    add_bullets(slide, bullets[split:], mid + 0.12, bullet_top, half, bullet_h, size=12, compact=True)
    add_footer(slide, footer)


def slide_table(prs, section, title, subtitle, key, headers, rows, footer=""):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle, section)
    add_key_message(slide, key)
    rs, cs = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(rs, cs, Inches(MARGIN), Inches(CONTENT_TOP), Inches(13.333 - 2 * MARGIN), Inches(min(0.42 * rs, FIG_BOTTOM - CONTENT_TOP))).table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = C_PRIMARY
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=12, bold=True, color=C_WHITE)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=11)
    add_footer(slide, footer)


def build(skip_extract: bool = False):
    if not skip_extract:
        run_prep()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== PROLOGUE =====
    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()
    for y, t, sz in [(2.0, "외삽 완전 정복", 40), (3.0, "이론 → OOD → 방법론 → 실전", 22), (4.1, "v3.1 | 국민대 IE Lab · 2026.07", 15)]:
        b = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.9))
        tf = b.text_frame
        tf.clear()
        r = tf.paragraphs[0].add_run()
        r.text = t
        set_run(r, size=sz, bold=sz > 20, color=C_WHITE)

    slide_text(
        prs, "", "왜 이 발표가 필요한가?", "문제 정의",
        "훈련 범위 밖 예측은 피할 수 없고, 잘못되면 위험하다",
        [
            "실험·공정 조건은 항상 훈련 범위를 벗어남 (RUL, 신소재, 극한 환경)",
            "일반 ML/DL은 train 분포 안에서만 잘 동작",
            "잘못된 외삽 → 제품 결함·안전사고",
            "필요: 이론 → OOD → 알고리즘 → NN 구조 → 실전 (N-CMAPSS)",
        ],
    )

    # ===== ACT 1: FOUNDATIONS =====
    slide_section(prs, "ACT 1", "외삽이란 무엇인가", "보간 vs 외삽 · Hull · 불확실성")

    if fig("pfister_extrapolation_fig2.png").exists():
        slide_fig_hero(
            prs, "ACT 1", "보간 vs 외삽", "Pfister 2024 Fig.2",
            "train 범위 안=보간(안전), 밖=외삽(위험)",
            fig("pfister_extrapolation_fig2.png"),
            "선형·비선형 조건부기댓값 · 외삽 구간에서 발산",
            [
                "보간: 훈련 범위 안 — 추가 가정 불필요",
                "외삽: 범위 밖 → 오차 급증",
                "test가 hull 안이면 보간 평가일 뿐",
            ],
            "Pfister & Bühlmann (2024)",
        )
    else:
        slide_text(
            prs, "ACT 1", "보간 vs 외삽", "가장 먼저 구분해야 할 개념",
            "train 범위 안=보간(안전), 밖=외삽(위험)",
            [
                "보간: 훈련 범위 안 예측 — 추가 가정 불필요",
                "외삽: 범위 밖 → 오차 급증 · 검증 불가",
                "2차 다항식: train 안 완벽 피팅 → 밖에서 발산 (Pfister 2024)",
                "test가 hull 안이면 보간 평가일 뿐",
            ],
            "Pfister & Bühlmann (2024)",
        )

    if fig("bonnasse_convex_hull_fig1.png").exists():
        slide_figure_full(
            prs, "ACT 1", "볼록 껍질 (Convex Hull)", "Bonnasse-Gahot 2022 Fig.1",
            "내재 공간 hull 안 ≠ 신경 표현 공간 hull 안",
            fig("bonnasse_convex_hull_fig1.png"),
            "Bonnasse-Gahot et al. (2022)",
        )
    else:
        slide_text(
            prs, "ACT 1", "볼록 껍질 (Convex Hull)", "훈련 데이터가 커버하는 범위",
            "hull 밖 = 추가 가정 없이는 추론 불가",
            [
                "Conv(X_train) = train 점들의 볼록 껍질",
                "내재 공간 hull 안 ≠ 신경 표현 공간 hull 안",
                "고차원일수록 hull 비율 ↓",
                "대책: 물리·구조적 제약",
            ],
            "Bonnasse-Gahot et al. (2022) · Pfister (2024)",
        )

    if fig("bartley_chla_tp.png").exists():
        slide_figure_full(
            prs, "ACT 1", "실제 데이터 예시", "Bartley 2019 Fig.1",
            "회귀선을 범위 밖까지 연장하면 위험",
            fig("bartley_chla_tp.png"),
            "Bartley et al. (2019) — Chl a vs TP",
        )
    else:
        slide_text(
            prs, "ACT 1", "실제 데이터 예시", "다변량에서 CI가 깨지는 이유",
            "회귀선을 범위 밖까지 연장하면 위험",
            [
                "Bartley 2019: Chl a–TP 회귀 · leverage 기반 외삽 구간",
                "단변량 CI는 hull 안에서만 의미",
                "다변량은 대부분 외삽 영역",
            ],
            "Bartley et al. (2019)",
        )

    slide_text(
        prs, "ACT 1", "불확실성 (UQ)이 필수인 이유", "숫자만 믿지 말 것",
        "hull 밖=train 데이터 없음 → epistemic UQ 폭증",
        [
            "우연적(Aleatoric): 측정 노이즈 — hull 안/밖 동일",
            "인식적(Epistemic): 모델 무지 — hull 밖에서 급증",
            "기존 CI는 외삽에서 coverage 깨짐",
            "extrapolation-aware CI 필요 (Pfister 2024)",
        ],
        "Ghahramani (2013) · Pfister (2024)",
    )

    slide_recap(
        prs, "ACT 1", "섹션 1 정리 — 외삽 기초",
        [
            "보간 ≠ 외삽. test가 hull 안이면 진짜 외삽 평가가 아님",
            "hull 밖은 Pfister (2024) 기준 추가 가정 없이 추론 불가",
            "외삽 예측에는 반드시 불확실성(UQ)을 함께 제시",
        ],
        "ACT 2: OOD 일반화 — ML 관점에서 외삽 이해",
    )

    # ===== ACT 2: OOD BASICS =====
    slide_section(prs, "ACT 2", "OOD 일반화 기초", "분포 이동 · 환경 · 목표")

    slide_text(
        prs, "ACT 2", "OOD란?", "Out-of-Distribution · 분포 밖",
        "train과 다른 분포의 test = OOD",
        [
            "OOD = Out-of-Distribution · P_test ≠ P_train",
            "Covariate shift: P(X) 변화, P(Y|X) 유지",
            "Concept shift: P(Y|X) 자체가 변화",
            "외삽과 가장 직접 연결: covariate shift",
        ],
        "Liu et al. (2023)",
    )

    if fig("ye_ood_failure_fig2.png").exists():
        slide_figure_full(
            prs, "ACT 2", "OOD 목표 vs ERM", "Ye 2021 Fig.2",
            "train 평균 최적 ≠ test(다른 도메인) 성공",
            fig("ye_ood_failure_fig2.png"),
            "Ye et al. (2021)",
        )
    else:
        slide_text(
            prs, "ACT 2", "OOD 목표 vs ERM", "평균 vs 최악 환경",
            "OOD = 최악 환경에서도 잘 동작 · ERM = train 평균만 최소화",
            [
                "ERM: train 평균 손실 ↓ → spurious feature 암기 가능",
                "OOD: 도메인/환경 shift 시 일반화 필요",
                "Ye 2021: train 평균 최적 ≠ test 성공",
            ],
            "Ye et al. (2021)",
        )

    slide_recap(
        prs, "ACT 2", "섹션 2 정리 — OOD 기초",
        [
            "OOD = Out-of-Distribution · train과 다른 test 분포",
            "Covariate shift = P(X)만 이동 — 외삽과 가장 직접 연결",
            "ERM = Empirical Risk Minimization · train 평균 손실 최소화",
            "ERM 평균 최적화 ≠ OOD(최악 환경) 성공",
        ],
        "ACT 3: OOD 알고리즘 — ERM, IRM, GroupDRO, DomainBed",
    )

    # ===== ACT 3: OOD ALGORITHMS =====
    slide_section(prs, "ACT 3", "OOD 알고리즘", "무엇을 쓸 수 있나")

    slide_text(
        prs, "ACT 3", "ERM — 반드시 비교할 baseline", "Empirical Risk Minimization",
        "ERM = min_θ (1/N)Σℓ(yᵢ,f_θ(xᵢ)) · OOD 논문의 필수 비교선",
        [
            "ERM = train N개 샘플 평균 손실 최소화 (기본 SGD/Adam)",
            "함정: spurious feature 암기 → OOD 급락",
            "baseline 요구: HP·모델 크기·탐색 횟수 동일",
            "DomainBed: ERM 없이는 '우리 방법이 더 낫다' 주장 불가",
        ],
        "Vapnik · Gulrajani & Lopez-Paz (2020)",
    )

    slide_figure_full(
        prs, "ACT 3", "IRM — SCM 구조", "Invariant Risk Minimization",
        "Z₁=인과 feature · Z₂=spurious feature (Colored MNIST)",
        fig("irm_scm_fig3.png"),
        "Arjovsky et al. (2019) Fig.3",
    )

    slide_text(
        prs, "ACT 3", "IRM — 결과", "인과 vs 비인과 weight",
        "IRM만 test 환경에서 non-causal weight 억제",
        [
            "Colored MNIST: 색-라벨 우연 상관",
            "ERM: 색 힌트 학습 → test(색 반전) 급락",
            "IRM: 환경 불변 Φ 학습 → test 유지",
            "한계: DomainBed에서 ERM ≥ IRM인 경우도",
        ],
        "Arjovsky et al. (2019)",
    )

    if fig("sagawa_groupdro_fig1.png").exists():
        slide_figure_full(
            prs, "ACT 3", "GroupDRO", "Sagawa 2020 Fig.1",
            "평균이 아닌 최악 그룹(worst-group) 보호",
            fig("sagawa_groupdro_fig1.png"),
            "Sagawa et al. (2020) ICLR — Waterbirds / CelebA",
        )
    else:
        slide_text(
            prs, "ACT 3", "GroupDRO", "Group Distributionally Robust Optimization",
            "평균이 아닌 최악 그룹(worst-group) 보호",
            [
                "Waterbirds: 새 종류 + 배경 spurious correlation",
                "GroupDRO: min max_g R_g(h)",
            ],
            "Sagawa et al. (2020)",
        )

    slide_text(
        prs, "ACT 3", "DomainBed", "공정 비교 결과",
        "HP·모델 선택 동일 시 ERM ≥ 기존 OOD SOTA",
        [
            "Rotated MNIST / PACS / VLCS / Office-Home",
            "ERM (공정 HP 탐색): 98.0 / 85.7 / 77.4 / 67.5%",
            "알고리즘 < HP 탐색 + baseline 비교",
            "Gulrajani & Lopez-Paz (2020)",
        ],
        "Gulrajani & Lopez-Paz (2020)",
    )

    slide_recap(
        prs, "ACT 3", "섹션 3 정리 — OOD 알고리즘",
        [
            "ERM = Empirical Risk Minimization (경험적 위험 최소화) — 모든 OOD 논문의 baseline",
            "IRM = Invariant Risk Minimization · GroupDRO = 최악 그룹 강건 최적화",
            "DomainBed: 공정 비교 시 ERM ≥ IRM인 경우도 많음 → baseline 필수",
            "알고리즘 < 실험 설계 + HP 탐색 + baseline 비교",
        ],
        "ACT 4: 신경망 외삽 — ReLU 한계와 구조적 해법",
    )

    # ===== ACT 4: NN EXTRAPOLATION =====
    slide_section(prs, "ACT 4", "신경망 외삽", "ReLU 한계 · 구조적 제약")

    slide_figure_full(
        prs, "ACT 4", "GNN 외삽 구조", "Xu et al. (2021) Fig.2",
        "아키텍처·입력 표현 = 암묵적 extrapolation prior",
        fig("xu_gnn_arch_fig2.png"),
        "Xu et al. (2021) ICLR",
    )

    slide_figure_full(
        prs, "ACT 4", "EQL — 수식 학습 구조", "Equation Learner",
        "sin·cos·×·÷ 연산 트리 → 해석 가능 수식",
        fig("eql_architecture_fig1.png"),
        "Martius & Lampert (2016) Fig.1",
    )

    slide_figure_full(
        prs, "ACT 4", "NALU — NAC/NALU 구조", "Neural Arithmetic Logic Units",
        "NAC(±) + exp/log gate(×÷) 구조 내장",
        fig("nalu_architecture_fig2.png"),
        "Trask et al. (2018) Fig.2",
    )

    slide_figure_full(
        prs, "ACT 4", "Monotonic Dense Unit", "Constrained Monotonic NN",
        "단조 제약 fully-connected layer 구조",
        fig("runje_monotonic_unit_fig3.png"),
        "Runje & Shankaranarayana (2023) Fig.3",
    )

    slide_figure_full(
        prs, "ACT 4", "Monotonic NN 아키텍처", "Neural architecture type 1",
        "Monotonic Dense Block × k + Final Activation",
        fig("runje_monotonic_arch_fig4.png"),
        "Runje & Shankaranarayana (2023) Fig.4",
    )

    if fig("xu_relu_extrapolation_fig1.png").exists():
        slide_figure_full(
            prs, "ACT 4", "ReLU MLP 외삽", "Xu 2021 Fig.1",
            "hull 밖 → affine(직선) 수렴",
            fig("xu_relu_extrapolation_fig1.png"),
            "Xu et al. (2021) ICLR",
        )
    if fig("fesser_pinn_failure_fig1.png").exists():
        slide_figure_full(
            prs, "ACT 4", "PINN 외삽 실패", "Fesser 2023 Fig.1",
            "보간 OK · 외삽 구간 오차·잔차 급증",
            fig("fesser_pinn_failure_fig1.png"),
            "Fesser et al. (2023) · Raissi PINN 한계",
        )
    elif fig("raissi_pinn_solution.png").exists():
        slide_figure_full(
            prs, "ACT 4", "PINN Burgers 해", "Raissi 2019 Fig.1",
            "L = L_data + λ·L_physics",
            fig("raissi_pinn_solution.png"),
            "Raissi et al. (2019)",
        )
    else:
        slide_text(
            prs, "ACT 4", "ReLU MLP · PINN", "구조 없으면 외삽 실패",
            "L_total = L_data + λ·L_physics (Raissi 2019)",
            [
                "ReLU MLP: hull 밖 → affine(직선) 수렴 (Xu 2021)",
                "PINN도 외삽 구간 실패 가능 (Fesser 2023)",
            ],
            "Xu · Raissi · Fesser",
        )

    slide_recap(
        prs, "ACT 4", "섹션 4 정리 — 신경망 외삽",
        [
            "ReLU MLP는 hull 밖에서 직선화 → 외삽 실패 (Xu 2021)",
            "구조적 제약(EQL, NALU, Monotonic)이 더 안전",
            "PINN도 외삽에서 실패 — UQ와 결합 필요",
        ],
        "ACT 5: N-CMAPSS — 이론을 실전에 적용",
    )

    # ===== ACT 5: CASE STUDY =====
    slide_section(prs, "ACT 5", "N-CMAPSS 실전", "APEX-Guard")

    slide_text(
        prs, "ACT 5", "왜 N-CMAPSS가 외삽인가?", "RUL + OOD 실험",
        "실험 설계가 외삽 연구의 출발점",
        [
            "목표: 항공 엔진 RUL 예측 (N-CMAPSS)",
            "Train: 저 TRA / Test: 고 TRA holdout → covariate shift",
            "Strict: ① TRA shift ② unit holdout ③ strict_late",
            "random split = 보간 평가일 뿐",
        ],
    )

    slide_text(
        prs, "ACT 5", "APEX-Guard 설계", "이론 → 구조 → 추론",
        "외삽을 '설계'로 방어",
        [
            "L1: strict_late — 순수 외삽 측정",
            "L2: CA-CSS (TRA(-), cycle_mono) — 구조적 제약",
            "L3: unit별 isotonic regression — 단조 외삽",
            "Dual Encoder: Health(TRA=0) + Load(TRA, cycle)",
        ],
    )

    slide_table(
        prs, "ACT 5", "strict_late 결과", "n=201 · 고 TRA OOD",
        "제약+단조가 TabPFN보다 OOD에서 우수",
        ["모델", "RMSE", "R²"],
        [
            ["APEX-Guard", "3.26", "0.947"],
            ["TabPFN", "3.80", "0.927"],
            ["Transformer", "7.31", "—"],
            ["TCN", "14.09", "—"],
        ],
        "ca-css-ncmapss strict_extrap_RESULTS",
    )

    slide_recap(
        prs, "ACT 5", "섹션 5 정리 — 실전 교훈",
        [
            "Hull 밖은 피할 수 없음 → CA-CSS + isotonic으로 방어",
            "TabPFN: 분포 내(in-distribution) 강, 고TRA OOD 급락",
            "도메인 지식(단조·물리) = 최고의 외삽 도구",
        ],
        "에필로그: 전체 핵심 정리",
    )

    # ===== EPILOGUE =====
    slide_recap(
        prs, "전체", "외삽을 정복하는 4가지 열쇠",
        [
            "① 이론: 보간≠외삽, hull 밖은 추가 가정 필요 (Pfister)",
            "② OOD: ERM 재검증 필수 (DomainBed)",
            "③ 구조: EQL·NALU·Monotonic·CA-CSS로 안전화",
            "④ UQ: 숫자 < 불확실성 — extrapolation-aware CI",
        ],
        "",
    )

    slide_text(
        prs, "", "필독 논문", "추가 읽을거리",
        "다음 부록에서 PDF를 더블클릭으로 열어보세요",
        [
            "★★★ Xu (2021) · Gulrajani (2020) · Pfister (2024)",
            "★★☆ Arjovsky IRM · Sagawa GroupDRO · Fesser PINN · Bonnasse-Gahot",
            "★★☆ Ye OOD · Trask NALU · Runje Monotonic · Raissi PINN · Bartley",
        ],
    )

    # ===== APPENDIX: papers + key figures + embedded PDFs =====
    slide_section(prs, "APPENDIX", "논문 PDF 첨부", "핵심 Figure + 원문 PDF · 더블클릭으로 열기")
    slide_paper_index(prs)
    for paper in PAPERS:
        slide_paper_detail(prs, paper)

    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()
    b = s.shapes.add_textbox(Inches(3.5), Inches(3.0), Inches(6), Inches(1.2))
    tf = b.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "Q & A"
    set_run(r, size=44, bold=True, color=C_WHITE)

    prs.save(OUT)
    import shutil
    shutil.copy(OUT, ROOT / "extrapolation-papers" / OUT.name)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    skip = "--skip-extract" in sys.argv
    build(skip_extract=skip)
